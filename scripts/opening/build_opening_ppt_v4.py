from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "docs" / "opening" / "opening_ppt_template_version_v3.pptx"
OUT = ROOT / "docs" / "opening" / "opening_ppt_template_version_v4.pptx"

TITLE = "面向 SQL+ 简化查询表达的多智能体自然语言数据库查询生成与反馈修正方法研究"
FONT = "Microsoft YaHei"

BLUE = RGBColor(24, 55, 92)
ORANGE = RGBColor(82, 92, 105)
GREEN = RGBColor(58, 78, 94)
RED = RGBColor(178, 38, 38)
DARK = RGBColor(31, 41, 51)
MUTED = RGBColor(102, 112, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(246, 248, 250)
PALE_BLUE = RGBColor(235, 240, 246)
PALE_ORANGE = RGBColor(241, 243, 245)
PALE_GREEN = RGBColor(235, 240, 246)
PALE_RED = RGBColor(250, 235, 235)


def clear_template_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_ids.remove(slide_id)


def set_text(shape, text: str, size: int = 15, bold: bool = False, color=DARK, align=PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.32), Inches(9.2), Inches(0.58))
    set_text(box, title, 24, True, BLUE)


def add_footer(slide, idx: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.03), Inches(8.9), Inches(0.22))
    set_text(box, f"开题汇报 | SQL+ / Text-to-SQL / Multi-Agent | {idx}", 8, False, MUTED)


def blank_slide(prs: Presentation, idx: int, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, title)
    add_footer(slide, idx)
    return slide


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text,
    fill=LIGHT,
    line=BLUE,
    size=11,
    bold=False,
    color=DARK,
    align=PP_ALIGN.CENTER,
):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    set_text(shape, text, size, bold, color, align)
    return shape


def add_label(slide, x, y, w, h, text, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, bold, color, align)
    return shape


def add_arrow(slide, x, y, text="→", size=18, color=BLUE):
    add_label(slide, x, y, 0.34, 0.26, text, size, True, color, PP_ALIGN.CENTER)


def add_bullets(slide, items: list[str], x=0.72, y=1.25, w=8.75, h=5.45, size=15) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def add_note(slide, text: str, y=6.35, color=ORANGE) -> None:
    add_box(slide, 0.72, y, 8.65, 0.48, text, PALE_RED, RED, 10, True, RED)


def add_focus(slide, text: str, y=5.85, size=15) -> None:
    add_box(slide, 0.88, y, 8.25, 0.58, text, PALE_RED, RED, size, True, RED)


def add_table(slide, rows: list[list[str]], x=0.45, y=1.18, w=9.1, h=5.25, size=8) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else (LIGHT if r % 2 == 0 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(size + 1 if r == 0 else size)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else DARK


def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        set_text(slide.shapes.title, TITLE, 27, True, BLUE, PP_ALIGN.CENTER)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            set_text(
                shape,
                "研究方向：Text-to-SQL / SQL+ 中间表示 / 多智能体 / 执行反馈修正\n汇报人：待填写    导师：待填写    时间：待填写",
                17,
                False,
                DARK,
                PP_ALIGN.CENTER,
            )
    add_footer(slide, 1)


def slide_outline(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "汇报结构")
    items = [
        ("1", "背景与问题", "自然语言数据库查询为什么仍不稳定"),
        ("2", "研究发展", "Text-to-SQL、IR、Agentic NL2SQL 的演进"),
        ("3", "动机测试", "按实验设置、方法、结果和分析说明为什么这样设计"),
        ("4", "系统设计", "SQL+ 中间表示与多智能体反馈闭环"),
        ("5", "可行性与初期验证", "转换、修复、repair skill 和 Spider 小子集"),
        ("6", "后续计划", "多组对比、消融、公开子集和达梦适配"),
    ]
    for i, (num, head, body) in enumerate(items):
        x = 0.65 + (i % 3) * 3.0
        y = 1.28 + (i // 3) * 2.15
        color = [BLUE, ORANGE, GREEN, RED, BLUE, ORANGE][i]
        fill = [PALE_BLUE, PALE_ORANGE, PALE_GREEN, PALE_RED, LIGHT, LIGHT][i]
        add_box(slide, x, y, 2.5, 1.5, "", fill, color)
        add_label(slide, x + 0.16, y + 0.14, 0.42, 0.34, num, 18, True, color, PP_ALIGN.CENTER)
        add_label(slide, x + 0.64, y + 0.18, 1.65, 0.3, head, 14, True, color)
        add_label(slide, x + 0.22, y + 0.68, 2.05, 0.52, body, 10, False, DARK)
    add_note(slide, "动机测试单独成段，先把实验怎么做、结果是什么、说明什么讲清楚。")


def slide_background(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "背景：自然语言数据库查询的真实需求")
    add_box(slide, 0.7, 1.25, 2.45, 1.15, "业务人员\n提出自然语言问题", PALE_BLUE, BLUE, 13, True)
    add_arrow(slide, 3.25, 1.65, "→", 18, MUTED)
    add_box(slide, 3.6, 1.25, 2.45, 1.15, "系统\n生成可执行查询", PALE_GREEN, GREEN, 13, True)
    add_arrow(slide, 6.15, 1.65, "→", 18, MUTED)
    add_box(slide, 6.5, 1.25, 2.45, 1.15, "数据库\n返回统计结果", PALE_ORANGE, ORANGE, 13, True)
    add_bullets(
        slide,
        [
            "企业数据分析中的很多需求本来就是自然语言表达，例如按地区统计销售额、找出消费最高客户、比较不同时间段的订单变化。",
            "Text-to-SQL 的目标是降低数据库使用门槛，但真实查询不仅要写对 SQL，还要理解 schema、字段值、业务约束和数据库方言。",
            "本课题关注的问题是：当第一次生成不可靠时，系统能否定位错误并做局部修复。",
        ],
        y=3.05,
        size=14,
    )


def slide_problem_anatomy(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "问题：SQL 错误通常不是单点语法错误")
    center = add_box(slide, 3.8, 2.75, 2.05, 0.8, "生成 SQL\n可执行但可能错", PALE_RED, RED, 13, True)
    problems = [
        ("schema linking", "表列选择错", 0.8, 1.2, BLUE),
        ("join path", "连接路径错", 6.8, 1.2, ORANGE),
        ("value linking", "过滤值不匹配", 0.8, 4.35, GREEN),
        ("aggregation", "聚合口径错", 6.8, 4.35, RED),
        ("projection/order", "输出列或排序错", 3.65, 5.35, BLUE),
    ]
    for name, desc, x, y, color in problems:
        add_box(slide, x, y, 2.15, 0.82, f"{name}\n{desc}", LIGHT, color, 10, True)
    add_label(slide, 3.26, 2.35, 0.5, 0.35, "→", 20, True, MUTED, PP_ALIGN.CENTER)
    add_label(slide, 5.9, 2.35, 0.5, 0.35, "←", 20, True, MUTED, PP_ALIGN.CENTER)
    add_label(slide, 3.26, 3.65, 0.5, 0.35, "→", 20, True, MUTED, PP_ALIGN.CENTER)
    add_label(slide, 5.9, 3.65, 0.5, 0.35, "←", 20, True, MUTED, PP_ALIGN.CENTER)
    add_note(slide, "如果每次都整条重写 SQL，容易修好一处又改坏另一处。")


def slide_timeline(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "研究发展：从 benchmark 到 agentic NL2SQL")
    years = [
        ("2018", "Spider\n跨数据库复杂 SQL"),
        ("2019", "IRNet/SemQL\n语义中间表示"),
        ("2020", "RAT-SQL\n关系感知 schema linking"),
        ("2021", "NatSQL / PICARD\n简化表示与约束解码"),
        ("2023", "BIRD / DIN-SQL / LEVER\n真实数据库与反馈验证"),
        ("2024", "Pipe Syntax / CHESS / CHASE-SQL\n线性 SQL 与多路径 agent"),
        ("2025", "Spider 2.0 / MAC-SQL / SQL-Factory\n企业工作流与多智能体生成"),
    ]
    add_label(slide, 0.8, 3.2, 8.5, 0.12, "────────────────────────────────────────", 13, True, MUTED)
    for i, (year, text) in enumerate(years):
        x = 0.52 + i * 1.3
        color = [BLUE, ORANGE, GREEN, BLUE, RED, ORANGE, GREEN][i]
        add_box(slide, x, 2.15, 1.05, 0.48, year, color, color, 12, True, WHITE)
        add_box(slide, x - 0.1, 3.58, 1.25, 1.12, text, LIGHT, color, 7, True)
        add_label(slide, x + 0.42, 2.74, 0.2, 0.35, "│", 14, True, color, PP_ALIGN.CENTER)
    add_note(slide, "发展趋势很清楚：任务越来越真实，系统越来越模块化，评价也从准确率扩展到反馈、效率和可修复性。")


def slide_benchmarks(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "相关技术一：公开 benchmark 的任务升级")
    rows = [
        ["代表工作", "年份", "关注点", "对本课题的启发"],
        ["Spider", "2018", "跨数据库、多表、复杂 SQL", "复杂 SQL 和 schema 泛化是基本难点"],
        ["BIRD", "2023", "真实数据库值、外部知识、执行效率", "value grounding 和效率反馈不能忽略"],
        ["Spider 2.0", "2025", "企业级 Text-to-SQL workflow", "真实系统需要方言、文档、代码库和多步执行环境"],
    ]
    add_table(slide, rows, y=1.2, h=3.2, size=9)
    add_box(slide, 0.8, 4.85, 8.4, 0.9, "本课题当前只做 Spider concert_singer 小子集可行性验证。\n后续再扩展到多数据库、BIRD 子集和达梦 SQL 方言。", PALE_RED, RED, 12, True)


def slide_ir_line(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "相关技术二：中间表示与 SQL 线性表达")
    items = [
        ("SemQL\n2019", "树形语义表示\n降低直接生成 SQL 压力", BLUE),
        ("NatSQL\n2021", "简化 SQL 外观\n弱化 FROM/JOIN 细节", ORANGE),
        ("GoogleSQL Pipe Syntax\n2024", "SQL 内部的管道式数据流", GREEN),
        ("SQL+\n本课题", "服务于生成、转换、诊断和局部修复", RED),
    ]
    for i, (head, body, color) in enumerate(items):
        x = 0.65 + i * 2.25
        add_box(slide, x, 1.45, 1.8, 2.45, "", LIGHT, color)
        add_label(slide, x + 0.18, 1.75, 1.4, 0.52, head, 13, True, color, PP_ALIGN.CENTER)
        add_label(slide, x + 0.18, 2.65, 1.42, 0.65, body, 10, False, DARK, PP_ALIGN.CENTER)
        if i < len(items) - 1:
            add_arrow(slide, x + 1.85, 2.48, "→", 16, MUTED)
    add_note(slide, "SQL+ 的目标不是重新发明 SQL，而是给错误定位和 repair skill 一个稳定的步骤载体。")


def slide_agentic_work(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "相关技术三：Agentic NL2SQL 的共同方向")
    rows = [
        ["系统", "年份", "核心做法", "仍留下的问题"],
        ["DIN-SQL", "2023", "schema linking、分解、生成、自修正", "修复粒度仍偏整体"],
        ["LEVER", "2023", "学习验证候选程序或 SQL", "更强调候选验证"],
        ["CHESS", "2024", "schema 检索、上下文组织、候选验证", "诊断和 patch 仍需细化"],
        ["CHASE-SQL", "2024", "多路径候选和偏好选择", "以候选选择为主"],
        ["MAC-SQL", "2025", "selector、decomposer、refiner 协作", "需要更清晰的局部修复接口"],
        ["SQL-Factory", "2025", "多智能体生成高质量 SQL 数据", "更适合作为后续数据扩充参考"],
    ]
    add_table(slide, rows, y=1.15, h=5.3, size=7)
    add_note(slide, "启发：Agent 的价值不在于数量，而在于每个模块是否有可检查输出。")


def slide_gap_summary(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "研究不足与本课题切入点")
    gaps = [
        ("中间表示", "多服务于初次生成，对执行反馈后的局部修复关注不足", BLUE),
        ("Agent 系统", "常用于分解、候选生成和选择，错误不一定映射到可修复步骤", ORANGE),
        ("评价指标", "execution accuracy 不够，需要 repairability 和成本指标", GREEN),
        ("真实环境", "公开 benchmark 与企业数据库方言、脏值、外部知识仍有差距", RED),
    ]
    for i, (head, body, color) in enumerate(gaps):
        x = 0.7 + (i % 2) * 4.45
        y = 1.35 + (i // 2) * 1.75
        add_box(slide, x, y, 3.75, 1.05, "", LIGHT, color)
        add_label(slide, x + 0.22, y + 0.16, 3.1, 0.26, head, 14, True, color)
        add_label(slide, x + 0.22, y + 0.52, 3.1, 0.34, body, 10, False, DARK)
    add_focus(slide, "切入点：让 Critic、Router、Repair Skill 和 Executor 围绕 SQL+ 步骤工作。", y=5.3, size=14)


def slide_system_overview(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "提出的系统：SQL+ 多智能体反馈修正框架")
    top = [
        ("用户问题", PALE_BLUE, BLUE),
        ("SQL+ Generator", PALE_BLUE, BLUE),
        ("Translator", PALE_GREEN, GREEN),
        ("Executor", PALE_GREEN, GREEN),
    ]
    for i, (text, fill, line) in enumerate(top):
        x = 0.65 + i * 2.25
        add_box(slide, x, 1.28, 1.65, 0.7, text, fill, line, 10, True)
        if i < len(top) - 1:
            add_arrow(slide, x + 1.7, 1.48, "→", 16, MUTED)
    bottom = [
        ("Critic Agent\n错误类型与步骤定位", PALE_ORANGE, ORANGE),
        ("Skill Router\n选择局部 repair skill", PALE_ORANGE, ORANGE),
        ("Repair Skill\n生成 SQL+ patch", PALE_RED, RED),
        ("Executor\n执行验证候选", PALE_GREEN, GREEN),
    ]
    for i, (text, fill, line) in enumerate(bottom):
        x = 0.65 + i * 2.25
        add_box(slide, x, 4.12, 1.65, 0.86, text, fill, line, 9, True)
        if i < len(bottom) - 1:
            add_arrow(slide, x + 1.7, 4.42, "→", 16, MUTED)
    add_label(slide, 7.85, 2.05, 0.3, 1.7, "↓\n反馈", 14, True, ORANGE, PP_ALIGN.CENTER)
    add_label(slide, 4.9, 3.25, 0.38, 0.38, "↺", 22, True, RED, PP_ALIGN.CENTER)
    add_focus(slide, "核心思路：不整条重写 SQL，把错误压回 SQL+ 的局部步骤中修。", y=6.12, size=14)


def slide_sqlplus_example(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "SQL+ 表达示例：把查询拆成步骤")
    sql = (
        "FROM orders o\n"
        "| JOIN order_items oi ON o.order_id = oi.order_id\n"
        "| WHERE o.status = 'paid'\n"
        "| GROUP o.order_id\n"
        "| AGG o.order_id, SUM(oi.quantity * oi.unit_price) AS amount\n"
        "| ORDER amount DESC\n"
        "| LIMIT 10"
    )
    add_box(slide, 0.75, 1.18, 4.5, 4.95, sql, WHITE, BLUE, 12, False, DARK, PP_ALIGN.LEFT)
    rows = [
        ["SQL+ 步骤", "可诊断问题"],
        ["WHERE", "状态值 paid 是否正确"],
        ["GROUP / AGG", "聚合口径和别名是否正确"],
        ["ORDER / LIMIT", "top-k 排序字段和方向是否正确"],
        ["JOIN", "连接路径是否必要且正确"],
    ]
    add_table(slide, rows, x=5.65, y=1.35, w=3.65, h=3.2, size=8)
    add_box(slide, 5.65, 5.0, 3.65, 0.82, "重点：每个步骤都是定位和局部修改的锚点。", PALE_RED, RED, 11, True, RED)


def slide_architecture_detail(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "系统架构设计：可观察的多智能体分工")
    rows = [
        ["模块", "输入", "输出", "评价方式"],
        ["Schema Agent", "问题、schema、值域样例", "相关表列、join path、候选值", "schema/value/join accuracy"],
        ["SQL+ Generator", "问题、schema 上下文", "初始 SQL+", "SQL+ valid, exec match"],
        ["Critic Agent", "SQL+、执行反馈、结果摘要", "错误类型、可疑步骤、证据", "localization accuracy"],
        ["Skill Router", "Critic 输出、SQL+ 结构", "repair skill 路由", "router accuracy"],
        ["Repair Skill", "SQL+ 局部步骤、反馈", "候选 patch", "patch minimality"],
        ["Executor", "候选 SQL", "执行结果、错误、最终选择", "repair success, latency"],
    ]
    add_table(slide, rows, y=1.05, h=5.65, size=7)


def slide_research_questions(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "研究问题与验证路径")
    rows = [
        ["研究问题", "验证路径", "关键指标"],
        ["RQ1 SQL+ 是否有必要", "与 SQL、SemQL-style、NatSQL-style、Pipe-style 对比", "跨子句引用、别名依赖、转换成功率"],
        ["RQ2 SQL+ 能否稳定生成", "Direct NL2SQL 与 NL2SQL+ 对比", "valid rate、execution match、token、latency"],
        ["RQ3 反馈能否定位到步骤", "Critic 与 oracle/人工差异对比", "error localization accuracy"],
        ["RQ4 局部 skill 是否优于整条重写", "SQL+ Router Skills 与 Direct SQL Refiner 对比", "repair success、patch minimality、repair rounds"],
        ["RQ5 能否迁移到公开子集", "Spider/BIRD/达梦子集逐步扩展", "exec match、方言错误、泛化失败类型"],
    ]
    add_table(slide, rows, y=1.05, h=5.6, size=7)


def slide_technical_route(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "研究方案与技术路线")
    phases = [
        ("1 SQL+ 表达层", "parser\nconverter\nAST 映射", BLUE),
        ("2 初始生成层", "Schema Agent\nSQL+ Generator\nTranslator", ORANGE),
        ("3 反馈诊断层", "Executor\nCritic Agent\n结果摘要", GREEN),
        ("4 局部修复层", "Skill Router\nRepair Skills\nPatch 验证", RED),
        ("5 实验评估层", "对比实验\n消融实验\n公开子集", BLUE),
    ]
    for i, (head, body, color) in enumerate(phases):
        x = 0.45 + i * 1.85
        add_box(slide, x, 1.45, 1.45, 2.7, "", LIGHT, color)
        add_label(slide, x + 0.12, 1.72, 1.2, 0.42, head, 11, True, color, PP_ALIGN.CENTER)
        add_label(slide, x + 0.12, 2.55, 1.2, 0.8, body, 9, False, DARK, PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            add_arrow(slide, x + 1.48, 2.62, "→", 16, MUTED)
    add_note(slide, "路线按可复现脚本推进：先可控自建数据集，再公开子集，最后考虑达梦方言。", y=5.55)


def slide_experiment_overview(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "前期可行性实验：总体设置")
    rows = [
        ["项目", "当前设置"],
        ["自建数据库", "企业订单分析样例库"],
        ["数据表", "customers、products、orders、order_items"],
        ["自然语言查询", "30 条"],
        ["SQL+ 标准样例", "30 条"],
        ["修复样例", "15 条规则修正样例；13 条 SQL+ prompt v2 已知失败样例"],
        ["执行环境", "SQLite 内存数据库"],
        ["模型", "gpt-5-mini"],
        ["评价方式", "执行生成 SQL，并与 gold SQL 执行结果比较"],
        ["gold 边界", "gold SQL 只用于离线评估，非 gold 修复不把 gold 差异输入模型"],
    ]
    add_table(slide, rows, y=1.05, h=5.75, size=8)


def slide_exp_conversion(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验一：SQL+ 表达与转换可行性")
    rows = [
        ["维度", "内容"],
        ["实验类型", "验证系统和方案可行性的基础实验"],
        ["实验设置", "30 条人工 SQL+ 标准样例，覆盖过滤、join、聚合、排序、top-k"],
        ["实验过程", "SQL+ parser 解析 -> converter 转 SQL -> SQLite 执行 -> 与 gold SQL 结果比较"],
        ["实验结果", "SQL+ 语法通过 30/30；SQL 可执行 30/30；执行一致 30/30"],
        ["实验结论", "SQL+ 已形成可解析、可转换、可执行的最小闭环"],
    ]
    add_table(slide, rows, y=1.1, h=5.25, size=9)
    add_focus(slide, "注意：这里验证转换链路，不是模型生成准确率。", y=6.18, size=14)


def slide_exp_baseline(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验二：Direct NL2SQL 与 NL2SQL+ baseline")
    rows = [
        ["方法", "样例数", "SQL/SQL+ 有效", "SQL 可执行", "执行一致"],
        ["Direct NL2SQL", "30", "-", "30/30", "16/30"],
        ["NL2SQL+ prompt v1", "30", "30/30", "30/30", "13/30"],
        ["NL2SQL+ prompt v2", "30", "30/30", "30/30", "17/30"],
    ]
    add_table(slide, rows, y=1.15, h=2.35, size=9)
    rows2 = [
        ["失败类型", "数量", "含义"],
        ["filter/value-linking", "5", "字段值、日期、状态值错误"],
        ["ORDER/LIMIT", "3", "排序字段、方向或 top-k 错误"],
        ["aggregation planning", "2", "GROUP、COUNT 或聚合别名错误"],
        ["schema/join planning", "2", "表连接路径错误"],
        ["projection mismatch", "1", "输出列不符合问题要求"],
    ]
    add_table(slide, rows2, y=3.78, h=2.65, size=8)
    add_note(slide, "结论：只换成 SQL+ 输出格式不够，主要错误是语义错误。", y=6.43)


def slide_exp_ir_complexity(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验三：IR 表达复杂度对比")
    rows = [
        ["表示形式", "平均 token", "步骤/子句", "嵌套深度", "别名依赖", "跨子句引用"],
        ["Standard SQL", "31.5333", "5.9", "0.6667", "2.0333", "2.3333"],
        ["SQL+", "35.0333", "6.1333", "0.6667", "0.7", "1.0"],
        ["SemQL-style proxy", "50.5667", "10.7333", "3.6667", "0.9", "1.2"],
        ["NatSQL-style proxy", "31.5", "5.4333", "0.9667", "1.3667", "1.6667"],
        ["Pipe-style proxy", "40.8", "6.1333", "0.6667", "1.3667", "1.6667"],
    ]
    add_table(slide, rows, y=1.05, h=3.75, size=8)
    add_box(slide, 0.75, 5.15, 8.55, 0.95, "结论：SQL+ 不是更短。它用显式步骤边界换取更低别名依赖和跨子句引用，更适合后续错误定位和局部修复。", PALE_ORANGE, ORANGE, 12, True)


def slide_exp_ir_generation(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验四：IR 生成成本与执行效果")
    rows = [
        ["方法", "表示有效", "SQL 可执行", "执行一致", "平均总 token", "平均延迟"],
        ["Direct SQL", "30/30", "30/30", "12/30", "599.1667", "6.5851s"],
        ["SQL+", "28/30", "28/30", "14/30", "813.0333", "9.2197s"],
        ["NatSQL-style proxy", "30/30", "30/30", "13/30", "740.7667", "6.2802s"],
        ["SemQL-style proxy", "30/30", "25/30", "12/30", "1028.9667", "9.9684s"],
    ]
    add_table(slide, rows, y=1.15, h=3.35, size=9)
    add_box(slide, 0.75, 4.95, 8.55, 1.1, "结论：SQL+ 初次生成略高，但 token 和 latency 也更高。开题论证不能停在生成阶段，必须继续看修复阶段收益。", PALE_GREEN, GREEN, 12, True)


def slide_exp_repair(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验二：反馈修正与 repairability")
    rows = [
        ["方法", "初始失败", "SQL+ 有效", "SQL 可执行", "修复成功", "说明"],
        ["SQL+ 诊断辅助 Refiner", "13", "13/13", "13/13", "13/13", "使用 gold-derived differences"],
        ["SQL+ 非 gold Refiner v2", "13", "13/13", "12/13", "4/13", "执行反馈和粗粒度诊断"],
        ["Direct SQL 非 gold Refiner", "14", "-", "14/14", "6/14", "直接修标准 SQL"],
        ["SQL+ Skill Router + Repair Skills v3", "13", "13/13", "13/13", "13/13", "Critic 路由到五类局部 skill"],
    ]
    add_table(slide, rows, y=1.05, h=3.55, size=7)
    rows2 = [
        ["方法", "修复成功", "定位准确率", "严格最小 patch", "平均 repair token"],
        ["Direct SQL Refiner", "6/14", "0.8571", "0.8571", "1609.3571"],
        ["SQL+ Router Skills", "13/13", "0.7692", "0.9231", "3813.9231"],
    ]
    add_table(slide, rows2, y=4.92, h=1.42, size=8)
    add_focus(slide, "结论：SQL+ 路线修复成功率和 patch 可控性更好，但 token 成本更高。", y=6.42, size=13)


def slide_exp_skills(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验三：Repair Skill 分治结果")
    rows = [
        ["Repair skill", "样例数", "修复成功", "覆盖的典型问题"],
        ["value-linking", "3", "3/3", "候选值替换、日期边界归一化、值过滤错误"],
        ["ORDER", "3", "3/3", "排序字段错误、升降序错误、LIMIT/Top-K 约束"],
        ["aggregation", "3", "3/3", "COUNT 口径、GROUP 维度、AGG 别名、ORDER/HAVING 引用"],
        ["join", "3", "3/3", "JOIN 方向、冗余 JOIN、缺失 JOIN、paid 过滤遗漏"],
        ["projection", "1", "1/1", "结果列多、列少或列顺序错误"],
    ]
    add_table(slide, rows, y=1.15, h=4.2, size=8)
    add_focus(slide, "结论：不同错误类型需要不同局部修复策略，整体重写不是最稳的路径。", y=5.76, size=13)


def slide_exp_spider(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验四：Spider 小规模公开子集")
    rows = [
        ["实验", "样例", "SQL+ 有效", "SQL 可执行", "执行一致", "说明"],
        ["conversion smoke test", "20", "20/20", "20/20", "20/20", "Spider gold SQL -> SQL+ -> SQL"],
        ["fresh e2e generation", "20", "19/20", "19/20", "19/20", "question + schema -> SQL+"],
        ["fresh e2e + semantic repair", "20", "20/20", "20/20", "20/20", "同一次 fresh 输出经 Router -> semantic repair skill"],
    ]
    add_table(slide, rows, y=1.15, h=2.8, size=7)
    add_box(slide, 0.75, 4.45, 8.55, 1.25, "结论：SQL+ 表达、转换、fresh 生成和 semantic repair 链路已在 concert_singer 小子集跑通。\n边界：conversion smoke 使用 gold SQL 改写，不是端到端准确率；整体不是完整 Spider benchmark 跑分。", PALE_RED, RED, 11, True, RED)


def slide_feasibility_summary(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "前期可行性分析小结")
    cards = [
        ("表达可行", "SQL+ -> SQL\n30/30 执行一致", BLUE),
        ("生成有潜力", "SQL+ v2 17/30\nDirect SQL 16/30", ORANGE),
        ("修复更关键", "单 Refiner 不稳\nRouter Skills 13/13", GREEN),
        ("成本需补齐", "SQL+ token 更高\nlatency 需重跑", RED),
        ("公开子集已跑通", "Spider 20 条\nfresh 19/20 -> 20/20", BLUE),
        ("不能过度外推", "小样例\n非完整 benchmark", ORANGE),
    ]
    for i, (head, body, color) in enumerate(cards):
        x = 0.65 + (i % 3) * 3.0
        y = 1.35 + (i // 3) * 1.8
        add_box(slide, x, y, 2.55, 1.25, "", LIGHT, color)
        add_label(slide, x + 0.18, y + 0.18, 2.1, 0.28, head, 12, True, color)
        add_label(slide, x + 0.18, y + 0.58, 2.1, 0.42, body, 10, False, DARK, PP_ALIGN.CENTER)
    add_focus(slide, "当前结果支撑开题可行性，但还不足以支撑泛化结论。", y=5.66, size=14)


def slide_future_experiments(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "后续实验设计：多组对比与消融")
    rows = [
        ["实验组", "目的", "主要指标"],
        ["Direct NL2SQL", "标准 SQL 直接生成基线", "exec accuracy, valid SQL, token, latency"],
        ["NL2SQL+ single agent", "检验只换中间表示是否有效", "SQL+ valid, exec match"],
        ["SemQL/NatSQL/Pipe-style proxy", "比较中间表示形态", "复杂度、转换、生成成本"],
        ["Standard SQL multi-agent", "排除多智能体本身带来的收益", "repair success, cost"],
        ["SQL+ multi-agent without feedback", "检验反馈闭环是否必要", "exec match, failure type"],
        ["SQL+ Router + Repair Skills", "验证局部技能路由", "localization, router accuracy, patch minimality"],
        ["Ablation", "去掉 Critic、Router、skill、executor", "性能下降幅度"],
    ]
    add_table(slide, rows, y=1.05, h=5.75, size=7)


def slide_metrics(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "评价指标：不只看最终准确率")
    groups = [
        ("结果质量", "execution accuracy\nvalid SQL rate\nSQL+ valid rate", BLUE),
        ("表达复杂度", "token length\nnesting depth\nalias dependency\ncross-clause reference", ORANGE),
        ("修复能力", "repair success\nlocalization accuracy\nrouter accuracy\npatch minimality", GREEN),
        ("成本效率", "prompt/completion tokens\nlatency\nrepair rounds\nconversion time", RED),
        ("可解释性", "步骤级诊断\n修改范围\n工具调用 trace", BLUE),
    ]
    for i, (head, body, color) in enumerate(groups):
        x = 0.65 + (i % 3) * 3.0
        y = 1.25 + (i // 3) * 2.05
        add_box(slide, x, y, 2.55, 1.45, "", LIGHT, color)
        add_label(slide, x + 0.18, y + 0.15, 2.1, 0.28, head, 13, True, color, PP_ALIGN.CENTER)
        add_label(slide, x + 0.18, y + 0.52, 2.1, 0.72, body, 9, False, DARK, PP_ALIGN.CENTER)
    add_note(slide, "只有把质量、成本、定位和 patch 范围一起看，才能判断 SQL+ 是否真的值得。", y=5.75)


def slide_limits(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "当前不足与风险控制")
    rows = [
        ["不足", "影响", "后续处理"],
        ["自建数据集较小", "不能证明大规模泛化", "扩展 Spider/BIRD 子集和更多查询类型"],
        ["已知失败集偏窄", "13/13 可能只覆盖当前错误类型", "增加未知错误、复合错误、无报错语义错"],
        ["semantic repair 仍是原型", "目前只在单库小子集验证", "抽象统一 repair 接口后跨数据集复测"],
        ["SQL+ 生成成本较高", "可能抵消修复收益", "补齐端到端 token、latency 和修复轮数"],
        ["达梦方言未完成", "不能直接说明生产数据库可用", "补充日期、分页、类型转换和字符串函数适配"],
    ]
    add_table(slide, rows, y=1.15, h=4.85, size=8)
    add_note(slide, "这些不是结论失败，而是后续实验必须回答的问题。", y=6.2)


def slide_schedule(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "后续进度安排")
    rows = [
        ["阶段", "主要工作"],
        ["开题后 1-2 个月", "完善 SQL+ 语法子集，扩展无报错语义错和 projection 诊断样例"],
        ["第 3-4 个月", "完善 Schema Agent、Critic Agent、Skill Router，引入 parser/executor/value lookup 工具"],
        ["第 5-6 个月", "接入 Spider/BIRD 小规模子集，扩展多数据库、多难度实验"],
        ["第 7-8 个月", "适配达梦 SQL 方言，测试函数、分页、日期和类型转换规则"],
        ["第 9-10 个月", "完成消融实验、结果分析、论文初稿和原型系统整理"],
        ["第 11-12 个月", "根据导师和中期反馈修改，完成论文定稿、答辩材料和代码归档"],
    ]
    add_table(slide, rows, y=1.2, h=5.35, size=8)


def slide_contributions(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "预期创新点与成果")
    items = [
        ("SQL+ 中间表示", "面向生成、转换、诊断和局部修复共同设计的步骤化查询表示。", BLUE),
        ("多智能体反馈修正框架", "Schema Agent、Critic Agent、Skill Router、Repair Skill、Executor 形成可检查闭环。", ORANGE),
        ("repairability 评价体系", "在 execution accuracy 外，比较定位、路由、patch 范围、修复轮数、token 和 latency。", GREEN),
        ("可复现实验与原型", "形成数据、脚本、实验日志、开题材料和后续达梦方言适配路线。", RED),
    ]
    for i, (head, body, color) in enumerate(items):
        x = 0.75 + (i % 2) * 4.35
        y = 1.35 + (i // 2) * 1.8
        add_box(slide, x, y, 3.75, 1.1, "", LIGHT, color)
        add_label(slide, x + 0.22, y + 0.16, 3.1, 0.28, head, 14, True, color)
        add_label(slide, x + 0.22, y + 0.54, 3.1, 0.34, body, 10, False, DARK)
    add_note(slide, "预期成果包括论文、原型系统、实验报告和可复现实验脚本。", y=5.55)


def slide_summary(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "总结")
    add_bullets(
        slide,
        [
            "本课题研究重点不是一次生成 SQL，而是 SQL+ 中间表示支撑下的执行反馈诊断和局部修复。",
            "相关研究说明 Text-to-SQL 正在从复杂 SQL 生成走向真实数据库、工具调用、多候选验证和企业工作流。",
            "现有不足在于错误反馈缺少稳定的步骤级落点，多智能体修复也缺少局部 patch 评价。",
            "前期实验已经证明 SQL+ 转换闭环和已知失败集上的局部修复可行，但当前结果仍属于开题阶段小规模证据。",
            "后续工作将集中在更大样例、更复杂错误、完整成本记录、公开子集迁移和达梦 SQL 方言适配。",
        ],
        size=14,
    )


def slide_refs(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "主要参考文献")
    refs = [
        "Yu et al. 2018. Spider: A Large-Scale Human-Labeled Dataset for Complex and Cross-Domain Text-to-SQL.",
        "Guo et al. 2019. IRNet: A General Purpose Text-to-SQL Parser with SemQL.",
        "Wang et al. 2020. RAT-SQL: Relation-Aware Schema Encoding and Linking.",
        "Gan et al. 2021. NatSQL: A Simplified SQL Intermediate Representation.",
        "Scholak et al. 2021. PICARD: Parsing Incrementally for Constrained Auto-Regressive Decoding.",
        "Li et al. 2023. BIRD: A Big Bench for Large-Scale Database Grounded Text-to-SQLs.",
        "Pourreza and Rafiei. 2023. DIN-SQL: Decomposed In-Context Learning of Text-to-SQL.",
        "Ni et al. 2023. LEVER: Learning to Verify Language-to-Code Generation.",
        "Shute et al. 2024. SQL Has Problems. We Can Fix Them: Pipe Syntax in SQL.",
        "Talaei et al. 2024. CHESS: Contextual Harnessing for Efficient SQL Synthesis.",
        "Pourreza et al. 2024. CHASE-SQL: Multi-Path Reasoning and Preference Optimized Candidate Selection.",
        "Lei et al. 2025. Spider 2.0: Evaluating Language Models on Enterprise Text-to-SQL Workflows.",
        "Wang et al. 2025. MAC-SQL: A Multi-Agent Collaborative Framework for Text-to-SQL.",
        "Li et al. 2025. SQL-Factory: A Multi-Agent Framework for High-Quality and Large-Scale SQL Generation.",
    ]
    add_bullets(slide, refs, y=1.05, size=9)


def slide_background_problem_merged(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "背景与问题：自然语言查询需要可修复")
    add_box(slide, 0.65, 1.2, 2.3, 0.88, "业务人员\n自然语言问题", PALE_BLUE, BLUE, 12, True)
    add_arrow(slide, 3.05, 1.45, "→", 16, MUTED)
    add_box(slide, 3.45, 1.2, 2.3, 0.88, "LLM / Agent\n生成 SQL", PALE_GREEN, GREEN, 12, True)
    add_arrow(slide, 5.86, 1.45, "→", 16, MUTED)
    add_box(slide, 6.25, 1.2, 2.6, 0.88, "数据库执行\n结果或错误反馈", PALE_ORANGE, ORANGE, 12, True)
    rows = [
        ["错误层级", "常见表现", "为什么难修"],
        ["schema linking", "表列选择、连接路径错误", "会传导到 WHERE、GROUP 和 SELECT"],
        ["value linking", "枚举值、日期、阈值错误", "执行可能不报错，但结果语义错"],
        ["aggregation", "COUNT、GROUP、HAVING 口径错误", "涉及跨子句依赖"],
        ["projection/order", "输出列、排序字段、top-k 错误", "结果看似合理但不符合问题"],
    ]
    add_table(slide, rows, y=2.55, h=3.15, size=8)
    add_focus(slide, "重点：本课题关注生成失败后的定位、路由和局部修复。", y=6.08)


def slide_related_tech_merged(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "当前研究与发展：benchmark、IR 与 agentic NL2SQL")
    rows = [
        ["方向", "代表工作与年份", "已解决问题", "对本课题的启发"],
        ["公开评测", "Spider 2018、BIRD 2023、Spider 2.0 2025", "复杂 schema、真实数据、企业工作流", "后续需要从小子集扩展到多库和真实方言"],
        ["中间表示", "SemQL 2019、NatSQL 2021、Pipe Syntax 2024", "降低 SQL 生成或表达复杂度", "SQL+ 要证明修复接口价值，而不只是换语法"],
        ["约束与验证", "PICARD 2021、LEVER 2023", "语法约束、执行验证、候选筛选", "执行反馈需要进入诊断和 patch 过程"],
        ["Agentic NL2SQL", "DIN-SQL 2023、CHESS 2024、CHASE-SQL 2024、MAC-SQL 2025", "分解、多候选、多角色协作", "需要更清楚的步骤级错误定位和 repair skill 评价"],
        ["数据生成", "SQL-Factory 2025", "多智能体生成高质量 SQL 数据", "可用于后续扩充训练和评测样例"],
    ]
    add_table(slide, rows, y=1.05, h=5.2, size=7)
    add_note(slide, "研究趋势是从单次生成走向真实环境、工具反馈、多智能体协作和可验证修复。", y=6.42)


def slide_research_route_merged(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "研究问题与技术路线")
    rows = [
        ["研究问题", "验证路径", "指标"],
        ["SQL+ 是否有必要", "与 SQL、SemQL-style、NatSQL-style、Pipe-style 对比", "依赖、复杂度、转换成功率"],
        ["只换表示是否足够", "Direct NL2SQL 与 NL2SQL+ baseline 对比", "valid rate、exec match、cost"],
        ["局部修复是否有效", "single refiner 与 Router + skills 对比", "repair success、patch minimality"],
        ["能否迁移", "Spider/BIRD/达梦子集逐步扩展", "exec match、失败类型、方言错误"],
    ]
    add_table(slide, rows, y=1.02, h=2.45, size=8)
    phases = [
        ("SQL+ 表达", "parser\nconverter", BLUE),
        ("初始生成", "Schema Agent\nSQL+ Generator", ORANGE),
        ("反馈诊断", "Executor\nCritic Agent", GREEN),
        ("局部修复", "Skill Router\nRepair Skills", RED),
        ("实验评估", "对比\n消融\n公开子集", BLUE),
    ]
    for i, (head, body, color) in enumerate(phases):
        x = 0.48 + i * 1.82
        add_box(slide, x, 4.1, 1.42, 1.45, "", LIGHT, color)
        add_label(slide, x + 0.1, 4.3, 1.18, 0.28, head, 10, True, color, PP_ALIGN.CENTER)
        add_label(slide, x + 0.12, 4.78, 1.15, 0.45, body, 8, False, DARK, PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            add_arrow(slide, x + 1.43, 4.75, "→", 14, MUTED)
    add_note(slide, "技术路线从可控 SQL+ 表达层开始，逐步扩展到反馈修复和公开数据集评估。", y=6.16)


def slide_exp_ir_merged(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验三：IR 复杂度、生成成本与边界")
    rows = [
        ["表示形式", "平均 token", "别名依赖", "跨子句引用", "执行一致 / 成本结论"],
        ["Standard SQL", "31.5333", "2.0333", "2.3333", "Direct SQL 12/30，平均 599.1667 tokens"],
        ["SQL+", "35.0333", "0.7", "1.0", "SQL+ 14/30，平均 813.0333 tokens"],
        ["NatSQL-style proxy", "31.5", "1.3667", "1.6667", "13/30，平均 740.7667 tokens"],
        ["SemQL-style proxy", "50.5667", "0.9", "1.2", "12/30，平均 1028.9667 tokens"],
        ["Pipe-style proxy", "40.8", "1.3667", "1.6667", "用于复杂度对照"],
    ]
    add_table(slide, rows, y=1.05, h=4.15, size=7)
    add_box(slide, 0.75, 5.55, 8.55, 0.72, "结论：SQL+ 不是更短，也不是无成本。它的价值主要在降低步骤耦合，为错误定位和局部 patch 提供接口。", PALE_ORANGE, ORANGE, 11, True)


def slide_future_metrics_merged(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "后续实验设计与评价指标")
    rows = [
        ["实验组", "目的", "主要指标"],
        ["Direct SQL / NL2SQL+", "比较直接生成和 SQL+ 生成", "execution accuracy、valid rate、token、latency"],
        ["SemQL/NatSQL/Pipe-style proxy", "比较中间表示形态", "复杂度、转换、生成成本"],
        ["whole-query rewrite / single refiner / Router skills", "比较修复策略", "repair success、repair rounds、patch minimality"],
        ["消融实验", "去掉 Critic、Router、skill、executor", "性能下降幅度、错误类型变化"],
        ["Spider/BIRD/达梦子集", "验证迁移和方言适配", "exec match、方言错误、泛化失败类型"],
    ]
    add_table(slide, rows, y=1.05, h=3.8, size=8)
    add_box(slide, 0.75, 5.25, 2.65, 0.72, "结果质量\nexecution accuracy / valid rate", PALE_BLUE, BLUE, 10, True)
    add_box(slide, 3.62, 5.25, 2.65, 0.72, "修复能力\nlocalization / router / patch", PALE_GREEN, GREEN, 10, True)
    add_box(slide, 6.49, 5.25, 2.65, 0.72, "成本效率\ntoken / latency / rounds", PALE_ORANGE, ORANGE, 10, True)


def slide_limits_schedule_merged(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "当前不足、风险控制与进度安排")
    rows = [
        ["风险", "当前情况", "控制方式"],
        ["样本规模小", "known-failure set 和 Spider 子集都较小", "扩展失败集、Spider 多库和 BIRD 子集"],
        ["结果被误读", "conversion smoke 来自 gold SQL", "PPT 和论文中始终标注边界"],
        ["SQL+ 成本较高", "token 和 latency 未必占优", "补齐成本实验，用 repairability 解释收益"],
        ["方言适配不足", "当前主要是 SQLite", "补达梦日期、分页、类型转换和函数规则"],
    ]
    add_table(slide, rows, y=1.05, h=3.05, size=8)
    timeline = [
        ("07-08", "统一生成入口\n扩充失败集"),
        ("09-10", "消融实验\n公开子集"),
        ("11-12", "达梦适配\n论文初稿"),
        ("后续", "复现实验\n答辩材料"),
    ]
    for i, (time, work) in enumerate(timeline):
        x = 0.75 + i * 2.2
        add_box(slide, x, 4.65, 1.75, 0.48, time, BLUE, BLUE, 11, True, WHITE)
        add_box(slide, x, 5.18, 1.75, 0.8, work, LIGHT, BLUE, 9, True)
    add_focus(slide, "后续重点：扩大评测范围，压实泛化边界。", y=6.25, size=14)


def slide_conclusion_merged(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "预期创新点与总结")
    items = [
        ("SQL+ 中间表示", "面向生成、转换、诊断和局部修复共同设计。", BLUE),
        ("多智能体反馈修正框架", "Critic、Router、Repair Skill 和 Executor 形成闭环。", ORANGE),
        ("repairability 评价视角", "不只看最终准确率，还看定位、路由、patch 范围和成本。", GREEN),
        ("可复现实验体系", "脚本、日志、开题材料和后续公开子集评测同步推进。", RED),
    ]
    for i, (head, body, color) in enumerate(items):
        x = 0.75 + (i % 2) * 4.35
        y = 1.25 + (i // 2) * 1.55
        add_box(slide, x, y, 3.75, 0.98, "", LIGHT, color)
        add_label(slide, x + 0.2, y + 0.14, 3.1, 0.26, head, 13, True, color)
        add_label(slide, x + 0.2, y + 0.48, 3.1, 0.32, body, 9, False, DARK)
    add_focus(slide, "总结：SQL+ 的价值在于把查询计划拆成可检查、可路由、可局部修复的步骤。", y=5.08, size=14)


def slide_motivation_scenario(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "动机测试一：实验背景与设置")
    rows = [
        ["项目", "设置"],
        ["实验背景", "第一次生成的 SQL 经常可执行但语义不对，需要判断问题来自表示形式、生成方式还是修复方式"],
        ["核心问题", "为什么设计 SQL+，为什么不能只靠单 Agent 或整条 SQL 重写"],
        ["实验环境", "SQLite 内存数据库，自建订单分析数据集，模型为 gpt-5-mini"],
        ["样例规模", "30 条自然语言查询；13 条 SQL+ prompt v2 已知失败样例；14 条 Direct SQL 失败样例"],
        ["对比对象", "Direct SQL、NL2SQL+、SQL+、NatSQL-style proxy、SemQL-style proxy、Direct SQL Refiner、SQL+ Refiner、Router + Repair Skills"],
        ["评估指标", "execution match、valid rate、token、latency、alias dependency、cross-clause reference、repair success、patch minimality"],
    ]
    add_table(slide, rows, y=1.02, h=5.05, size=7)
    add_focus(slide, "动机测试先证明设计必要性，再进入系统架构。", y=6.18, size=14)


def slide_motivation_process(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "动机测试二：实验过程")
    rows = [
        ["步骤", "具体做法", "回答的问题"],
        ["1 初次生成对比", "同一批 30 条问题和同一 schema，分别生成 Direct SQL 与 SQL+，转换后执行并对齐 gold 结果", "只换成 SQL+ 是否就能明显提升准确率"],
        ["2 表示复杂度对比", "把同一批 gold 查询整理成 SQL、SQL+、SemQL-style、NatSQL-style、Pipe-style，统计 token、别名依赖和跨子句引用", "SQL+ 的优势是不是来自更短表达"],
        ["3 生成成本对比", "用同一模型生成多种目标表示，记录 valid rate、execution match、总 token 和平均延迟", "SQL+ 是否有额外生成成本"],
        ["4 修复策略对比", "收集失败输出，分别交给 SQL+ Refiner、Direct SQL Refiner、Critic-Refiner 和 Router + Repair Skills 修复", "多 Agent 与局部 repair skill 是否必要"],
    ]
    add_table(slide, rows, y=1.05, h=4.85, size=7)
    add_focus(slide, "实验过程围绕一个问题展开：SQL+ 和多 Agent 是否真的让错误更容易被修。", y=6.12, size=13)


def slide_motivation_test(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "动机测试三：结果表一，为什么设计 SQL+")
    rows = [
        ["测试项", "关键结果", "结果分析"],
        ["Baseline", "Direct SQL 16/30；NL2SQL+ prompt v2 17/30", "只换输出格式提升很小，SQL+ 不能只靠初次生成准确率立论"],
        ["IR 复杂度", "SQL+ 平均 token 35.0333，高于 Standard SQL 31.5333", "SQL+ 不是更短的表示，不能用压缩长度解释价值"],
        ["结构依赖", "SQL+ alias dependency 0.7，cross-clause reference 1.0；SQL 分别为 2.0333 和 2.3333", "SQL+ 的优势在于步骤边界更清楚，跨区域耦合更少"],
        ["生成成本", "SQL+ execution match 14/30，平均 813.0333 tokens，9.2197s", "SQL+ 有成本，必须在修复阶段体现收益"],
    ]
    add_table(slide, rows, y=1.05, h=4.75, size=8)
    add_focus(slide, "结论：SQL+ 不是为了更短，而是为了把错误落到可定位、可修改的步骤上。", y=6.06, size=13)


def slide_motivation_agents(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "动机测试四：结果表二，为什么使用多 Agent")
    rows = [
        ["方法", "样例", "修复结果", "结果分析"],
        ["SQL+ non-gold Refiner v2", "13 条 SQL+ 已知失败样例", "4/13", "粗粒度反馈不足，模型容易改不到关键步骤"],
        ["Direct SQL non-gold Refiner", "14 条 Direct SQL 失败样例", "6/14", "整条 SQL 修复能工作，但修改范围不稳定，解释性较弱"],
        ["Schema-Critic-Refiner", "13 条 SQL+ 已知失败样例", "3/13", "先诊断再重写不一定有效，诊断没有转化为明确 repair 动作"],
        ["Step-wise Critic-Refiner", "13 条 SQL+ 已知失败样例", "3/13", "步骤级诊断更细，但没有限制 patch 范围时仍不稳定"],
        ["Skill Router + Repair Skills v3", "13 条 known-failure set", "13/13", "按错误类型路由到局部 skill 后，修复范围更可控"],
    ]
    add_table(slide, rows, y=1.0, h=4.98, size=7)
    add_focus(slide, "结论：Agent 数量不是重点，诊断、路由、局部 patch 和执行验证才是关键。", y=6.16, size=13)


def slide_motivation_requirements(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "动机测试五：结果分析与设计结论")
    rows = [
        ["实验观察", "推出的设计要求"],
        ["初次生成阶段，SQL+ 只比 Direct SQL 略高，且 token 和 latency 更高", "不能把 SQL+ 设计成单纯的生成格式，必须服务于反馈修复"],
        ["SQL+ 的别名依赖和跨子句引用更少", "需要把查询拆成步骤，让错误定位能落到 FROM、JOIN、WHERE、AGG、ORDER 等局部"],
        ["单 Refiner 与简单 Critic 串联效果不稳", "需要明确分工，Critic 负责诊断，Router 负责分流，Repair Skill 负责局部 patch"],
        ["分治 repair skill 在当前 known-failure set 上达到 13/13", "需要保留 skill 化修复，并用 Executor 对候选 patch 做执行验证"],
        ["Spider fresh e2e 小子集可跑通但规模有限", "后续必须扩展多库、多难度和更多错误类型，不能过度外推"],
    ]
    add_table(slide, rows, y=1.1, h=4.8, size=8)
    add_focus(slide, "因此，本课题设计的是 SQL+ 支撑下的多智能体反馈修正闭环。", y=6.15, size=14)


def slide_experiment_design(prs: Presentation, idx: int) -> None:
    slide = blank_slide(prs, idx, "实验结构设计：动机、可行性、对比、泛化")
    rows = [
        ["实验层次", "对应实验", "回答的问题"],
        ["动机测试", "baseline、IR 复杂度、单 Refiner 对比", "为什么不能只靠直接生成或整体重写"],
        ["方案可行性", "SQL+ parser、translator、executor 闭环", "SQL+ 是否能稳定解析、转换和执行"],
        ["机制可行性", "Critic -> Router -> Repair Skill -> Executor", "局部 skill 分治是否能修复已知错误类型"],
        ["公开子集验证", "Spider concert_singer 小规模 fresh e2e 与 semantic repair", "能否离开自建数据集初步跑通"],
        ["后续对比与消融", "多表示、多修复策略、去模块实验", "各模块是否真的贡献效果"],
    ]
    add_table(slide, rows, y=1.15, h=4.7, size=8)
    add_focus(slide, "结构原则：动机测试说明为什么做，可行性实验说明能不能做。", y=6.08, size=14)


def build() -> None:
    prs = Presentation(str(TEMPLATE))
    clear_template_slides(prs)

    title_slide(prs)
    slide_outline(prs, 2)
    slide_background_problem_merged(prs, 3)
    slide_timeline(prs, 4)
    slide_related_tech_merged(prs, 5)
    slide_motivation_scenario(prs, 6)
    slide_motivation_process(prs, 7)
    slide_motivation_test(prs, 8)
    slide_motivation_agents(prs, 9)
    slide_motivation_requirements(prs, 10)
    slide_gap_summary(prs, 11)
    slide_system_overview(prs, 12)
    slide_sqlplus_example(prs, 13)
    slide_architecture_detail(prs, 14)
    slide_experiment_design(prs, 15)
    slide_experiment_overview(prs, 16)
    slide_exp_conversion(prs, 17)
    slide_exp_repair(prs, 18)
    slide_exp_skills(prs, 19)
    slide_exp_spider(prs, 20)
    slide_feasibility_summary(prs, 21)
    slide_future_metrics_merged(prs, 22)
    slide_limits_schedule_merged(prs, 23)
    slide_conclusion_merged(prs, 24)
    slide_refs(prs, 25)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
