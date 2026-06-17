from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "assets" / "opening_template.pptx"
OUT = ROOT / "docs" / "opening" / "opening_ppt_template_version_v2.pptx"

FONT = "Microsoft YaHei"
TITLE = "面向 SQL+ 简化查询表达的多智能体自然语言数据库查询生成与反馈修正方法研究"

PRIMARY = RGBColor(31, 78, 121)
ACCENT = RGBColor(198, 89, 17)
GREEN = RGBColor(84, 130, 53)
RED = RGBColor(170, 60, 60)
DARK = RGBColor(40, 40, 40)
MUTED = RGBColor(95, 95, 95)
LIGHT = RGBColor(238, 242, 247)
PALE_BLUE = RGBColor(221, 235, 247)
PALE_ORANGE = RGBColor(251, 229, 214)
PALE_GREEN = RGBColor(226, 239, 218)


def clear_template_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_ids.remove(slide_id)


def set_run(run, size: int, bold: bool = False, color=DARK) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def put_text(shape, text: str, size: int = 18, bold: bool = False, color=DARK, align=PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)


def add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.48), Inches(0.32), Inches(9.1), Inches(0.55))
    put_text(box, title, 25, True, PRIMARY)


def add_footer(slide, idx: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.03), Inches(8.9), Inches(0.22))
    put_text(box, f"开题汇报 | SQL+ / Text-to-SQL / Multi-Agent | {idx}", 8, False, MUTED)


def add_box(slide, x, y, w, h, text, fill=LIGHT, line=PRIMARY, size=12, bold=False, color=DARK):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    put_text(shape, text, size, bold, color, PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x, y, text="→", size=20, color=PRIMARY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.35), Inches(0.25))
    put_text(box, text, size, True, color, PP_ALIGN.CENTER)


def add_bullets(slide, items: list[str], x=0.72, y=1.25, w=8.7, h=5.2, size=17, color=DARK) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)


def add_note(slide, text: str, y=6.42) -> None:
    box = slide.shapes.add_textbox(Inches(0.72), Inches(y), Inches(8.7), Inches(0.4))
    put_text(box, text, 11, False, ACCENT)


def add_table(slide, rows: list[list[str]], x=0.45, y=1.25, w=9.1, h=5.2, size=9) -> None:
    tbl_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = tbl_shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY if r == 0 else (LIGHT if r % 2 == 0 else RGBColor(255, 255, 255))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    set_run(run, size if r else size + 1, r == 0, RGBColor(255, 255, 255) if r == 0 else DARK)


def slide_blank(prs: Presentation, idx: int, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, title)
    add_footer(slide, idx)
    return slide


def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        put_text(slide.shapes.title, TITLE, 28, True, PRIMARY, PP_ALIGN.CENTER)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            put_text(
                shape,
                "研究方向：达梦 SQL+ / Text-to-SQL / 多智能体 / 执行反馈修正\n汇报人：待填写    导师：待填写    时间：待填写",
                17,
                False,
                DARK,
                PP_ALIGN.CENTER,
            )
    add_footer(slide, 1)


def content_slide(prs: Presentation, idx: int, title: str, bullets: list[str], note: str = "") -> None:
    slide = slide_blank(prs, idx, title)
    add_bullets(slide, bullets)
    if note:
        add_note(slide, note)


def direct_architecture_slide(prs: Presentation, idx: int) -> None:
    slide = slide_blank(prs, idx, "传统 Text-to-SQL 流程及其问题")
    labels = ["自然语言问题", "LLM / 解析器", "标准 SQL", "数据库执行", "查询结果"]
    xs = [0.45, 2.15, 3.85, 5.55, 7.25]
    fills = [PALE_BLUE, PALE_BLUE, PALE_ORANGE, PALE_GREEN, PALE_GREEN]
    for i, label in enumerate(labels):
        add_box(slide, xs[i], 1.55, 1.35, 0.75, label, fills[i], PRIMARY, 12, True)
        if i < len(labels) - 1:
            add_arrow(slide, xs[i] + 1.42, 1.78)
    add_box(slide, 2.55, 3.0, 4.9, 0.75, "错误反馈通常只作用在整条 SQL 上", PALE_ORANGE, ACCENT, 13, True)
    add_arrow(slide, 5.7, 2.47, "↓", 22, ACCENT)
    add_bullets(
        slide,
        [
            "标准 SQL 的 SELECT, FROM, WHERE, GROUP BY, HAVING, ORDER BY 相互耦合。",
            "复杂查询中的 join path、聚合口径、别名引用和排序条件容易互相影响。",
            "报错后直接重写整条 SQL，可能修好一个错误，又破坏原本正确的部分。",
        ],
        x=0.85,
        y=4.05,
        size=15,
    )


def proposed_architecture_slide(prs: Presentation, idx: int) -> None:
    slide = slide_blank(prs, idx, "本文拟研究的 SQL+ 多智能体闭环")
    top = [("Intent Agent", "意图"), ("Schema Agent", "表字段"), ("Planner Agent", "步骤计划"), ("SQL+ Generator", "SQL+")]
    for i, (name, sub) in enumerate(top):
        add_box(slide, 0.55 + i * 2.15, 1.2, 1.7, 0.75, f"{name}\n{sub}", PALE_BLUE, PRIMARY, 10, True)
        if i < 3:
            add_arrow(slide, 2.28 + i * 2.15, 1.42)
    mid = [("Translator", "SQL+ → SQL"), ("Executor", "执行反馈"), ("Critic Agent", "错误定位"), ("Skill Router", "技能路由")]
    for i, (name, sub) in enumerate(mid):
        add_box(slide, 0.55 + i * 2.15, 2.75, 1.7, 0.75, f"{name}\n{sub}", PALE_GREEN, GREEN, 10, True)
        if i < 3:
            add_arrow(slide, 2.28 + i * 2.15, 2.97)
    add_box(slide, 2.2, 4.35, 2.0, 0.75, "Repair Skill\n局部 patch", PALE_ORANGE, ACCENT, 11, True)
    add_box(slide, 5.3, 4.35, 2.0, 0.75, "SQL+ Patch\n重新验证", PALE_ORANGE, ACCENT, 11, True)
    add_arrow(slide, 4.45, 4.55)
    add_arrow(slide, 7.35, 3.15, "↺", 24, ACCENT)
    add_note(slide, "关键点：不是把多个 prompt 串起来，而是让每个 Agent 产生可检查的中间产物。")


def sqlplus_steps_slide(prs: Presentation, idx: int) -> None:
    slide = slide_blank(prs, idx, "SQL+ 的表达逻辑：把复杂 SQL 拆成可修复步骤")
    steps = ["FROM", "JOIN", "WHERE", "GROUP", "AGG", "HAVING", "SELECT", "ORDER", "LIMIT"]
    for i, step in enumerate(steps):
        add_box(slide, 0.35 + i * 1.05, 1.55, 0.85, 0.55, step, PRIMARY if i % 2 == 0 else ACCENT, PRIMARY, 10, True, RGBColor(255, 255, 255))
        if i < len(steps) - 1:
            add_arrow(slide, 1.18 + i * 1.05, 1.68, "→", 13)
    rows = [
        ["SQL+ 步骤", "主要诊断对象", "对应 repair skill"],
        ["WHERE", "值链接、日期边界、隐含过滤", "value-linking skill"],
        ["JOIN", "连接路径、连接方向、冗余 join", "join repair skill"],
        ["GROUP / AGG", "聚合口径、别名、HAVING 引用", "aggregation repair skill"],
        ["SELECT", "输出列、别名、列顺序", "projection repair skill"],
        ["ORDER / LIMIT", "排序字段、方向、top-k", "order repair skill"],
    ]
    add_table(slide, rows, y=2.75, h=3.15, size=10)
    add_note(slide, "SQL+ 的价值不只在生成阶段，更在于给错误定位和局部修复提供稳定锚点。")


def experiment_logic_slide(prs: Presentation, idx: int) -> None:
    slide = slide_blank(prs, idx, "实验组织逻辑：围绕四个研究假设展开")
    items = [
        ("H1 表达可行", "SQL+ 能否稳定转换并执行", "SQL+ 转换实验\nSpider smoke test"),
        ("H2 表达差异", "SQL+ 与 SQL、SemQL、NatSQL 的区别", "IR 复杂度实验\nIR 生成成本实验"),
        ("H3 生成效果", "中间表示是否改善初次生成", "Direct SQL vs SQL+\n执行一致率"),
        ("H4 修复收益", "修复收益能否抵消额外开销", "repairability 指标\nSkill Router 实验"),
    ]
    for i, (h, q, e) in enumerate(items):
        y = 1.2 + i * 1.25
        add_box(slide, 0.65, y, 1.45, 0.72, h, PALE_BLUE, PRIMARY, 12, True)
        add_arrow(slide, 2.25, y + 0.22)
        add_box(slide, 2.65, y, 2.45, 0.72, q, LIGHT, PRIMARY, 11, False)
        add_arrow(slide, 5.25, y + 0.22)
        add_box(slide, 5.65, y, 3.25, 0.72, e, PALE_GREEN, GREEN, 10, True)


def exp_card_slide(prs: Presentation, idx: int, title: str, purpose: str, setup: str, metrics: str, result: str, conclusion: str) -> None:
    slide = slide_blank(prs, idx, title)
    cards = [
        ("实验目的", purpose, PALE_BLUE, PRIMARY),
        ("实验条件", setup, LIGHT, PRIMARY),
        ("评价指标", metrics, PALE_GREEN, GREEN),
        ("主要结果", result, PALE_ORANGE, ACCENT),
        ("当前结论", conclusion, RGBColor(245, 245, 245), DARK),
    ]
    positions = [(0.6, 1.2, 4.1, 1.1), (5.0, 1.2, 4.3, 1.1), (0.6, 2.65, 4.1, 1.1), (5.0, 2.65, 4.3, 1.1), (0.6, 4.15, 8.7, 1.55)]
    for (head, body, fill, line), (x, y, w, h) in zip(cards, positions):
        add_box(slide, x, y, w, h, f"{head}\n{body}", fill, line, 11, True if head in ["实验目的", "主要结果"] else False)


def build() -> None:
    prs = Presentation(str(TEMPLATE))
    clear_template_slides(prs)

    title_slide(prs)
    content_slide(
        prs,
        2,
        "汇报结构",
        [
            "第一部分：研究背景、问题来源和国内外研究现状。",
            "第二部分：为什么使用 SQL+，以及它和标准 SQL、SemQL、NatSQL、Pipe-style 表示的区别。",
            "第三部分：面向 SQL+ 的多智能体生成与反馈修正框架。",
            "第四部分：初步实验设计、实验条件、核心结果和当前结论。",
            "第五部分：后续实验计划、风险边界和预期创新点。",
        ],
    )
    direct_architecture_slide(prs, 3)
    content_slide(
        prs,
        4,
        "问题背景：Text-to-SQL 在真实场景中仍不稳定",
        [
            "企业数据库通常存在多表关系、字段别名、业务规则、时间口径和 SQL 方言差异。",
            "LLM 能生成看似合理的 SQL，但语义错误常常不会触发语法报错。",
            "复杂查询中，join、filter、aggregation、projection、order 之间存在强耦合。",
            "因此，研究重点应从“能否生成一条 SQL”转向“能否生成、验证、定位和修复”。",
        ],
    )
    content_slide(
        prs,
        5,
        "国内外研究现状：从直接生成到过程化生成",
        [
            "传统 Text-to-SQL 研究关注语义解析、schema linking、语法约束解码和跨数据库泛化。",
            "LLM 方法引入 few-shot prompting、任务分解、候选生成、执行反馈和 self-correction。",
            "BIRD、Spider 2.0 等新 benchmark 表明，真实任务还需要外部知识、方言文档和工作流理解。",
            "已有工作提供了生成和选择机制，但对“可定位、可路由、可局部修复”的中间层研究仍不足。",
        ],
    )
    content_slide(
        prs,
        6,
        "为什么不是直接采用 SemQL 或 NatSQL",
        [
            "SemQL 更偏语义解析树，结构约束强，利于消除部分 SQL 语法复杂性，但和 SQL 执行链路之间仍需要较重映射。",
            "NatSQL 更接近自然语言表达，降低部分 SQL 子查询和 join 的书写难度，但局部执行和错误步骤定位并不是其主要目标。",
            "GoogleSQL Pipe Syntax 证明 SQL 生态内部可以引入线性数据流表达，但它是特定 SQL 方言扩展。",
            "本课题中的 SQL+ 定位为 NL2SQL 中间表示：保留 SQL 可转换性，同时显式服务于错误定位和 repair skill 路由。",
        ],
    )
    sqlplus_steps_slide(prs, 7)
    proposed_architecture_slide(prs, 8)
    content_slide(
        prs,
        9,
        "研究问题与评价维度",
        [
            "RQ1：SQL+ 是否能在不脱离 SQL 生态的前提下降低表达耦合。",
            "RQ2：多智能体是否能把意图理解、schema linking、SQL+ 生成、转换、执行和修复拆成可观察子任务。",
            "RQ3：执行反馈能否映射回 SQL+ 步骤，并触发局部 repair skill。",
            "RQ4：SQL+ 修复收益能否抵消生成阶段增加的 token 和 latency 成本。",
            "评价维度包括 execution accuracy、valid rate、token cost、latency、error localization accuracy、patch minimality 和 repair rounds。",
        ],
    )
    experiment_logic_slide(prs, 10)
    add_table(slide_blank(prs, 11, "实验设置总览"), [
        ["实验", "目的", "数据/样例", "对比对象", "核心指标"],
        ["SQL+ 转换", "验证中间表示可执行", "自建 30 条订单样例", "gold SQL", "转换成功、执行一致"],
        ["IR 复杂度", "解释为什么使用 SQL+", "同一批 30 条查询", "SQL、SQL+、SemQL-style、NatSQL-style、Pipe-style", "token、步骤、嵌套、别名依赖"],
        ["IR 生成成本", "比较初次生成代价", "同一批 30 条查询", "Direct SQL、SQL+、NatSQL-style、SemQL-style", "token、latency、valid rate、accuracy"],
        ["反馈修复", "验证局部修复价值", "SQL+ 已知失败集", "SQL Refiner、SQL+ Refiner、Skill Router", "repair success、patch minimality"],
        ["Spider smoke", "验证公开数据集迁移", "Spider dev concert_singer 20 条", "gold SQL", "SQL+ 有效、SQL 可执行、执行一致"],
    ], size=8)
    exp_card_slide(
        prs,
        12,
        "实验一：SQL+ 转换可行性",
        "验证 SQL+ 不是只停留在表示层，而是能稳定转换为可执行 SQL。",
        "自建订单分析数据集，30 条自然语言问题、SQL+、gold SQL 配对。",
        "SQL+ valid rate、SQL executable rate、execution match。",
        "30/30 SQL+ 可解析，30/30 SQL 可执行，30/30 与 gold SQL 执行一致。",
        "SQL+ 最小闭环成立。后续实验可以把问题转向生成质量和修复能力。",
    )
    exp_card_slide(
        prs,
        13,
        "实验二：IR 表达复杂度对比",
        "回答为什么使用 SQL+，而不是直接 SQL 或其它中间表示。",
        "同一批 30 条查询构造 Standard SQL、SQL+、SemQL-style、NatSQL-style、Pipe-style proxy。proxy 仅用于受控比较表达形态。",
        "平均 token、步骤/子句数、嵌套深度、别名依赖、跨子句引用、转换成功。",
        "SQL+ 平均 token 35.03，高于 SQL 31.53，但别名依赖 0.70、跨子句引用 1.00，低于标准 SQL 的 2.03 和 2.33。",
        "SQL+ 不应被表述为更短，而应表述为依赖更少、步骤更清晰，更适合定位和局部修复。",
    )
    add_table(slide_blank(prs, 14, "实验二结果：IR 表达复杂度"), [
        ["表示形式", "平均 token", "步骤/子句", "嵌套深度", "别名依赖", "跨子句引用", "转换成功"],
        ["Standard SQL", "31.5333", "5.9", "0.6667", "2.0333", "2.3333", "30/30"],
        ["SQL+", "35.0333", "6.1333", "0.6667", "0.7", "1.0", "30/30"],
        ["SemQL-style proxy", "50.5667", "10.7333", "3.6667", "0.9", "1.2", "N/A"],
        ["NatSQL-style proxy", "31.5", "5.4333", "0.9667", "1.3667", "1.6667", "N/A"],
        ["Pipe-style proxy", "40.8", "6.1333", "0.6667", "1.3667", "1.6667", "N/A"],
    ], size=9)
    exp_card_slide(
        prs,
        15,
        "实验三：IR 生成成本与执行效果",
        "判断 SQL+ 初次生成是否更省成本、更准确。",
        "同一批 30 条查询，统一模型和 prompt 框架，比较 Direct SQL、SQL+、NatSQL-style、SemQL-style。",
        "valid rate、SQL executable rate、execution accuracy、平均 token、latency。",
        "SQL+ 执行一致 14/30，Direct SQL 为 12/30；SQL+ 平均 token 813.03，latency 9.22s，高于 Direct SQL 的 599.17 和 6.59s。",
        "SQL+ 在初次生成阶段没有明显成本优势，研究重点应放在修复收益是否能抵消额外开销。",
    )
    add_table(slide_blank(prs, 16, "实验三结果：生成成本与执行效果"), [
        ["方法", "表示有效", "SQL 可执行", "执行一致", "平均总 token", "平均延迟"],
        ["Direct SQL", "30/30", "30/30", "12/30", "599.1667", "6.5851s"],
        ["SQL+", "28/30", "28/30", "14/30", "813.0333", "9.2197s"],
        ["NatSQL-style proxy", "30/30", "30/30", "13/30", "740.7667", "6.2802s"],
        ["SemQL-style proxy", "30/30", "25/30", "12/30", "1028.9667", "9.9684s"],
    ], size=10)
    exp_card_slide(
        prs,
        17,
        "实验四：Direct SQL 与 NL2SQL+ baseline",
        "观察单 Agent 条件下 SQL+ 是否改善初次生成。",
        "自建 30 条订单查询样例，使用 Direct NL2SQL prompt 与 NL2SQL+ prompt v2。",
        "SQL 可执行率、执行一致率、失败类型。",
        "Direct NL2SQL 为 16/30，NL2SQL+ prompt v2 为 17/30。",
        "仅替换输出形式收益有限。后续必须引入 schema、critic、router 和 repair skill，而不是只优化 prompt。",
    )
    exp_card_slide(
        prs,
        18,
        "实验五：反馈修复与 Skill Router",
        "验证 SQL+ 层局部修复是否优于整条 SQL 重写。",
        "当前 SQL+ 已知失败集 13 条，包含 value、ORDER、aggregation、join、projection 等错误。",
        "repair success、SQL+ valid、SQL executable、error localization、patch minimality。",
        "SQL+ 非 gold 单 Refiner 为 4/13，Direct SQL Refiner 为 6/14，SQL+ Skill Router + Repair Skills v3 为 13/13。",
        "多 Agent 简单串联不必然有效。错误分类、技能路由和局部 patch 是当前更有效的修复路线。",
    )
    add_table(slide_blank(prs, 19, "实验五结果：反馈修复"), [
        ["方法", "样例", "SQL+ 有效", "SQL 可执行", "修复成功"],
        ["SQL+ 非 gold 单 Refiner v2", "13", "13/13", "12/13", "4/13"],
        ["Direct SQL 非 gold Refiner", "14", "-", "14/14", "6/14"],
        ["SQL+ Schema-Critic-Refiner", "13", "13/13", "13/13", "3/13"],
        ["SQL+ Step-wise Critic-Refiner", "13", "13/13", "12/13", "3/13"],
        ["SQL+ Skill Router + Repair Skills v3", "13", "13/13", "13/13", "13/13"],
    ], size=9)
    add_table(slide_blank(prs, 20, "实验六：repairability 指标"), [
        ["方法", "样例", "修复成功", "定位准确率", "最小 patch 率", "平均 repair token"],
        ["Direct SQL Refiner", "14", "6/14", "0.8571", "0.8571", "1609.3571"],
        ["SQL+ Critic Router Skills", "13", "13/13", "0.7692", "0.9231", "3813.9231"],
        ["Direct SQL overlap", "9", "4/9", "0.8889", "0.8889", "1583.2222"],
        ["SQL+ Router overlap", "9", "9/9", "0.7778", "0.8889", "4001.7778"],
    ], size=10)
    add_note(prs.slides[-1], "解释：SQL+ 的修复成功率更高，patch 更可控，但 Critic 与 Router 引入了更高 token 成本。")
    add_table(slide_blank(prs, 21, "实验七：repair skill 分治结果"), [
        ["Repair Skill", "样例数", "修复成功", "覆盖问题"],
        ["value-linking", "3", "3/3", "值、日期边界、状态过滤"],
        ["ORDER", "3", "3/3", "排序字段、方向、LIMIT"],
        ["aggregation", "3", "3/3", "聚合口径、AGG 别名、HAVING/ORDER 引用"],
        ["join", "3", "3/3", "join 路径、冗余 join、缺少 paid 过滤、连接方向"],
        ["projection", "1", "1/1", "输出列多、少或顺序错误"],
    ], size=10)
    exp_card_slide(
        prs,
        22,
        "实验八：Spider smoke test",
        "验证 SQL+ 转换链路能否迁移到公开 Text-to-SQL 数据集子集。",
        "Spider dev 中 concert_singer 数据库的 20 条受支持样例。",
        "SQL+ valid、SQL executable、execution match。",
        "20/20 SQL+ 有效，20/20 SQL 可执行，20/20 与 gold SQL 执行一致。",
        "该结果只说明小规模公开子集可迁移，不能作为完整 Spider benchmark 分数。",
    )
    content_slide(
        prs,
        23,
        "初步实验的综合判断",
        [
            "SQL+ 的基本转换链路已经可用，具备继续研究的工程基础。",
            "SQL+ 初次生成并不天然更便宜，也不应被包装成“显著提高准确率”。",
            "SQL+ 的阶段性价值主要体现在错误定位、局部 patch 和 repair skill 路由。",
            "当前最需要补齐的是更大规模的非 gold 诊断实验、端到端 latency、token 成本和达梦 SQL 方言适配。",
        ],
    )
    content_slide(
        prs,
        24,
        "后续实验计划",
        [
            "扩展 Spider 子集：从单数据库 smoke test 扩展到多数据库、多难度、多类型查询。",
            "补齐端到端修复成本：重新记录 Critic、Router、Repair、Executor 各阶段 latency 和 token。",
            "做无 gold 语义错诊断：不依赖 gold SQL，观察 Critic Agent 能否定位结果为空、聚合口径错、projection 错。",
            "适配达梦 SQL：补充日期函数、分页、类型转换、字符串函数和方言兼容测试。",
            "与 SemQL/NatSQL 的对比保持边界：开题阶段先做 controlled proxy，后续若时间允许再接入原系统代码。",
        ],
    )
    content_slide(
        prs,
        25,
        "预期创新点",
        [
            "面向生成和修复的 SQL+ 中间表示：重点不是更短，而是步骤化、可转换、可定位、可 patch。",
            "面向 SQL+ 的多智能体反馈闭环：把 schema、planning、translation、critic、router、repair skill 和 executor 组织成可检查流程。",
            "面向 repairability 的评价体系：在 execution accuracy 之外，引入错误定位、技能路由、最小 patch、修复轮数、token 和 latency。",
        ],
    )
    content_slide(
        prs,
        26,
        "当前风险与应对",
        [
            "样例规模偏小：后续用 Spider 多数据库子集和 BIRD 子集补充外部有效性。",
            "proxy 对比不能替代完整复现：PPT 和开题报告中明确标注 controlled proxy，避免夸大。",
            "Skill Router 当前依赖已知失败集：后续要增加自动诊断和未知错误样例。",
            "SQL+ 成本更高：后续实验重点比较修复收益，而不是单纯比较初次生成。",
            "达梦方言尚未充分验证：后续增加方言转换层和兼容性测试集。",
        ],
    )
    content_slide(
        prs,
        27,
        "总结",
        [
            "本课题不是单纯做 Text-to-SQL demo，而是研究 SQL+ 中间表示与多智能体反馈修正机制。",
            "已有实验支持 SQL+ 转换可行、公开小子集可迁移、局部 repair skill 对已知失败集有效。",
            "已有实验也提示 SQL+ 不应被简单表述为更省 token 或初次生成更准。",
            "后续研究重点是扩大样例、强化无 gold 诊断、补齐成本指标，并验证 SQL+ 修复收益是否稳定存在。",
        ],
    )
    content_slide(
        prs,
        28,
        "主要参考文献",
        [
            "Spider, BIRD, Spider 2.0：Text-to-SQL benchmark 与真实场景评估。",
            "IRNet/SemQL, NatSQL：Text-to-SQL 中间表示研究。",
            "RAT-SQL, PICARD, RESDSQL：schema linking、语法约束与结构化解码。",
            "DAIL-SQL, DIN-SQL, CHESS, CHASE-SQL：LLM-based Text-to-SQL、任务分解与候选选择。",
            "MAC-SQL, SQLCritic, Tool-Assisted Agent, LEVER：多智能体、执行反馈和修复评估。",
            "GoogleSQL Pipe Syntax：SQL 线性表达与 SQL 语言扩展。",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
