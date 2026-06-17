from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
PPTS = [
    ROOT / "docs" / "opening" / "opening_ppt_template_version_v2.pptx",
    ROOT / "docs" / "opening" / "opening_ppt_template_version_v3.pptx",
]


NOTES = [
    """【汇报讲稿】
这一页先说明课题定位。我的题目聚焦 SQL+、Text-to-SQL、多智能体和执行反馈修正。这里需要强调，本课题不是做一个单轮自然语言转 SQL 的工程 demo，而是研究如何通过一个更适合生成和修复的中间表示，把查询生成、SQL 转换、执行验证和错误修正连接成闭环。课题背景一方面来自达梦 SQL+ 的扩展语法需求，另一方面来自真实数据库自然语言查询中“能生成但难验证、难定位、难修复”的问题。

【答辩备注】
如果老师问“SQL+ 是不是重新发明 SQL”，回答：不是。SQL+ 在本课题中是 NL2SQL 的中间表示层，目标是降低表达耦合和提高可修复性，最终仍要转换成标准 SQL 或达梦 SQL 执行。""",
    """【汇报讲稿】
这一页介绍汇报结构。我会按五个部分展开：第一，问题来源和研究背景；第二，已有研究和研究空间；第三，本文方法，包括 SQL+ 和多智能体闭环；第四，初步实验，重点说明实验目的、数据集、评价指标和当前结果；第五，后续计划、风险边界和预期创新点。这个顺序的目的，是先说明为什么需要这个课题，再说明我准备怎么做，最后用初步实验说明它是可落地的。

【答辩备注】
这一页不要展开讲技术细节，控制在 30 秒左右。重点是让老师知道实验部分不是单独堆结果，而是围绕研究问题设计。""",
    """【汇报讲稿】
这一页讲传统 Text-to-SQL 的基本流程。自然语言问题经过模型或解析器生成标准 SQL，再交给数据库执行，最后返回查询结果。这个流程很直接，但问题是反馈通常只落在最终 SQL 上。SQL 一旦出错，很难判断是表选错、字段选错、join 错、聚合错，还是排序和过滤条件错。复杂查询不是一个简单的句子生成问题，而是多个相互耦合的查询决策问题。

【答辩备注】
这一页是后面引出 SQL+ 的铺垫。可以补充一个例子：同样是“销售额最高的客户”，错误可能来自 paid 过滤缺失，也可能来自 SUM 口径错误，还可能来自 ORDER BY 字段错误。最终 SQL 层只能看到结果错，不容易定位步骤。""",
    """【汇报讲稿】
这一页补充真实场景中的困难。企业数据库中常见多表关系、字段别名、时间口径和 SQL 方言差异。大模型生成的 SQL 即使能执行，也可能在语义上不符合用户意图，比如统计口径错、过滤条件漏掉、join 路径不合理。这些错误不会一定报语法错。因此本课题关注的不只是“生成 SQL”，而是生成之后能不能执行、诊断和修复。

【答辩备注】
这一页可以强调四类难点：schema 复杂、方言差异、可执行但语义错误、整条重写修复不稳定。它们对应后续 SQL+ 的四个价值：步骤化、可转换、可定位、可局部 patch。""",
    """【汇报讲稿】
这一页概括国内外研究现状。传统方法强调语义解析、schema linking 和语法约束；LLM 方法进一步引入 few-shot、任务分解和执行反馈；Spider、BIRD、Spider 2.0 等 benchmark 说明任务越来越接近真实数据分析场景。这里的“当前缺口”不是说已有工作没有价值，而是说已有工作多数仍围绕最终 SQL 进行生成、选择或修复，对中间查询步骤级错误定位、repair skill 路由和局部修复研究还不充分。

【答辩备注】
“已有工作已证明过程化生成和执行反馈有效”可以具体指 DIN-SQL、MAC-SQL、CHASE-SQL、SQLCritic 等。DIN-SQL 支持任务分解，MAC-SQL 支持多智能体协作，CHASE-SQL 支持候选生成与选择，SQLCritic 支持分子句 critic 和反馈修正。""",
    """【汇报讲稿】
这一页回答为什么不直接采用 SemQL、NatSQL 或 Pipe Syntax。SemQL 更像语义解析树，结构约束强；NatSQL 更接近自然语言化 SQL，降低部分 SQL 书写难度；GoogleSQL Pipe Syntax 是 SQL 生态中的线性数据流扩展。但本课题需要的是面向 NL2SQL 生成、执行反馈和局部修复的中间表示。SQL+ 的设计重点不是最短表达，而是保留 SQL 可转换性，同时让错误可以映射到具体查询步骤上。

【答辩备注】
如果老师问“为什么不复现 SemQL/NatSQL”，回答：开题阶段先做 controlled proxy 对比，用来说明表达形态差异，不声称完整复现原系统。后续如果时间允许，可以接入原系统代码做更严格对比。""",
    """【汇报讲稿】
这一页说明 SQL+ 的基本表达逻辑。SQL+ 按 FROM、JOIN、WHERE、GROUP、AGG、HAVING、SELECT、ORDER、LIMIT 组织查询步骤。每个步骤都对应一类常见错误和修复技能。例如 WHERE 对应值链接和过滤条件错误，JOIN 对应连接路径错误，AGG 对应聚合口径和别名引用错误，SELECT 对应投影列错误，ORDER 对应排序字段和 top-k 错误。

【答辩备注】
这一页要强调 SQL+ 的步骤不是展示标签，而是后续 Critic Agent 定位、Skill Router 路由和 Repair Skill 局部 patch 的操作单元。""",
    """【汇报讲稿】
这一页讲本文拟研究的多智能体闭环。前半部分是生成链路：用户问题经过 Intent Agent、Schema Agent 和 Planner Agent，生成 SQL+ 表达。随后 Translator 将 SQL+ 转换成 SQL，Executor 执行 SQL。如果执行失败或者结果异常，反馈进入 Critic Agent，Critic 输出错误类型和可疑 SQL+ 步骤，Skill Router 再把错误路由给对应的 Repair Skill。Repair Skill 只修改 SQL+ 的局部步骤，形成 SQL+ Patch，然后回到 Translator 和 Executor 重新转换、重新执行。

【答辩备注】
新版图的关键是循环位置：`Repair Skill -> SQL+ Patch -> Translator -> Executor`。这里不能画成 Router 到 Repair 后结束。要说明本方法不是多个 prompt 串联，而是每个环节都有可检查中间产物。""",
    """【汇报讲稿】
这一页提出研究问题和评价维度。RQ1 关注 SQL+ 是否降低表达耦合；RQ2 关注多智能体是否能拆分成可观察子任务；RQ3 关注执行反馈能否映射到 SQL+ 步骤；RQ4 关注修复收益能否抵消额外成本。评价不只看 execution accuracy，还包括 valid rate、token、latency、error localization accuracy、patch minimality 和 repair rounds。

【答辩备注】
这页回应导师反馈：研究内容不能只写“先做 parser，再做 agent”，而要写可验证问题、技术难点和评价指标。""",
    """【汇报讲稿】
这一页把实验和研究假设对应起来。H1 验证 SQL+ 表达和转换是否可行；H2 比较 SQL+ 与标准 SQL、SemQL-style、NatSQL-style、Pipe-style 的表达差异；H3 看初次生成是否有收益；H4 看修复收益是否成立。后面的实验不是单独堆结果，而是逐步回答为什么要 SQL+、它的代价是什么，以及它是否能在修复阶段体现价值。

【答辩备注】
如果老师问“为什么实验这么多”，回答：因为本课题不是只优化 accuracy，而是同时研究表达复杂度、生成成本、执行一致性和 repairability。""",
    """【汇报讲稿】
这一页是实验设置总览。当前实验主要使用三类数据：第一，自建订单分析数据集，共 30 条查询样例，每条包含自然语言问题、gold SQL 和 gold SQL+；第二，SQL+ 已知失败集，来自 SQL+ prompt v2 的失败样例，用于反馈修复实验；第三，Spider dev 中 concert_singer 数据库的 20 条受支持样例，用于公开 benchmark smoke test。SemQL-style、NatSQL-style 和 Pipe-style 在这里是受控 proxy，只用于比较表达形态，不代表完整复现原系统。

【答辩备注】
要主动说明数据边界：自建数据用于机制验证，已知失败集用于修复链路验证，Spider smoke test 用于迁移可行性，不是完整 Spider benchmark。""",
    """【汇报讲稿】
实验一验证 SQL+ 转换可行性。数据集是自建订单分析数据集，包含 30 条查询样例，覆盖单表过滤、多表 join、聚合、分组、排序和 top-k 等常见分析查询。实验条件是使用本地 SQL+ parser 和 SQL converter，将 gold SQL+ 转换成 SQL，并在 SQLite 内存数据库上执行。评价指标包括 SQL+ valid rate、SQL executable rate 和 execution match。SQL+ valid 由 parser 是否成功解析判断；SQL executable 由转换后的 SQL 是否能在 SQLite 上执行判断；execution match 由转换 SQL 和 gold SQL 的执行结果是否完全一致判断。当前结果是 30/30、30/30、30/30。

【答辩备注】
这项实验只证明 SQL+ 表达和转换链路可行，不证明模型生成一定准确。数据文件主要对应 `data/sqlplus_cases.jsonl`、`data/schema.sql` 和 `src/sqlplus.py`，运行逻辑在 SQL+ 转换实验脚本中。""",
    """【汇报讲稿】
实验二是 IR 表达复杂度对比，目的是回答“为什么使用 SQL+”。数据仍然是自建 30 条订单分析查询。对同一批查询构造五种表示：Standard SQL、SQL+、SemQL-style proxy、NatSQL-style proxy 和 Pipe-style proxy。这里的 proxy 是为了受控比较表达形式，不代表完整复现 SemQL、NatSQL 或 GoogleSQL Pipe Syntax。指标包括 token 数、步骤或子句数、嵌套深度、别名依赖和跨子句引用。结果显示 SQL+ 平均 token 为 35.0333，高于标准 SQL 的 31.5333，但别名依赖为 0.7，跨子句引用为 1.0，低于标准 SQL 的 2.0333 和 2.3333。

【答辩备注】
这页要强调结论不是“SQL+ 更短”，而是“SQL+ 依赖更少、步骤边界更清楚”。指标来源是 `scripts/sqlplus/run_ir_complexity_eval.py`，汇总文件是 `data/ir_complexity_summary.csv`。""",
    """【汇报讲稿】
这一页展示 IR 复杂度的具体结果。token_count 是把表示文本按标识符、操作符和标点切分后的数量；step_or_clause_count 是 SQL 子句数或 SQL+ 步骤数；nesting_depth 是括号嵌套深度；alias_dependency_count 是统计 AS 别名在后续区域被引用的次数；cross_clause_reference_count 是对跨子句引用的近似统计。SQL+ 的步骤数略高，说明它通过显式拆步换取可解释性；但 alias dependency 和 cross-clause reference 更低，说明它降低了后续修复时的跨区域耦合。

【答辩备注】
如果老师质疑指标是否严谨，可以说这些是开题阶段的结构复杂度 proxy，用于比较表达形态。后续会结合更多真实查询和修复指标验证。""",
    """【汇报讲稿】
实验三比较 IR 生成成本和执行效果。数据仍是同一批 30 条自建订单查询，模型使用 gpt-5-mini，四种生成目标分别是 Direct SQL、SQL+、NatSQL-style proxy 和 SemQL-style proxy。实验记录 OpenAI API 返回的 input_tokens、output_tokens、total_tokens，以及本地记录的 latency_seconds。生成结果先检查表示是否有效，再转换成 SQL，在同一 SQLite 数据库上执行，并与 gold SQL 的执行结果比较。结果显示 SQL+ execution match 为 14/30，Direct SQL 为 12/30；但 SQL+ 平均 total token 为 813.0333，平均延迟 9.2197 秒，高于 Direct SQL 的 599.1667 和 6.5851 秒。

【答辩备注】
这项实验说明 SQL+ 初次生成没有明显成本优势。开题中应谨慎表述：SQL+ 的价值需要在修复收益上继续验证，而不是说它初次生成显著更准或更省。""",
    """【汇报讲稿】
这一页展示生成成本表格。Valid repr 表示模型输出是否符合该表示的格式要求；Valid SQL 表示输出经过转换后能否在 SQLite 上执行；Exec match 表示生成 SQL 与 gold SQL 在同一数据库上的执行结果是否一致；平均 token 直接来自 OpenAI API usage 字段；平均 latency 来自脚本对 API 调用耗时的记录。SQL+ 比 Direct SQL 稍高 2 条 execution match，但 token 和 latency 更高。SemQL-style proxy 的平均 token 最高，且部分转换后 SQL 不可执行。

【答辩备注】
如果老师问为什么 SQL+ valid 只有 28/30，原因是当时模型生成了当前 SQL+ 子集暂不支持的 LEFT 步骤。这也说明 SQL+ 语法子集边界需要继续完善。""",
    """【汇报讲稿】
实验四是 Direct SQL 与 NL2SQL+ baseline。数据集仍是 30 条自建订单查询。Direct NL2SQL 让模型直接从自然语言生成 SQL；NL2SQL+ prompt v2 让模型先生成 SQL+，再由转换器转 SQL 执行。评价方式与前面一致，主要看 SQL 可执行率和执行结果是否与 gold SQL 一致。结果是 Direct NL2SQL 为 16/30，NL2SQL+ prompt v2 为 17/30。提升很小，说明仅靠把输出格式换成 SQL+ 不足以解决问题。

【答辩备注】
这一页是引出多智能体和 repair skill 的关键。结论要讲稳：SQL+ 不是单靠 prompt 就能显著提升初次生成准确率，它需要配合 schema linking、critic、router 和 repair skill 才能体现中间表示价值。""",
    """【汇报讲稿】
实验五比较反馈修复方法。数据来自 SQL+ prompt v2 的 13 条已知失败样例，以及 Direct SQL 的 14 条失败样例。失败样例主要包括 value-linking、ORDER、aggregation、join 和 projection 问题。对比方法包括 SQL+ 非 gold 单 Refiner、Direct SQL 非 gold Refiner、SQL+ Schema-Critic-Refiner、SQL+ Step-wise Critic-Refiner 和 SQL+ Skill Router + Repair Skills v3。修复成功的判定方式是修复后的 SQL 或 SQL+ 转 SQL 在 SQLite 上执行，并与 gold SQL 执行结果一致。结果显示单 Refiner 效果有限，而 Skill Router + Repair Skills v3 在当前 13 条已知失败集上达到 13/13。

【答辩备注】
必须主动说明：13/13 只限当前已知失败集，不是大规模 benchmark 结果，也不是完整自主诊断系统性能。""",
    """【汇报讲稿】
这一页展示反馈修复对比结果。可以看到，简单增加多个 agent 并不一定有效，Schema-Critic-Refiner 和 Step-wise Critic-Refiner 都只有 3/13。真正提升来自错误类型路由和局部 repair skill：先由 Critic 给出错误类型和可疑步骤，再由 Skill Router 调用对应技能，最后由 Executor 验证候选 patch。这个结果说明本课题不能只讲“多智能体”，还要讲“多智能体如何分工、如何使用中间表示进行局部修复”。

【答辩备注】
如果老师问为什么 Step-wise Critic 反而低，可以回答：更细诊断如果不和可执行修复技能结合，可能仍然误导 Refiner。后续方向是 Critic + Router + Skill，而不是单纯扩大 prompt。""",
    """【汇报讲稿】
实验六是 repairability 指标对比。这个实验基于已有修复输出进行离线统计。Direct SQL 使用单 Refiner 输出，SQL+ 使用 Critic Agent、Skill Router 和 Repair Skills 的输出。repair success 仍然由执行结果是否与 gold SQL 一致判断；localization accuracy 是看模型定位到的步骤是否与 gold 差异中的预期错误步骤有交集；patch minimality 是比较实际修改步骤与预期修改步骤的重合程度；repair rounds 对 SQL+ 是 routing plan 的长度；repair token 对 SQL+ 统计 Step-wise Critic Agent 的 token，用来估计诊断成本。结果显示 SQL+ 修复成功率更高，patch minimality 更好，但 repair token 成本也更高。

【答辩备注】
要特别说明 latency 边界：SQL+ 表中的 repair latency 只包含本地 deterministic router 和 repair skill 执行，不包含 Critic API latency；Direct SQL 早期运行没有记录 latency，所以标为 N/A。""",
    """【汇报讲稿】
实验七展示 repair skill 分治结果。当前已经实现五类 skill：value-linking、ORDER、aggregation、join 和 projection。每类 skill 都在对应小样例上验证：value-linking 3/3，ORDER 3/3，aggregation 3/3，join 3/3，projection 1/1。这里的数据集来自 SQL+ 已知失败集的错误类型切分。每个 skill 会生成或选择局部 patch，然后通过 SQL+ parser、SQL converter 和 SQLite executor 检查可执行性与结果一致性。

【答辩备注】
这一页的意义不是样例规模大，而是证明按错误类型拆修复策略是可行的。后续需要扩大每类错误样例，加入复合错误和无报错语义错。""",
    """【汇报讲稿】
实验八是 Spider smoke test。数据来自 Spider dev 中 concert_singer 数据库，当前选取 20 条 SQL+ 子集能够覆盖的样例。实验做法是将这些 gold SQL 改写或映射为 SQL+，再通过 SQL+ 转换器生成 SQL，在对应 SQLite 数据库上执行，并与 Spider gold SQL 的执行结果比较。指标包括 SQL+ valid、SQL executable 和 execution match，三者当前都是 20/20。这个结果说明 SQL+ 表达与转换机制在公开 benchmark 子集上有初步迁移可行性。

【答辩备注】
一定要说清楚：这不是完整 Spider benchmark 分数，只是 concert_singer 数据库中受支持小子集的 smoke test。后续需要扩展到多数据库、多难度和更多 SQL 结构。""",
    """【汇报讲稿】
这一页给出综合判断。已有实验支持三点：第一，SQL+ 转换链路已经跑通；第二，SQL+ 不一定更短，但具有较低的别名依赖和跨子句引用，更适合作为错误定位和局部 patch 的锚点；第三，Skill Router + Repair Skills 在当前已知失败集上明显优于单 Refiner。与此同时，也要承认 SQL+ 初次生成成本更高，样例规模仍小，端到端 latency 还需要补齐。

【答辩备注】
这页适合主动降低过强结论。不要说“SQL+ 已经证明有效”，而说“当前结果支持继续研究 SQL+ 的 repairability 价值”。""",
    """【汇报讲稿】
这一页讲后续实验计划。第一，扩展 Spider 子集，从单数据库扩展到多数据库、多难度、多查询类型。第二，补齐端到端成本，分别记录 Critic、Router、Repair、Executor 的 token 和 latency。第三，做无 gold 语义诊断，处理 SQL 能执行但结果为空、聚合口径错、projection 错等问题。第四，适配达梦 SQL 方言，包括日期函数、分页、类型转换和字符串函数。第五，如果时间允许，再接入 SemQL 或 NatSQL 原系统做更严格对比。

【答辩备注】
这些计划要和现有不足对应：规模不足、成本不足、诊断不足、方言不足、对比不足。""",
    """【汇报讲稿】
这一页讲预期创新点。第一是面向生成和修复的 SQL+ 中间表示，它不是更短的 SQL，而是步骤化、可转换、可定位和可 patch 的表示。第二是面向 SQL+ 的多智能体反馈闭环，把 Critic、Router、Repair Skill 和 Executor 组织成可检查流程。第三是面向 repairability 的评价体系，在 execution accuracy 之外加入错误定位、技能路由、最小 patch、修复轮数、token 和 latency。

【答辩备注】
表述要谨慎，不用“首创”“显著领先”。可以说“本课题尝试把 SQL+ 中间表示设计和反馈修正机制结合起来”。""",
    """【汇报讲稿】
这一页讲风险和应对。样例规模偏小，对应后续扩展 Spider 和 BIRD 子集；proxy 对比有限，对应后续视时间接入原系统；已知失败集偏窄，对应增加未知错误和无 gold 语义错；SQL+ 成本更高，对应继续比较修复收益与端到端成本；达梦方言不足，对应补充方言兼容测试。

【答辩备注】
老师通常会关心可行性和边界。这里要表现出你知道当前工作还处在开题阶段，已有实验是可行性证据，后续实验计划已经针对不足展开。""",
    """【汇报讲稿】
这一页总结。第一，本课题研究的不是单次生成 SQL，而是生成、执行、诊断和修复闭环。第二，当前实验说明 SQL+ 在转换可行性、步骤化表达和已知失败集局部修复上有初步证据。第三，后续要扩大公开数据集子集，补齐无 gold 诊断和端到端成本，并适配达梦 SQL 方言。

【答辩备注】
最后一句可以落到课题意义：通过 SQL+ 和多智能体反馈修正，提高自然语言数据库查询的可解释性、可修复性和工程可落地性。""",
    """【汇报讲稿】
这一页列参考文献。参考文献覆盖四类：第一，Spider、BIRD、Spider 2.0 等 benchmark；第二，SemQL、NatSQL、GoogleSQL Pipe Syntax 等中间表示或 SQL 扩展；第三，DAIL-SQL、DIN-SQL、CHESS、CHASE-SQL 等 LLM-based Text-to-SQL；第四，MAC-SQL、SQLCritic、Tool-Assisted Agent、LEVER 等多智能体、执行反馈和修复相关工作。

【答辩备注】
如果老师问具体支撑关系，可以这样回答：DIN-SQL 支撑任务分解，MAC-SQL 支撑多智能体协作，CHASE-SQL 支撑候选生成与选择，SQLCritic 支撑分子句诊断和反馈修正，Pipe Syntax 支撑 SQL 线性表达的合理性。""",
]


def main() -> None:
    for ppt in PPTS:
        if not ppt.exists():
            continue
        prs = Presentation(str(ppt))
        if len(prs.slides) != len(NOTES):
            raise RuntimeError(f"{ppt}: slide count mismatch: {len(prs.slides)} slides, {len(NOTES)} notes")
        for slide, note in zip(prs.slides, NOTES):
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.clear()
            paragraphs = [line.strip() for line in note.splitlines() if line.strip()]
            notes_tf.text = "\n".join(paragraphs)
        prs.save(ppt)
        print(f"wrote detailed notes to {ppt}")


if __name__ == "__main__":
    main()
