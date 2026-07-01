from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
SRC = ROOT / "docs" / "opening" / "opening_ppt_template_version_v5_charts.pptx"
OUT = ROOT / "docs" / "opening" / "opening_ppt_template_version_v6_long_notes.pptx"

NOTES: dict[int, str] = {
1: """【汇报讲稿】
各位老师好，我汇报的题目是“面向 SQL+ 简化查询表达的多智能体自然语言数据库查询生成与反馈修正方法研究”。这个题目表面上属于 Text-to-SQL，也就是把自然语言问题转换成数据库可以执行的 SQL，但我这里不是单纯做一个自然语言转 SQL 的工程 demo。我的研究重点是：当复杂查询生成失败之后，系统能不能知道错在哪里，能不能把错误定位到具体查询步骤，能不能用局部修复而不是整条 SQL 重写来完成反馈修正。
所以本课题的核心对象有三个：第一是 SQL+ 中间表示，它把传统 SQL 中 SELECT、FROM、WHERE、JOIN、GROUP、ORDER 等相互耦合的结构拆成更线性的步骤；第二是多智能体流程，包括 Intent、Schema、Planner、Generator、Translator、Executor、Critic、Router 和 Repair Skill；第三是执行反馈和局部修复机制，也就是把数据库返回的错误或结果异常反馈到 SQL+ 层，而不是只在最终 SQL 文本层重写。
这也是我和普通 Text-to-SQL demo 的区别：我关心的不是“模型一次能不能生成一条看起来正确的 SQL”，而是“生成、转换、执行、诊断、修复”这一整条链路是否可解释、可检查、可迭代。
【答辩备注】
可能问题：SQL+ 是不是要替代 SQL？回答：不是。SQL+ 在本课题里是中间表示，最终仍然要转换成标准 SQL 或达梦 SQL 执行。它的作用是服务于大模型生成、错误定位和局部修复。
可能问题：为什么题目里强调多智能体？回答：因为复杂查询错误通常不只是一类错误，意图理解、schema linking、步骤规划、转换执行和错误修复需要不同角色处理。多智能体不是为了堆概念，而是为了让每个阶段产出可检查的中间结果。""",

2: """【汇报讲稿】
这次汇报按照六个部分展开。第一部分是研究背景与问题来源，主要说明自然语言数据库查询为什么在真实场景里仍然困难。第二部分是国内外研究现状，我会从 benchmark、中间表示、SQL 扩展、LLM Text-to-SQL、多智能体和反馈修正几个方向梳理已有工作。第三部分是研究空间和本文切入点，说明为什么现有方法还不能充分解决 SQL 层错误定位和局部修复问题。第四部分是本文拟研究的方法，包括 SQL+ 表达设计、多智能体生成框架、Critic-Agent 诊断、Skill Router 路由和 Repair Skill 局部修复。第五部分是目前已经完成的初步实验，重点不是报一个很高的分数，而是说明这个研究方向的可行性和边界。第六部分是后续计划、风险控制和预期创新。
在讲实验时，我会特别说明每组实验回答的问题、数据集来源、变量设置、指标定义和结论边界。因为导师之前也提醒过，研究内容不能写成“先做 parser、再做 agent、再做系统”的工程步骤，而要体现技术难点、评价方法和对比对象，所以后面实验会围绕两个核心问题展开：为什么需要 SQL+，以及为什么需要多智能体加局部 repair skill。
【答辩备注】
这页不要展开技术细节，只把汇报路线交代清楚。答辩时如果老师觉得内容多，可以说明本次重点在“研究动机、研究空间和前期可行性实验”，完整系统实现会在开题后继续推进。""",

3: """【汇报讲稿】
这一页先讲问题背景。自然语言数据库查询的目标是让业务人员或非数据库专业用户直接用自然语言提问，例如“查询已支付订单中金额最高的十个订单”，系统自动生成 SQL 并返回结果。这个方向近几年因为大语言模型的发展有了明显进步，但真实数据库场景仍然不稳定。
主要困难有四类。第一是 schema linking，也就是模型要从很多表和字段中选出正确的表、字段和 join 路径。第二是 value linking，例如自然语言里说“已支付”，数据库里实际可能是 status='paid'，也可能是 payment_state='SUCCESS'，这个映射不一定直接。第三是复杂 SQL 结构，尤其是多表 join、group by、having、order by、top-k、子查询和窗口函数。第四是反馈修正困难，SQL 能执行不代表语义正确，很多错误不会报语法错误，而是返回一个看似正常但语义不对的结果。
所以本课题关注的不是生成一条 SQL 就结束，而是希望在生成后继续做执行验证、错误诊断和局部修复。传统做法通常让模型看到错误后重新生成整条 SQL，但这种方式不稳定，可能改坏原本正确的部分。SQL+ 的价值就在于把查询拆成可定位的步骤，让错误可以落到 WHERE、JOIN、AGG、ORDER、SELECT 等局部单元上。
【答辩备注】
可能问题：为什么“可执行 SQL”还不够？回答：因为很多语义错误不会导致数据库报错，例如漏掉 paid 过滤、聚合口径不对、排序字段不对，SQL 仍然能执行，但结果不符合用户意图。
可能问题：为什么要局部修复？回答：整条 SQL 重写容易引入新错误，局部修复可以保留正确部分，只修改可疑步骤，便于验证和解释。""",

4: """【汇报讲稿】
这一页梳理研究发展脉络。早期 Text-to-SQL 研究主要围绕语义解析、Seq2Seq、语法约束和 schema linking 展开。Spider 数据集推动了跨数据库复杂查询研究，它要求模型面对未见过的数据库 schema 生成 SQL。BIRD 和 Spider 2.0 进一步把任务推向真实场景，强调真实数据库、外部知识、执行效率、多 SQL 方言和企业级 workflow。
第二条线是中间表示和 SQL 扩展。SemQL 通过结构化语义表示降低 SQL 生成复杂度，NatSQL 试图用更接近自然语言的 SQL 形式减少生成难度，GoogleSQL Pipe Syntax 则说明 SQL 可以通过管道式表达变得更线性、更接近数据流。它们给本课题提供了一个重要启发：并不是所有任务都必须直接生成最终 SQL，中间表示可以降低复杂度并改善可解释性。
第三条线是 LLM-based Text-to-SQL 和 agentic NL2SQL。DAIL-SQL 关注 prompt 和样例组织，DIN-SQL 强调任务分解和自修正，MAC-SQL 使用多智能体协作，CHASE-SQL 使用多路径候选生成和偏好选择，SQLCritic 关注分子句 critic 和反馈修正。这些工作说明，过程化生成、执行反馈和多角色协作是有效方向。
本课题的研究空间在于：已有工作大多仍然围绕最终 SQL 的生成、选择或修正，而我希望把错误反馈压回 SQL+ 这种步骤化中间表示，在 SQL+ 层做错误定位、skill 路由和局部 patch。
【答辩备注】
可能问题：已有多智能体 Text-to-SQL 已经做了，本文区别是什么？回答：区别不在“也用了多个 agent”，而在于本文把 SQL+ 作为可修复中间层，让 Critic 输出错误类型和 SQL+ 步骤，Router 选择修复技能，Repair Skill 只改局部步骤。
可能问题：为什么提 Pipe Syntax？回答：Pipe Syntax 证明 SQL 线性化和数据流式表达有现实基础，但本文不是照搬 GoogleSQL，而是面向 NL2SQL 生成与反馈修正设计 SQL+。""",

5: """【汇报讲稿】
这一页对研究现状做一个归纳。可以看到，Text-to-SQL 的研究对象已经从早期的单句 SQL 生成，逐渐发展到真实数据库、多表 schema、多轮反馈、多模型协作和执行验证。也就是说，当前任务越来越接近真实数据分析流程，系统也越来越复杂。
但是从我的课题角度看，仍然存在一个没有充分展开的问题：当生成结果失败时，错误如何被定位到一个可以修复的结构单元上。直接在最终 SQL 上修复有两个问题。第一，SQL 的表达顺序和逻辑构造顺序不完全一致，例如 SELECT 写在前面，但很多字段依赖 FROM、JOIN、GROUP 和 AGG 的结果。第二，复杂 SQL 中别名、聚合、排序和过滤经常跨子句引用，模型看到一整条 SQL 时很难判断到底是哪一步错了。
因此，我把研究空间总结为一句话：已有工作证明了过程化生成、执行反馈和多智能体有价值，但面向 SQL+ 中间步骤的错误定位、技能路由和局部修复仍有研究空间。后续的实验也会围绕这个空间展开，不只比较准确率，还比较 token、latency、别名依赖、跨子句引用、error localization、patch minimality 和 repair rounds。
【答辩备注】
这里回应导师反馈：研究现状不能只罗列论文，要说明论文解决了什么、还留下什么问题。答辩时可强调“我的研究不是否定已有工作，而是在已有分解、反馈和 agent 方法基础上，补充 SQL+ 步骤层的 repairability 研究”。""",

6: """【汇报讲稿】
这一页进入动机测试。为了避免研究内容像工程步骤，我把前期实验设计成几个可检验的问题。第一个问题是 SQL+ 是否真的提供了更适合修复的表达结构；第二个问题是 SQL+ 的表达成本和生成成本是多少；第三个问题是多智能体和 repair skill 是否比单一 Refiner 更稳定；第四个问题是这套机制能否迁移到公开数据集的小子集。
这里的实验不是为了证明 SQL+ 在所有方面都优于标准 SQL，而是为了弄清楚它的价值边界。比如，如果 SQL+ 只是 token 更短，那意义并不充分；如果 SQL+ 虽然更长，但能降低别名依赖、减少跨子句引用、提高错误定位和局部 patch 的稳定性，那么它作为中间表示才有研究价值。
因此，动机测试分成两类：一类是表示层测试，比较 Standard SQL、SQL+、SemQL-style proxy、NatSQL-style proxy 和 Pipe-style proxy 的结构复杂度；另一类是修复层测试，比较 Direct SQL Refiner、SQL+ Refiner、Critic-Refiner、Skill Router + Repair Skills 等方法的修复结果。
【答辩备注】
可能问题：为什么不直接做完整 benchmark？回答：开题阶段首先要证明问题定义和方法路线可行，所以先做受控实验；完整 benchmark 是后续扩展目标。
可能问题：proxy 是什么？回答：proxy 是受控近似表示，用同一批样例构造类似 SemQL、NatSQL、Pipe-style 的表达形态，用于比较结构指标，不代表完整复现原系统。""",

7: """【汇报讲稿】
这一页说明统一实验设置。当前实验主要使用三类数据。第一类是自建订单分析数据集，共 30 条查询，每条包含自然语言问题、gold SQL 和 gold SQL+，用于 SQL+ 转换、IR 复杂度、初次生成和生成成本对比。第二类是 SQL+ 已知失败集，来自 NL2SQL+ prompt v2 的失败样例，当前用于修复实验的核心集合是 13 条，覆盖 value-linking、ORDER、aggregation、join、projection 等错误类型。第三类是 Spider dev 中 concert_singer 数据库的 20 条受支持小子集，用于公开 benchmark 场景下的迁移可行性验证。
模型设置方面，初次生成和 IR 生成成本实验使用 gpt-5-mini，通过 OpenAI API 记录 input tokens、output tokens、total tokens 和本地 latency。执行环境使用 SQLite，SQL+ 先经过本地 parser 和 translator 转为 SQL，再由 Executor 执行。评价中 gold SQL 只用于最后的 execution match，不作为真实生成或 repair 的输入。
核心指标包括：SQL+ valid rate、SQL executable rate、execution match、token cost、latency、alias dependency、cross-clause reference、repair success、error localization accuracy、patch minimality 和 average repair rounds。这样设计的目的，是把“是否正确”扩展为“为什么错、错在哪里、修得是否局部、代价是否可接受”。
【答辩备注】
一定要说明实验边界：自建数据用于机制验证，Spider 小子集用于迁移 smoke test，不是完整 Spider 排名；known-failure set 用于验证修复链路，不是大规模真实错误分布。""",

8: """【汇报讲稿】
这一页回答“为什么使用 SQL+”的第一个问题：SQL+ 是否在结构上更适合错误定位和局部修复。实验数据来自自建订单分析数据集的 30 条查询。对同一批 gold 查询，我构造五种表示：Standard SQL、SQL+、SemQL-style proxy、NatSQL-style proxy 和 Pipe-style proxy。这里的 SemQL-style、NatSQL-style 和 Pipe-style 都是 controlled proxy，意思是它们只是为了在同一批样例上比较表达形态，不代表完整复现原论文系统。
实验变量是中间表示形式；固定条件是同一批查询、同一 schema、同一套统计脚本。统计指标包括 token、步骤或子句数量、嵌套深度、别名依赖和跨子句引用。当前图中重点展示别名依赖和跨子句引用，因为这两个指标直接影响修复难度。别名依赖指 AS 别名在后续 SELECT、ORDER、HAVING 等位置被引用的次数；跨子句引用指同一个逻辑字段或表达式跨多个 SQL 区域重复出现的次数。
结果显示，SQL+ 的平均别名依赖为 0.7，跨子句引用为 1.0，低于 Standard SQL 的 2.03 和 2.33。这个结果说明 SQL+ 并不是更短的表示，实际上它的平均 token 比标准 SQL 略高；但它把查询逻辑拆成更清楚的步骤边界，减少跨区域耦合。
所以这一页的结论是：SQL+ 的价值不是压缩长度，而是把错误更容易落到 WHERE、JOIN、GROUP/AGG、ORDER、SELECT 等局部步骤上，为 Critic Agent 和 Repair Skill 提供可操作的修复入口。
【答辩备注】
可能问题：这些指标是不是人为设计的？回答：是开题阶段用于结构复杂度比较的 proxy 指标，不是最终唯一评价。它们服务于 repairability 假设，后续还要结合真实修复成功率、patch minimality 和 repair rounds 验证。
可能问题：为什么 SemQL proxy 嵌套更高？回答：SemQL-style 更接近树状语义结构，在当前 proxy 实现中会显式表示层级关系，因此嵌套深度更高；这说明树状中间表示和步骤化局部修复目标并不完全一致。""",

9: """【汇报讲稿】
这一页继续回答 SQL+ 的成本问题。因为即使 SQL+ 更适合修复，如果它在生成阶段带来很高成本，也需要说明这种成本是否有可能被修复收益抵消。实验仍然使用同一批 30 条自建订单分析查询，模型使用 gpt-5-mini。对比方法包括 Direct SQL、SQL+、NatSQL-style proxy 和 SemQL-style proxy 四种生成目标。
实验过程是：第一步，把相同的自然语言问题和 schema 输入模型，让模型生成对应目标表示；第二步，检查输出格式是否有效，例如 SQL+ 是否能被 parser 解析；第三步，如果不是最终 SQL，就通过对应转换逻辑转成 SQL；第四步，在同一 SQLite 数据库上执行；第五步，将执行结果与 gold SQL 的执行结果比较；同时记录 OpenAI API usage 中的 input_tokens、output_tokens、total_tokens，以及脚本记录的 latency_seconds。
图中柱状图表示 execution match，折线表示平均 token 成本。结果显示，SQL+ 的 execution match 是 14/30，Direct SQL 是 12/30，NatSQL-style proxy 是 13/30，SemQL-style proxy 是 12/30。SQL+ 略高，但并不显著。同时 SQL+ 平均 token 为 813.03，高于 Direct SQL 的 599.17，平均延迟也更高。
因此这里的结论要谨慎：SQL+ 不能被论证为初次生成阶段更省成本或显著更准。它的研究价值应该转向后续修复阶段，也就是看它是否能提高错误定位、减少 patch 范围、提升修复成功率，并且这些收益能否抵消额外 token 和 latency 成本。
【答辩备注】
可能问题：SQL+ token 更高是不是缺点？回答：是成本，需要承认。开题后的实验要继续评估 repair token cost 和 end-to-end latency，不能只看 accuracy。
可能问题：为什么当前 SQL+ valid 不是 30/30？回答：部分模型输出包含当时 SQL+ 子集暂不支持的结构，例如不规范步骤或方言形式，这也说明 SQL+ 语法子集和 prompt 约束需要继续完善。""",

10: """【汇报讲稿】
这一页把前面的动机测试收束到系统设计。前面的实验说明，SQL+ 不是为了让初次生成一定更短或更准，而是为了给错误定位和局部修复提供结构化入口。因此系统不能停留在“自然语言直接生成 SQL”，而要形成“自然语言 -> SQL+ -> SQL -> 执行反馈 -> SQL+ 层修复”的闭环。
在这个闭环里，Intent Agent 负责解析用户问题中的查询目标、条件、聚合和排序意图；Schema Agent 负责找相关表、字段和可能的连接关系；Planner Agent 把查询拆成步骤；SQL+ Generator 生成中间表示；Translator 把 SQL+ 转成可执行 SQL；Executor 执行 SQL 并返回结果或错误；Critic Agent 负责判断错误类型和可疑步骤；Skill Router 根据错误类型选择修复技能；Repair Skill 只对 SQL+ 的局部步骤生成 patch；最后再回到 Translator 和 Executor 重新验证。
这里的重点不是 agent 数量，而是每个环节都要产出可检查的中间结果。比如 Critic 不能只说“SQL 错了”，而要输出可能是 WHERE 值链接错误、JOIN 路径错误、AGG 口径错误或 ORDER 排序错误。Router 也不是再调用一个大模型重写，而是选择对应的局部 repair skill。
【答辩备注】
可能问题：这是不是工程流程？回答：不是只描述工程步骤。每个模块对应一个研究问题和评价指标，例如 Schema Agent 对应 schema linking accuracy，Critic 对应 localization accuracy，Router 对应 routing accuracy，Repair Skill 对应 patch minimality 和 repair success。""",

11: """【汇报讲稿】
这一页说明研究不足与本文切入点。已有工作已经在多个方向上取得进展。中间表示方面，SemQL、NatSQL 等说明将 SQL 转化为更容易生成的形式是有价值的；SQL 扩展方面，Pipe Syntax 说明线性化表达可以提高可读性和可维护性；多智能体方面，MAC-SQL、CHASE-SQL 等说明多角色协作和候选生成选择能够改善复杂查询；反馈修正方面，SQLCritic 等工作说明分子句 critic 和执行反馈有助于修正错误。
但是这些工作多数仍然围绕最终 SQL 的生成、选择或整体修正。我的切入点是把错误定位和修复的对象前移到 SQL+ 中间表示层。也就是说，不是等 SQL 错了以后让模型重新写一整条 SQL，而是先判断错误对应 SQL+ 中的哪个步骤，再由对应 repair skill 修改局部步骤。
这就形成本文的研究空间：通过 SQL+ 提供可诊断步骤，通过 Critic 定位错误，通过 Router 路由技能，通过 Repair Skill 做局部 patch，通过 Executor 验证修复结果。这个空间既和已有 Text-to-SQL、agentic NL2SQL 工作相关，又有自己的评价重点：可定位性、可修复性和修复代价。
【答辩备注】
如果老师问“创新性在哪里”，回答时不要夸大。可以说：本课题尝试把 SQL+ 中间表示和多智能体反馈修正结合起来，研究错误如何从最终 SQL 映射回中间步骤，并用 repair skill 做局部修复。""",

12: """【汇报讲稿】
这一页展示 SQL+ 与反馈修正的关系。SQL+ 的核心不是把 SQL 换一种写法，而是把查询逻辑拆成一系列可以诊断的步骤。例如 FROM 表示起始数据源，JOIN 表示表连接关系，WHERE 表示过滤条件，GROUP 和 AGG 表示分组和聚合口径，SELECT 表示最终投影列，ORDER 和 LIMIT 表示 top-k 排序逻辑。
这样设计之后，每个步骤都可以对应一类常见错误。WHERE 可以定位值链接和过滤条件错误，JOIN 可以定位连接路径、连接方向和冗余连接错误，AGG 可以定位聚合函数、分组维度和别名引用错误，SELECT 可以定位输出列多列、少列或顺序错误，ORDER 可以定位排序字段、排序方向和 top-k 逻辑错误。
因此，SQL+ 不只是 translator 的输入，也是 Critic 和 Repair Skill 的操作对象。Critic Agent 可以输出“可疑步骤是 AGG，错误类型是 aggregation”，Router 就能把这个错误交给 aggregation repair skill。Repair Skill 不需要重写整条查询，只需要修改 AGG 或 ORDER 这些局部步骤，再交给 Translator 和 Executor 验证。
【答辩备注】
可能问题：SQL+ 会不会增加系统复杂度？回答：会增加一层中间表示和转换成本，但换来的是步骤边界、局部修复入口和可解释诊断信息。后续实验会继续评估这个成本是否被修复收益抵消。""",

13: """【汇报讲稿】
这一页通过一个具体例子说明 SQL 和 SQL+ 的差别。自然语言问题是查询已支付订单中金额最高的 10 个订单，返回订单编号和订单金额。标准 SQL 中 SELECT、FROM、JOIN、WHERE、GROUP BY、ORDER BY、LIMIT 等子句同时出现，而且 SELECT 中的 amount 别名还会在 ORDER BY 中被引用。对于模型来说，一旦结果错了，很难判断问题来自 paid 过滤、orders 与 order_items 的 join、SUM 金额口径，还是 ORDER BY 排序字段。
SQL+ 把同一个查询拆成更接近数据流的步骤：先 FROM orders，再 JOIN order_items，然后 WHERE status='paid'，接着 GROUP 按订单分组，再 AGG 计算金额，随后 SELECT 输出订单号和金额，最后 ORDER 和 LIMIT 取 top-k。这样每一步都对应一个可诊断问题。
这页右侧表格说明了 SQL+ 的锚点：WHERE 检查状态值是否正确，JOIN 检查连接路径，GROUP/AGG 检查聚合维度和金额口径，SELECT 检查最终输出列，ORDER/LIMIT 检查 top-k 排序逻辑。后续 Repair Skill 修改的也是这些局部步骤，而不是整体重写 SQL。
【答辩备注】
可能问题：这个例子是不是太简单？回答：它是讲解用例，目的是说明 SQL+ 的步骤边界。后续实验中样例覆盖多表 join、聚合、排序、时间范围、projection 和 value-linking 等错误类型。""",

14: """【汇报讲稿】
这一页是系统架构。上半部分是生成链路：用户输入自然语言问题后，Intent Agent 识别查询目标、过滤条件、聚合需求和排序限制；Schema Agent 根据问题和数据库 schema 选择相关表、字段和连接关系；Planner Agent 形成 SQL+ 步骤草图；SQL+ Generator 输出完整 SQL+ 表达。
中间部分是转换与执行：Translator 将 SQL+ 转换为标准 SQL 或后续适配达梦 SQL 的方言 SQL；Executor 执行 SQL，得到结果、错误信息或执行异常。如果 SQL 可执行但结果明显异常，也可以进入语义反馈流程。
下半部分是修复闭环：Critic Agent 根据执行反馈、SQL+ 步骤和 schema 信息定位错误类型和可疑步骤；Skill Router 根据 Critic 的输出选择 value、order、aggregation、join、projection 或 semantic repair skill；Repair Skill 只生成 SQL+ 层局部 patch；patch 之后重新进入 Translator 和 Executor 验证。如果验证通过，输出最终 SQL 和结果；如果失败，可以继续下一轮修复。
这张图的关键是反馈回路的位置：不是失败后重新从自然语言生成一整条 SQL，而是失败反馈进入 Critic、Router 和 Repair Skill，把修改限制在 SQL+ 局部步骤中。
【答辩备注】
可能问题：多智能体是不是会增加延迟？回答：会增加调用成本，所以后续实验会记录 token 和 latency。本课题关注的是复杂查询失败后的修复收益是否能抵消额外成本。""",

15: """【汇报讲稿】
这一页说明可行性测试的整体路径。开题阶段我没有直接追求完整 Spider 或 BIRD 排名，而是先验证链路是否跑通、指标是否可度量、研究假设是否有初步依据。
第一步是 SQL+ 转换可行性测试：使用自建 30 条查询，验证 gold SQL+ 是否能稳定转换为 SQL，并与 gold SQL 执行结果一致。第二步是表达复杂度和生成成本测试：比较 SQL、SQL+、SemQL-style proxy、NatSQL-style proxy 和 Pipe-style proxy 的结构指标、token、latency 和 execution match。第三步是修复链路测试：收集 SQL+ 和 Direct SQL 的失败样例，比较单 Refiner、Critic-Refiner 和 Skill Router + Repair Skills。第四步是 Spider 小规模迁移测试：在 concert_singer 20 条受支持子集上验证 conversion smoke 和 fresh e2e。
这些测试共同回答三个问题：SQL+ 是否可转换，SQL+ 的结构价值在哪里，多智能体局部修复是否有初步收益。
【答辩备注】
如果老师问为什么实验规模不大，可以回答：开题阶段重点是证明研究路线和评价框架可行；后续计划扩展多数据库、多难度、多 SQL 结构，并补齐端到端 token 和 latency。""",

16: """【汇报讲稿】
这一页列出实验数据、模型和评价指标。自建订单分析数据集包含 30 条自然语言查询，每条有 gold SQL 和 gold SQL+，覆盖单表过滤、多表 join、聚合、分组、排序、top-k 和 projection 等常见分析查询。SQL+ 已知失败集主要来自 NL2SQL+ prompt v2 的失败输出，目前用于修复实验的是 13 条典型失败样例。Spider 小子集来自 dev 集中的 concert_singer 数据库，当前选取 20 条 SQL+ 子集能够覆盖的查询。
模型方面，初次生成和 IR 生成成本实验使用 gpt-5-mini，记录 API 返回的 token usage 和脚本侧 latency。执行环境使用 SQLite，SQL+ 先通过本地 parser 和 converter 转成 SQL，再执行。评价时，gold SQL 只用于比较执行结果，不作为生成或修复过程的输入。
主要指标分三类。生成指标包括 SQL+ valid、SQL executable、execution match、token cost 和 latency。结构指标包括 token length、step/clause count、nesting depth、alias dependency 和 cross-clause reference。修复指标包括 repair success、error localization accuracy、router accuracy、patch minimality、average repair rounds 和 repair token cost。
【答辩备注】
这里要主动说明两个边界：第一，Spider conversion smoke test 是 gold SQL -> SQL+ -> SQL，不是端到端生成；第二，13/13 repair 是已知失败集上的结果，不是完整 benchmark。""",

17: """【汇报讲稿】
这一页展示初次生成 baseline。实验数据是自建订单分析数据集的 30 条自然语言问题。Direct SQL 组让模型直接生成 SQL；NL2SQL+ prompt v1 和 v2 让模型先生成 SQL+，再由本地 SQL+ converter 转为 SQL。固定条件是同一批问题、同一 schema、同一模型 gpt-5-mini、同一执行数据库和同一评价脚本。
评价流程是：先检查生成表示是否符合格式；再检查转换后 SQL 是否能在 SQLite 上执行；最后将生成 SQL 的执行结果与 gold SQL 的执行结果进行比较，得到 execution match。图中展示的是最关键的执行一致数。
结果是 Direct SQL 为 16/30，NL2SQL+ prompt v1 为 13/30，NL2SQL+ prompt v2 为 17/30。v2 相比 Direct SQL 只多 1 条，提升幅度很小。这个结果很重要，因为它说明只把输出格式从 SQL 换成 SQL+，并不能直接带来显著准确率提升。
因此，SQL+ 的研究重点不能停留在初次生成，而要和诊断、路由、修复结合起来。SQL+ 的价值更可能体现在失败之后：它能不能让 Critic 更容易定位错误，能不能让 Repair Skill 只改局部步骤，能不能提高修复成功率和 patch minimality。
【答辩备注】
可能问题：为什么初次生成没有明显提升？回答：因为当前 SQL+ 语法和 prompt 只是把结构线性化，模型仍然需要正确理解 schema、值、join 和聚合。SQL+ 的主要假设不是“第一次一定更准”，而是“失败后更容易修”。""",

18: """【汇报讲稿】
这一页是反馈修正策略对比，也是目前最能体现多智能体和 repair skill 价值的一组实验。数据来自两个失败集合：一个是 SQL+ prompt v2 的 13 条已知失败样例，一个是 Direct SQL 的 14 条失败样例。失败类型包括 value-linking、ORDER、aggregation、join 和 projection。实验的固定条件是同一执行数据库、同一 gold SQL 评价方式、同一 SQL+ parser/translator 和 Executor。
对比方法包括五类。第一是 SQL+ 非 gold 单 Refiner，也就是只给模型执行反馈和当前 SQL+，让它修复。第二是 Direct SQL 非 gold Refiner，在最终 SQL 层修。第三是 Schema-Critic-Refiner，把 schema 诊断加入修复。第四是 Step-wise Critic-Refiner，让 Critic 按 SQL+ 步骤给出诊断。第五是 SQL+ Skill Router + Repair Skills v3，也就是 Critic 输出错误类型和可疑步骤后，由 Router 选择 value、order、aggregation、join、projection 等局部 skill，再由 Executor 验证。
修复成功的判定是修复后的 SQL 或 SQL+ 转 SQL 能执行，并且执行结果与 gold SQL 一致。结果显示，SQL+ 单 Refiner 为 4/13，Direct SQL Refiner 为 6/14，Schema-Critic 和 Step-wise Critic 都是 3/13，而 SQL+ Skill Router + Repair Skills v3 在当前 13 条已知失败集上达到 13/13。
这个结果说明，多 Agent 数量本身不是重点。真正有效的是把诊断、路由、局部 patch 和执行验证分开，让每一步都有明确职责和可检查中间产物。
【答辩备注】
必须强调边界：13/13 不是完整 benchmark 成绩，只是当前已知失败集上的链路可行性证据。它说明“错误类型路由 + 局部 repair skill”值得继续研究，但还需要扩展未知错误、多数据库样例和无 gold 语义诊断。""",

19: """【汇报讲稿】
这一页展示分治 repair skill 的设计。前面的实验发现，单一 Refiner 不稳定，单纯增加 Critic 文本也不一定提高成功率。因此我把常见错误拆成几类，并为每类设计局部修复技能。
当前已经实现的技能包括 value-linking、order、aggregation、join、projection，以及后续拆出的 semantic repair skill。value-linking 处理状态值、日期边界和枚举值错误；order 处理排序字段、排序方向和 top-k；aggregation 处理 GROUP BY、COUNT/SUM 口径、HAVING 和 AGG 别名；join 处理连接路径、连接方向、冗余 join 和缺失 join；projection 处理输出列多列、少列、顺序错误和明细标识保留问题。
每个 skill 的运行方式不是让模型重写整条 SQL，而是根据 Critic 输出的错误类型和 SQL+ 步骤，生成或选择候选 patch，然后由 SQL+ parser、SQL converter 和 Executor 检查有效性与执行结果。这种分治方式的意义在于缩小修改范围，让修复更可解释，也方便记录每次修复到底修改了哪个步骤。
【答辩备注】
可能问题：skill 是不是硬编码？回答：当前开题阶段有规则化和模型辅助两类实现，目的是先验证按错误类型局部修复是否有效。后续可以把部分规则升级为工具增强或模型生成候选 patch。""",

20: """【汇报讲稿】
这一页展示公开数据集小规模验证。数据来自 Spider dev 中的 concert_singer 数据库，当前选取 20 条 SQL+ 子集能够覆盖的查询。这里要特别区分 conversion smoke test 和端到端生成测试。
conversion smoke test 的过程是：从 Spider gold SQL 出发，人工或规则映射成 SQL+，再由 SQL+ converter 转回 SQL，最后在 Spider 提供的 SQLite 数据库上比较执行结果。因此 20/20 只说明 SQL+ 表达与转换机制在这个小子集上可迁移，并不是模型从自然语言端到端生成的准确率。
端到端测试则是从自然语言问题和 schema 出发，让系统 fresh 生成 SQL+，再转换成 SQL 执行。当前 fresh e2e 结果是 19/20；对同一批 fresh 输出，再通过 Skill Router -> semantic repair skill 进行局部修复后达到 20/20。这说明 SQL+ 链路在公开 benchmark 小子集上有初步可行性，同时 semantic repair skill 能处理至少一部分端到端生成偏差。
但这个结果仍然有明显边界：只有一个 Spider 数据库、20 条样例、受 SQL+ 子集覆盖限制，不能表述为完整 Spider benchmark 成绩。后续需要补齐多数据库、多 SQL 结构和更大规模测试。
【答辩备注】
可能问题：为什么 conversion smoke test 不算端到端？回答：因为它从 gold SQL 出发，不经过模型自然语言生成，所以只能验证转换覆盖性，不能验证生成能力。
可能问题：20/20 是否说明系统已经很好？回答：不能。它只是小子集可行性证据，开题报告中必须写清楚边界。""",

21: """【汇报讲稿】
这一页总结目前可行性。第一，SQL+ 表达和转换链路已经在自建 30 条样例上跑通，SQL+ 转 SQL 后执行结果与 gold SQL 一致，说明中间表示和转换器具备基本可行性。第二，IR 复杂度实验说明 SQL+ 并不一定更短，但别名依赖和跨子句引用更少，支持“步骤边界更清晰”的假设。第三，初次生成实验说明 SQL+ prompt 单独使用提升有限，因此后续不能只优化 prompt。第四，修复实验说明 Skill Router + Repair Skills 在当前已知失败集上有明显收益。第五，Spider 小子集实验说明 SQL+ 转换和 fresh e2e 在公开数据集小范围内具备迁移可能性。
同时也要承认不足：样例规模还小，proxy 不是完整系统复现，SQL+ 生成成本和 Critic token 成本偏高，Spider 目前只有单数据库小子集，多数据库泛化和达梦 SQL 方言适配还没有完成。
因此，目前结论应表述为：已有结果支持继续研究 SQL+ 的 repairability 价值，而不是证明 SQL+ 已经全面优于 Direct SQL 或已有系统。
【答辩备注】
这页要主动降低结论强度。可以说“前期实验提供了可行性证据和方向判断”，不要说“已经证明方法有效”。""",

22: """【汇报讲稿】
这一页说明后续实验计划和指标体系。后续实验首先要扩展公开数据集，从当前 Spider concert_singer 单数据库扩展到多数据库小子集，再逐步增加 query 难度和 SQL 结构。其次，要完善端到端多智能体生成，不只测试 SQL+ converter，而是从自然语言、schema、SQL+ 生成、SQL 转换、执行反馈到修复的完整链路。
指标上不能只看 execution accuracy。生成阶段继续记录 SQL+ valid、SQL executable、execution match、token 和 latency；诊断阶段记录 error localization accuracy，也就是 Critic 定位的步骤是否与真实错误步骤一致；路由阶段记录 router accuracy，也就是 Router 是否选择了正确 repair skill；修复阶段记录 repair success、average repair rounds、patch minimality、repair token cost 和 repair latency。
此外，还要做消融实验：去掉 SQL+、去掉 Critic、去掉 Skill Router、只做 SQL 层整体修复、只做 SQL+ 单 Refiner、使用不同中间表示 proxy。这样才能说明提升来自 SQL+ 中间表示和多智能体局部修复机制，而不是单纯来自模型能力或 prompt 调整。
【答辩备注】
这页回应导师关于“研究内容像工程步骤”的反馈。讲法应强调技术难点和评价：如何设计可转换 SQL+、如何定位错误步骤、如何选择 repair skill、如何评估修复收益是否抵消成本。""",

23: """【汇报讲稿】
这一页说明当前进度和风险控制。目前已经完成 SQL+ parser 和 converter 的初版，自建订单数据集和 30 条样例，Direct SQL 与 NL2SQL+ baseline，IR 复杂度和生成成本实验，五类 repair skill 初版，Skill Router v3，以及 Spider concert_singer 20 条小子集的转换和端到端验证。开题文档、PPT 和实验记录也已经同步整理。
后续风险主要有四个。第一是样例规模风险，当前样例还小，需要扩展 Spider/BIRD 小子集和更多业务查询。第二是泛化风险，known-failure set 上的 13/13 不能代表未知错误，需要增加新错误类型和多数据库测试。第三是成本风险，SQL+ 和 Critic 会增加 token 和 latency，需要通过 repair收益进行平衡。第四是达梦方言适配风险，目前主要在 SQLite 上验证，后续要补充达梦 SQL 函数、分页、日期和类型转换差异。
进度安排上，开题后先扩展数据集和多数据库实验，再完善无 gold 语义诊断和端到端成本记录，最后做达梦 SQL 方言适配和论文写作。
【答辩备注】
如果老师问“最大风险是什么”，可以回答：不是 parser 实现，而是 SQL+ repairability 是否能在未知数据和多数据库场景稳定体现。因此后续重点是扩展公开子集和记录更完整指标。""",

24: """【汇报讲稿】
这一页总结预期创新点。第一个创新点是面向生成与修复的 SQL+ 中间表示。它不是为了替代 SQL，也不是为了追求更短，而是把复杂查询拆成可转换、可诊断、可局部修复的步骤。
第二个创新点是面向 SQL+ 的多智能体反馈修正闭环。Intent、Schema、Planner、Generator 负责生成，Translator 和 Executor 负责转换执行，Critic 负责错误定位，Skill Router 负责选择修复技能，Repair Skill 负责局部 patch，形成可检查的中间产物链。
第三个创新点是 repairability 导向的评价体系。除了 execution accuracy，还评估 SQL+ valid、SQL executable、error localization、router accuracy、patch minimality、repair rounds、token cost 和 latency。这样可以更完整地回答 SQL+ 的收益是否真实存在，以及是否值得付出额外成本。
这些创新点目前还处于开题阶段的设计和初步验证，不夸大为已经全面领先已有系统。后续研究会通过更大样例、多数据库和更多消融实验继续验证。
【答辩备注】
表述要稳健，建议用“拟提出”“尝试构建”“初步验证”这类措辞。不要说“首次”或“显著领先”，除非后续有充分实验支撑。""",

25: """【汇报讲稿】
最后这一页列出主要参考文献。文献可以分为几类。第一类是 benchmark，包括 Spider、BIRD 和 Spider 2.0，用来说明 Text-to-SQL 正从传统跨数据库任务走向真实企业级数据分析 workflow。第二类是中间表示和 SQL 扩展，包括 IRNet/SemQL、NatSQL、GoogleSQL Pipe Syntax，用来支撑“直接生成最终 SQL 不是唯一选择，中间表示和线性化表达有研究价值”。第三类是 LLM-based Text-to-SQL，包括 DAIL-SQL、DIN-SQL 等，用来支撑 prompt、任务分解和自修正方法。第四类是多智能体和反馈修正，包括 MAC-SQL、CHASE-SQL、SQL-Factory、SQLCritic、LEVER 等，用来支撑多角色协作、候选验证和执行反馈的重要性。
这些文献共同说明，本课题不是凭空提出 SQL+ 和多智能体，而是在已有 benchmark、中间表示、SQL 扩展、LLM 生成和反馈修正研究基础上，进一步聚焦 SQL+ 步骤层的错误定位、技能路由和局部修复。
【答辩备注】
如果老师问某篇文献和本课题关系，可以这样回答：SemQL/NatSQL 支撑中间表示方向；Pipe Syntax 支撑 SQL 线性化表达；DIN-SQL 支撑任务分解和自修正；MAC-SQL/CHASE-SQL 支撑多智能体协作；SQLCritic 支撑结构化反馈修正；Spider/BIRD/Spider 2.0 支撑真实复杂场景的评价需求。""",
}


def set_notes(slide, note: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = "\n".join(line.strip() for line in note.splitlines() if line.strip())


def main() -> None:
    if not SRC.exists():
        raise FileNotFoundError(SRC)
    shutil.copyfile(SRC, OUT)
    prs = Presentation(str(OUT))
    if len(prs.slides) != 25:
        raise RuntimeError(f"expected 25 slides, got {len(prs.slides)}")
    for slide_no, note in NOTES.items():
        set_notes(prs.slides[slide_no - 1], note)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
