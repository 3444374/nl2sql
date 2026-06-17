from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
PPTS = [
    ROOT / "docs" / "opening" / "opening_ppt_template_version_v2.pptx",
    ROOT / "docs" / "opening" / "opening_ppt_template_version_v3.pptx",
]
NOTE_DIR = ROOT / "docs" / "opening" / "notes"

NOTE_OVERRIDES = {
    11: NOTE_DIR / "slide11_experiment_setup_note.md",
    12: NOTE_DIR / "slide12_sqlplus_conversion_note.md",
    13: NOTE_DIR / "slide13_ir_complexity_note.md",
    14: NOTE_DIR / "slide14_ir_metrics_note.md",
    15: NOTE_DIR / "slide15_ir_generation_note.md",
    16: NOTE_DIR / "slide16_generation_result_note.md",
    17: NOTE_DIR / "slide17_baseline_note.md",
    18: NOTE_DIR / "slide18_repair_experiment_note.md",
    19: NOTE_DIR / "slide19_repair_result_note.md",
    20: NOTE_DIR / "slide20_repairability_note.md",
    21: NOTE_DIR / "slide21_repair_skill_note.md",
    22: NOTE_DIR / "slide22_spider_smoke_note.md",
}


def apply_note(slide, note: str) -> None:
    text_frame = slide.notes_slide.notes_text_frame
    text_frame.clear()
    paragraphs = [line.strip() for line in note.splitlines() if line.strip()]
    text_frame.text = "\n".join(paragraphs)


def main() -> None:
    for ppt in PPTS:
        if not ppt.exists():
            continue
        prs = Presentation(str(ppt))
        for slide_no, note_path in NOTE_OVERRIDES.items():
            if slide_no > len(prs.slides):
                raise RuntimeError(f"{ppt} has no slide {slide_no}")
            note = note_path.read_text(encoding="utf-8")
            apply_note(prs.slides[slide_no - 1], note)
        prs.save(ppt)
        print(f"updated experiment notes in {ppt}")


if __name__ == "__main__":
    main()
