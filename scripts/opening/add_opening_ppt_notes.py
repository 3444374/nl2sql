from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
PPT = ROOT / "docs" / "opening" / "opening_ppt_template_version_v2.pptx"


NOTES = [
    """这一页先说明课题定位。我的题目聚焦 SQL+、Text-to-SQL、多智能体和执行反馈修正。这里需要强调，本课题不是做一个单轮自然语言转 SQL 的工程 demo，而是研究如何通过一个更适合生成和修复的中间表示，把查询生成、执行验证和错误修正连接成闭环。汇报时可以先交代研究背景来自达梦 SQL+ 和自然语言数据库查询两个方向。""",
    """这一页介绍汇报结构。建议控制在半分钟内，告诉老师后面会按“问题来源、已有研究、方法设计、初步实验、后续计划”的顺序展开。这里的重点是让实验部分显得不是突然出现，而是服务于前面提出的研究问题。""",
    """这一页讲传统 Text-to-SQL 的基本流程。自然语言问题经过模型或解析器生成标准 SQL，再交给数据库执行。这个流程看起来直接，但问题是反馈通常只落在最终 SQL 上。SQL 一旦出错，很难判断是表选错、字段选错、join 错、聚合错，还是排序和过滤条件错。这里可以指出，复杂查询不是一个线性生成问题，而是多个相互耦合的决策问题。""",
    """这一页补充真实场景中的困难。企业数据库中常见多表关系、字段别名、时间口径和方言差异。大模型生成的 SQL 即使能执行，也可能在语义上不符合用户意图。比如统计口径错、过滤条件漏掉、join 路径不合理，这些错误不会一定报语法错。由此引出本课题关注的不只是生成 SQL，还包括验证、诊断和修复。""",
    """这一页概括国内外研究现状。传统方法强调语义解析、schema linking 和语法约束，LLM 方法进一步引入 few-shot、任务分解和执行反馈。Spider、BIRD、Spider 2.0 等 benchmark 也说明任务越来越接近真实数据分析场景。这里要讲清楚研究演进：从直接生成，到过程化生成，再到执行反馈和 agent 化。""",
    """这一页回答为什么不直接用已有中间表示。SemQL 更像语义解析树，NatSQL 更接近自然语言 SQL，Pipe Syntax 则是 GoogleSQL 的方言扩展。它们都提供了参考，但本课题需要的是面向 NL2SQL 生成和反馈修正的中间层。SQL+ 的设计目标不是最短表达，而是保留 SQL 可转换性，同时让错误可以映射到步骤上。""",
    """这一页说明 SQL+ 的基本表达方式。可以按 FROM、JOIN、WHERE、GROUP、AGG、HAVING、SELECT、ORDER、LIMIT 的顺序解释。每一个步骤都对应一类可能的错误和修复技能。比如 WHERE 对应值链接错误，JOIN 对应连接路径错误，AGG 对应聚合口径和别名引用错误。这里要突出 SQL+ 是 repair 的锚点。""",
    """这一页讲本文拟研究的多智能体闭环。前半部分由 Intent、Schema、Planner 和 SQL+ Generator 完成理解和生成，后半部分由 Translator、Executor、Critic、Skill Router 和 Repair Skill 完成执行、定位和修复。这里要避免把它讲成普通流程图。重点是每个 agent 都有可检查的中间输出，错误不是简单让模型重新生成，而是路由到对应的修复技能。""",
    """这一页集中提出研究问题和评价维度。可以按四个问题讲：SQL+ 是否必要，SQL+ 如何设计，多智能体如何协作，修复收益能否抵消额外成本。评价指标也要比单纯准确率更丰富，包括 valid rate、execution accuracy、token、latency、error localization accuracy、patch minimality 和 repair rounds。这样能回应导师说的“研究内容不能只是工程步骤”。""",
    """这一页把实验和研究假设对应起来。H1 验证 SQL+ 表达和转换是否可行，H2 比较 SQL+ 和其它表示的表达差异，H3 看初次生成是否有收益，H4 看修复收益是否成立。这里要说明，实验不是堆结果，而是逐步回答为什么用 SQL+、SQL+ 的代价是什么、修复是否有价值。""",
    """这一页是实验设计总览。汇报时不需要逐格念完，可以说当前实验分为五类：转换可行性、IR 复杂度、IR 生成成本、反馈修复和公开数据集 smoke test。每类实验都有目的、对比对象和指标。这里也要提醒，SemQL-style、NatSQL-style、Pipe-style 是 controlled proxy，不是完整复现原系统。""",
    """这一页讲第一个实验。目的很简单：证明 SQL+ 不是停留在设想层，而是能解析、转换并执行。当前自建订单数据集 30 条样例，SQL+ 全部可解析，转换 SQL 全部可执行，并且和 gold SQL 执行结果一致。这个结果支撑后续实验，但不能直接说明生成质量已经好。""",
    """这一页讲 IR 表达复杂度实验的目的。这个实验回答“为什么用 SQL+”。结果显示 SQL+ 的 token 不一定更少，平均 token 比标准 SQL 略高。但是 SQL+ 的别名依赖和跨子句引用明显更低。这意味着它不一定降低输入长度，但降低了表达耦合，更适合定位错误和局部修复。""",
    """这一页展示 IR 复杂度表格。这里要解释几个指标：步骤或子句数表示表达被拆成多少操作单元，嵌套深度表示结构复杂程度，别名依赖表示后续步骤依赖前面别名的程度，跨子句引用表示一个语义项跨多个子句相互引用的情况。SQL+ 的优势体现在依赖更少，而不是更短。""",
    """这一页讲 IR 生成成本实验。这个实验是为了避免只讲 SQL+ 的好处。结果显示 SQL+ 初次生成准确率略高于 Direct SQL，但 token 和 latency 也更高。因此当前不能说 SQL+ 在初次生成阶段显著优于 SQL。更合理的结论是：SQL+ 的价值需要在后续修复阶段继续验证。""",
    """这一页展示生成成本表格。可以重点讲三点：Direct SQL 成本最低，但执行一致率只有 12/30；SQL+ 执行一致率为 14/30，但 token 和延迟更高；SemQL-style proxy 成本最高且可执行率下降。这个结果对课题是有价值的，因为它逼迫我们把研究重心放在修复收益，而不是单纯包装 SQL+。""",
    """这一页讲 baseline 对比。Direct NL2SQL 是 16/30，NL2SQL+ prompt v2 是 17/30，提升很小。这个结果说明只把目标从 SQL 换成 SQL+ 并不够。后续必须引入 Schema Agent、Critic Agent、Skill Router 和 repair skill，才能体现 SQL+ 作为中间表示的价值。""",
    """这一页讲反馈修复实验。单个 Refiner 的表现并不好，SQL+ 非 gold 单 Refiner 是 4/13，Direct SQL Refiner 是 6/14。后来引入 Skill Router 和分治 repair skill 后，当前已知失败集达到 13/13。这里一定要说明这个结果只是在当前已知失败集上成立，不能说成大规模 benchmark 结果。""",
    """这一页展示不同修复方法的对比。重点不是简单说 13/13，而是解释为什么前面的多 agent 串联效果不稳定。Schema-Critic-Refiner 和 Step-wise Critic-Refiner 都没有明显提升，说明“多个 agent”本身不是答案。真正有效的是把错误类型和 repair skill 对齐，让每种错误有对应的局部修复策略。""",
    """这一页讲 repairability 指标。SQL+ Router 在修复成功率上更好，patch minimality 也较高，但平均 repair token 明显更高。这个结果比较客观：SQL+ 带来了修复上的收益，也带来了诊断和路由成本。后续需要补齐端到端 latency，判断修复收益能否抵消生成和诊断阶段的额外开销。""",
    """这一页讲 repair skill 分治结果。当前已实现 value-linking、ORDER、aggregation、join、projection 五类 skill。每类 skill 对应不同错误，不再让模型整条重写。比如 aggregation skill 处理 AGG 别名、HAVING 和 ORDER 引用，join skill 处理连接路径和 paid 过滤条件。这里体现课题的可解释性和可修复性。""",
    """这一页讲 Spider smoke test。当前只在 Spider dev 中 concert_singer 数据库的 20 条受支持样例上测试，SQL+ 有效、SQL 可执行、执行一致都是 20/20。这个结果的意义是验证 SQL+ 转换链路能迁移到公开数据集子集，不是完整 Spider benchmark 分数。答辩时一定要主动说明边界。""",
    """这一页给出综合判断。当前结果支持三点：SQL+ 转换链路可行，SQL+ 的优势主要在表达依赖和修复锚点，Skill Router 加 repair skill 是当前更有效的修复路线。同时也要承认 SQL+ 初次生成成本更高，样例规模还小，不能过度夸大。""",
    """这一页讲后续实验计划。建议按优先级讲：第一，扩展 Spider 多数据库子集；第二，补齐端到端 token 和 latency；第三，做无 gold 的语义错误诊断；第四，适配达梦 SQL 方言；第五，如果时间允许，再接入 SemQL 或 NatSQL 原系统做更严格对比。""",
    """这一页讲预期创新点。第一是面向生成和修复的 SQL+ 中间表示。第二是 SQL+ 多智能体反馈闭环。第三是面向 repairability 的评价体系。这里不要说“首创”或“显著领先”，更稳妥的说法是：本课题尝试把中间表示设计和反馈修正机制结合起来。""",
    """这一页讲风险和应对。风险包括样例规模偏小、proxy 对比不能代替完整复现、Skill Router 当前依赖已知失败集、SQL+ 成本更高、达梦方言还没充分验证。每个风险都对应后续实验安排。这样老师会看到你知道问题在哪里，也知道如何补。""",
    """这一页总结。可以用三句话收束：本课题研究的不是单次生成 SQL，而是生成、执行、诊断和修复闭环；当前实验说明 SQL+ 在转换和局部修复上有初步可行性；后续要扩大样例、强化无 gold 诊断、补齐成本指标和达梦方言适配。""",
    """这一页列参考文献。汇报时不需要逐条念，可以说明参考文献覆盖四类：Text-to-SQL benchmark，SQL 中间表示，LLM-based Text-to-SQL，多智能体和执行反馈修复。GoogleSQL Pipe Syntax 用来支撑 SQL+ 的线性表达设计，不是直接照搬 GoogleSQL 语法。""",
]


def main() -> None:
    prs = Presentation(str(PPT))
    if len(prs.slides) != len(NOTES):
        raise RuntimeError(f"slide count mismatch: {len(prs.slides)} slides, {len(NOTES)} notes")

    for slide, note in zip(prs.slides, NOTES):
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.clear()
        notes_tf.text = note

    prs.save(PPT)
    print(f"wrote notes to {PPT}")


if __name__ == "__main__":
    main()
