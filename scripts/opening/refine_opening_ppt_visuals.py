from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
IN_PPT = ROOT / "docs" / "opening" / "opening_ppt_template_version_v2.pptx"
OUT_PPT = ROOT / "docs" / "opening" / "opening_ppt_template_version_v3.pptx"

FONT = "Microsoft YaHei"
BLUE = RGBColor(31, 78, 121)
ORANGE = RGBColor(198, 89, 17)
GREEN = RGBColor(84, 130, 53)
RED = RGBColor(170, 60, 60)
DARK = RGBColor(45, 45, 45)
MUTED = RGBColor(100, 100, 100)
LIGHT = RGBColor(238, 242, 247)
PALE_BLUE = RGBColor(221, 235, 247)
PALE_ORANGE = RGBColor(251, 229, 214)
PALE_GREEN = RGBColor(226, 239, 218)
WHITE = RGBColor(255, 255, 255)


def set_text(shape, text: str, size: int = 16, bold: bool = False, color=DARK, align=PP_ALIGN.CENTER) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def cover_body(slide) -> None:
    # Remove the original dense body elements before rebuilding the slide.
    # This keeps the refined deck easier to edit than simply covering old text.
    for shape in list(slide.shapes):
        if Inches(1.0) <= shape.top < Inches(6.95):
            slide.shapes._spTree.remove(shape._element)
    shape = slide.shapes.add_shape(1, Inches(0.38), Inches(1.08), Inches(9.25), Inches(5.82))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = WHITE


def box(slide, x, y, w, h, text, fill=LIGHT, stroke=BLUE, size=13, bold=False, color=DARK):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    set_text(shape, text, size, bold, color)
    return shape


def label(slide, x, y, w, h, text, size=13, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, bold, color, align)
    return shape


def arrow(slide, x, y, text="→", color=BLUE):
    return label(slide, x, y, 0.35, 0.24, text, 20, True, color, PP_ALIGN.CENTER)


def bullets(slide, x, y, items: list[str], size=12, color=DARK) -> None:
    text = "\n".join(f"• {item}" for item in items)
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(8.6), Inches(1.2))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)


def slide_2(slide) -> None:
    cover_body(slide)
    steps = [
        ("01", "问题来源", "为什么直接 NL2SQL 不稳定"),
        ("02", "已有研究", "Text-to-SQL、IR、Agent"),
        ("03", "本文方法", "SQL+ 与多智能体闭环"),
        ("04", "初步实验", "目的、条件、指标、结果"),
        ("05", "计划与风险", "后续实验和边界说明"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = 0.55 + i * 1.78
        box(slide, x, 2.0, 1.42, 1.65, "", PALE_BLUE if i % 2 == 0 else PALE_ORANGE, BLUE if i % 2 == 0 else ORANGE)
        label(slide, x + 0.12, 2.12, 0.45, 0.32, num, 14, True, BLUE)
        label(slide, x + 0.12, 2.55, 1.15, 0.35, title, 15, True, DARK)
        label(slide, x + 0.12, 3.02, 1.15, 0.45, desc, 10, False, MUTED)
        if i < len(steps) - 1:
            arrow(slide, x + 1.45, 2.68, "→", ORANGE)
    box(slide, 1.0, 4.55, 8.0, 0.72, "汇报逻辑：先说明为什么需要 SQL+，再用实验解释它的收益和代价", PALE_GREEN, GREEN, 14, True)


def slide_4(slide) -> None:
    cover_body(slide)
    cards = [
        ("Schema 复杂", "多表、字段别名、外键路径不透明", PALE_BLUE, BLUE),
        ("SQL 方言差异", "函数、分页、日期处理在不同数据库中不一致", PALE_ORANGE, ORANGE),
        ("语义错不报错", "SQL 能执行，但聚合口径或过滤条件不符合意图", PALE_GREEN, GREEN),
        ("修复不稳定", "整条重写可能破坏原本正确的 join 或 filter", RGBColor(244, 230, 230), RED),
    ]
    for i, (head, body, fill, stroke) in enumerate(cards):
        x = 0.72 + (i % 2) * 4.55
        y = 1.45 + (i // 2) * 1.75
        box(slide, x, y, 3.95, 1.22, "", fill, stroke)
        label(slide, x + 0.22, y + 0.18, 3.4, 0.28, head, 15, True, stroke)
        label(slide, x + 0.22, y + 0.58, 3.42, 0.48, body, 12, False, DARK)
    box(slide, 1.0, 5.35, 8.0, 0.72, "研究焦点：从“生成一条 SQL”转向“生成、执行、定位和修复”", LIGHT, BLUE, 15, True)


def slide_5(slide) -> None:
    cover_body(slide)
    items = [
        ("传统方法", "语义解析\nschema linking\n语法约束", BLUE),
        ("LLM 方法", "few-shot\n任务分解\n执行反馈", ORANGE),
        ("Benchmark", "Spider\nBIRD\nSpider 2.0", GREEN),
        ("当前缺口", "可定位\n可路由\n可局部修复", RED),
    ]
    for i, (head, body, color) in enumerate(items):
        x = 0.65 + i * 2.25
        box(slide, x, 1.65, 1.78, 2.05, "", LIGHT, color)
        label(slide, x + 0.17, 1.88, 1.4, 0.3, head, 15, True, color)
        label(slide, x + 0.18, 2.42, 1.34, 0.9, body, 12, False, DARK, PP_ALIGN.CENTER)
        if i < 3:
            arrow(slide, x + 1.86, 2.48, "→", MUTED)
    box(slide, 0.9, 4.75, 8.2, 0.88, "研究空间：已有工作证明生成和执行反馈有价值，但中间步骤级错误定位与 repair skill 路由仍有进一步研究空间", PALE_ORANGE, ORANGE, 13, True)


def slide_6(slide) -> None:
    cover_body(slide)
    rows = [
        ("SemQL", "语义解析树", "结构约束强", "执行链路映射较重"),
        ("NatSQL", "自然语言化 SQL", "降低部分 SQL 书写难度", "局部修复不是主要目标"),
        ("Pipe Syntax", "SQL 方言扩展", "线性数据流", "依赖 GoogleSQL 生态"),
        ("SQL+", "NL2SQL 中间表示", "可转换、可定位、可 patch", "面向反馈修正设计"),
    ]
    for i, (name, pos, pro, caveat) in enumerate(rows):
        x = 0.55 + i * 2.27
        fill = PALE_BLUE if i < 2 else (PALE_ORANGE if i == 2 else PALE_GREEN)
        stroke = BLUE if i < 2 else (ORANGE if i == 2 else GREEN)
        box(slide, x, 1.4, 1.9, 3.55, "", fill, stroke)
        label(slide, x + 0.15, 1.65, 1.55, 0.35, name, 16, True, stroke, PP_ALIGN.CENTER)
        label(slide, x + 0.15, 2.25, 1.55, 0.35, pos, 11, True, DARK, PP_ALIGN.CENTER)
        label(slide, x + 0.15, 3.0, 1.55, 0.55, pro, 10, False, DARK, PP_ALIGN.CENTER)
        label(slide, x + 0.15, 3.95, 1.55, 0.55, caveat, 10, False, MUTED, PP_ALIGN.CENTER)
    box(slide, 1.1, 5.55, 7.8, 0.55, "SQL+ 的区别：不是更短的 SQL，而是为执行反馈和局部修复设计的步骤化表示", LIGHT, BLUE, 13, True)


def slide_9(slide) -> None:
    cover_body(slide)
    rqs = [
        ("RQ1", "SQL+ 是否降低表达耦合"),
        ("RQ2", "多智能体如何拆分可观察子任务"),
        ("RQ3", "反馈能否映射到 SQL+ 步骤"),
        ("RQ4", "修复收益能否抵消额外成本"),
    ]
    for i, (rq, text) in enumerate(rqs):
        x = 0.7 + (i % 2) * 4.35
        y = 1.35 + (i // 2) * 1.25
        box(slide, x, y, 3.7, 0.86, "", PALE_BLUE if i < 2 else PALE_GREEN, BLUE if i < 2 else GREEN)
        label(slide, x + 0.2, y + 0.18, 0.55, 0.3, rq, 15, True, ORANGE)
        label(slide, x + 0.86, y + 0.18, 2.55, 0.35, text, 12, True, DARK)
    metrics = ["valid rate", "execution accuracy", "token", "latency", "error localization", "patch minimality", "repair rounds"]
    for i, metric in enumerate(metrics):
        x = 0.55 + (i % 4) * 2.25
        y = 4.35 + (i // 4) * 0.72
        box(slide, x, y, 1.85, 0.42, metric, LIGHT, ORANGE, 10, True)
    label(slide, 0.75, 3.7, 8.0, 0.28, "评价不只看准确率，还看成本、定位和修复过程是否可控", 13, True, BLUE)


def slide_8(slide) -> None:
    cover_body(slide)
    label(slide, 0.72, 1.10, 2.0, 0.28, "生成链路", 12, True, MUTED)
    label(slide, 0.72, 2.85, 2.0, 0.28, "执行与诊断链路", 12, True, MUTED)
    label(slide, 0.72, 4.72, 2.0, 0.28, "局部修复闭环", 12, True, MUTED)

    top = [
        (0.55, "用户问题\nNatural language"),
        (2.35, "Intent Agent\n查询意图"),
        (4.15, "Schema Agent\n表字段关系"),
        (5.95, "Planner Agent\nSQL+ 步骤"),
        (7.75, "SQL+ Generator\nSQL+ 表达"),
    ]
    for i, (x, text) in enumerate(top):
        box(slide, x, 1.40, 1.38, 0.68, text, PALE_BLUE, BLUE, 9, True)
        if i < len(top) - 1:
            arrow(slide, x + 1.43, 1.57, "→", BLUE)

    box(slide, 7.75, 2.92, 1.38, 0.68, "Translator\nSQL+ → SQL", PALE_GREEN, GREEN, 9, True)
    box(slide, 5.95, 2.92, 1.38, 0.68, "Executor\n执行 SQL", PALE_GREEN, GREEN, 9, True)
    box(slide, 4.15, 2.92, 1.38, 0.68, "Critic Agent\n定位错误步骤", PALE_ORANGE, ORANGE, 9, True)
    box(slide, 2.35, 2.92, 1.38, 0.68, "Skill Router\n选择修复技能", PALE_ORANGE, ORANGE, 9, True)

    label(slide, 8.18, 2.18, 0.55, 0.35, "↓", 22, True, ORANGE, PP_ALIGN.CENTER)
    arrow(slide, 7.38, 3.10, "→", GREEN)
    arrow(slide, 5.58, 3.10, "→", ORANGE)
    arrow(slide, 3.78, 3.10, "→", ORANGE)

    box(slide, 2.35, 4.70, 1.65, 0.68, "Repair Skill\n局部修改 SQL+ 步骤", PALE_ORANGE, ORANGE, 9, True)
    box(slide, 4.75, 4.70, 1.65, 0.68, "SQL+ Patch\n更新中间表示", PALE_ORANGE, ORANGE, 9, True)
    box(slide, 7.25, 4.70, 1.65, 0.68, "Re-validate\n重新转换并执行", PALE_GREEN, GREEN, 9, True)
    label(slide, 2.83, 3.67, 0.55, 0.42, "↓", 22, True, ORANGE, PP_ALIGN.CENTER)
    arrow(slide, 4.08, 4.88, "→", ORANGE)
    arrow(slide, 6.48, 4.88, "→", ORANGE)

    # The feedback edge is deliberately large: SQL+ Patch returns to Translator,
    # then the converted SQL is executed again. This avoids showing repair as a
    # one-way extra stage after execution.
    label(slide, 8.13, 4.03, 0.55, 0.45, "↺", 28, True, ORANGE, PP_ALIGN.CENTER)
    label(slide, 6.75, 3.98, 2.7, 0.32, "patch 后回到 Translator / Executor", 10, True, ORANGE, PP_ALIGN.CENTER)
    label(slide, 8.18, 4.33, 0.55, 0.32, "↑", 22, True, ORANGE, PP_ALIGN.CENTER)

    box(slide, 5.85, 5.88, 2.1, 0.45, "最终 SQL / 查询结果", PALE_GREEN, GREEN, 10, True)
    label(slide, 6.55, 3.67, 0.55, 0.42, "↓", 22, True, GREEN, PP_ALIGN.CENTER)
    label(slide, 0.78, 6.45, 8.5, 0.42, "关键点：失败反馈进入 Critic 和 Router；Repair Skill 只修改 SQL+ 局部步骤；patch 后回到 Translator 和 Executor 重新验证。", 10, False, ORANGE)

def slide_23(slide) -> None:
    cover_body(slide)
    box(slide, 0.65, 1.35, 4.1, 4.2, "", PALE_GREEN, GREEN)
    label(slide, 0.9, 1.62, 3.5, 0.32, "已有支撑", 16, True, GREEN)
    bullets(slide, 0.92, 2.15, [
        "SQL+ 转换链路可行",
        "SQL+ 的低依赖特征支持修复定位",
        "Skill Router + Repair Skills 在已知失败集上有效",
    ], 12)
    box(slide, 5.15, 1.35, 4.1, 4.2, "", PALE_ORANGE, ORANGE)
    label(slide, 5.4, 1.62, 3.5, 0.32, "仍需补齐", 16, True, ORANGE)
    bullets(slide, 5.42, 2.15, [
        "SQL+ 初次生成成本更高",
        "样例规模仍偏小",
        "非 gold 诊断和端到端 latency 还要继续验证",
    ], 12)
    box(slide, 1.2, 5.95, 7.5, 0.48, "结论口径：SQL+ 的阶段性价值主要在 repairability，而不是初次生成成本", LIGHT, BLUE, 12, True)


def slide_24(slide) -> None:
    cover_body(slide)
    steps = [
        ("1", "扩展 Spider 子集", "多库、多难度、多类型查询"),
        ("2", "补齐成本指标", "Critic / Router / Repair / Executor latency"),
        ("3", "无 gold 语义诊断", "结果为空、聚合错、projection 错"),
        ("4", "达梦方言适配", "日期、分页、函数、类型转换"),
        ("5", "更严格 IR 对比", "条件允许时接入原系统代码"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.65 + i * 1.78
        box(slide, x, 1.65, 1.42, 3.35, "", PALE_BLUE if i % 2 == 0 else PALE_GREEN, BLUE if i % 2 == 0 else GREEN)
        label(slide, x + 0.43, 1.92, 0.55, 0.55, num, 22, True, ORANGE, PP_ALIGN.CENTER)
        label(slide, x + 0.16, 2.72, 1.1, 0.42, head, 11, True, DARK, PP_ALIGN.CENTER)
        label(slide, x + 0.16, 3.45, 1.1, 0.82, body, 9, False, MUTED, PP_ALIGN.CENTER)
        if i < 4:
            arrow(slide, x + 1.45, 3.02, "→", ORANGE)
    box(slide, 1.0, 5.75, 8.0, 0.55, "后续实验围绕“规模扩大、成本补齐、真实诊断、达梦适配”展开", LIGHT, BLUE, 12, True)


def slide_25(slide) -> None:
    cover_body(slide)
    items = [
        ("创新点一", "SQL+ 中间表示", "面向生成、转换、定位和局部 patch"),
        ("创新点二", "多智能体反馈闭环", "Critic、Router、Repair Skill 与 Executor 协作"),
        ("创新点三", "Repairability 评价", "定位、路由、patch、轮数、token、latency"),
    ]
    for i, (tag, head, body) in enumerate(items):
        x = 0.85 + i * 3.0
        box(slide, x, 1.75, 2.45, 3.5, "", PALE_BLUE if i == 0 else (PALE_GREEN if i == 1 else PALE_ORANGE), BLUE if i == 0 else (GREEN if i == 1 else ORANGE))
        label(slide, x + 0.25, 2.05, 1.85, 0.35, tag, 12, True, MUTED, PP_ALIGN.CENTER)
        label(slide, x + 0.25, 2.75, 1.85, 0.45, head, 16, True, DARK, PP_ALIGN.CENTER)
        label(slide, x + 0.25, 3.75, 1.85, 0.75, body, 11, False, DARK, PP_ALIGN.CENTER)
    box(slide, 1.0, 5.75, 8.0, 0.52, "表述边界：强调“尝试结合”，不使用“首创”或“显著领先”等过强说法", LIGHT, RED, 12, True)


def slide_26(slide) -> None:
    cover_body(slide)
    rows = [
        ("样例规模偏小", "扩展 Spider 多库子集和 BIRD 子集"),
        ("proxy 对比有限", "明确边界，后续视时间接入原系统"),
        ("已知失败集偏窄", "增加未知错误和无 gold 语义错"),
        ("SQL+ 成本更高", "比较修复收益与端到端成本"),
        ("达梦方言不足", "补充函数、分页、日期和类型测试"),
    ]
    label(slide, 0.75, 1.25, 3.5, 0.3, "风险", 15, True, RED)
    label(slide, 5.15, 1.25, 3.5, 0.3, "应对", 15, True, GREEN)
    for i, (risk, response) in enumerate(rows):
        y = 1.75 + i * 0.82
        box(slide, 0.65, y, 3.8, 0.55, risk, RGBColor(244, 230, 230), RED, 11, True)
        arrow(slide, 4.65, y + 0.15, "→", MUTED)
        box(slide, 5.1, y, 3.95, 0.55, response, PALE_GREEN, GREEN, 11, True)


def slide_27(slide) -> None:
    cover_body(slide)
    cards = [
        ("研究对象", "不是单次生成 SQL，而是生成、执行、诊断和修复闭环"),
        ("当前证据", "SQL+ 转换可行，局部 repair skill 对已知失败集有效"),
        ("后续重点", "扩大样例、强化无 gold 诊断、补齐成本和方言适配"),
    ]
    for i, (head, body) in enumerate(cards):
        y = 1.45 + i * 1.45
        box(slide, 1.0, y, 8.0, 0.95, "", PALE_BLUE if i == 0 else (PALE_GREEN if i == 1 else PALE_ORANGE), BLUE if i == 0 else (GREEN if i == 1 else ORANGE))
        label(slide, 1.25, y + 0.2, 1.35, 0.35, head, 15, True, BLUE if i != 2 else ORANGE)
        label(slide, 2.85, y + 0.2, 5.7, 0.42, body, 13, True, DARK)


def slide_28(slide) -> None:
    cover_body(slide)
    groups = [
        ("Benchmark", "Spider\nBIRD\nSpider 2.0"),
        ("中间表示", "SemQL\nNatSQL\nPipe Syntax"),
        ("LLM NL2SQL", "DAIL-SQL\nDIN-SQL\nCHESS\nCHASE-SQL"),
        ("Agent / Repair", "MAC-SQL\nSQLCritic\nTool-Assisted Agent\nLEVER"),
    ]
    for i, (head, body) in enumerate(groups):
        x = 0.75 + (i % 2) * 4.3
        y = 1.55 + (i // 2) * 1.9
        box(slide, x, y, 3.75, 1.35, "", LIGHT, BLUE if i % 2 == 0 else ORANGE)
        label(slide, x + 0.22, y + 0.2, 3.0, 0.28, head, 15, True, BLUE if i % 2 == 0 else ORANGE)
        label(slide, x + 0.22, y + 0.58, 3.0, 0.55, body, 11, False, DARK)
    box(slide, 1.0, 5.75, 8.0, 0.5, "参考文献覆盖：任务评测、表示设计、LLM 生成、多智能体反馈修复", PALE_GREEN, GREEN, 12, True)


REFINERS = {
    2: slide_2,
    4: slide_4,
    5: slide_5,
    6: slide_6,
    8: slide_8,
    9: slide_9,
    23: slide_23,
    24: slide_24,
    25: slide_25,
    26: slide_26,
    27: slide_27,
    28: slide_28,
}


def main() -> None:
    prs = Presentation(str(IN_PPT))
    for number, fn in REFINERS.items():
        fn(prs.slides[number - 1])
    prs.save(OUT_PPT)
    print(f"wrote visual refined deck to {OUT_PPT}")


if __name__ == "__main__":
    main()
