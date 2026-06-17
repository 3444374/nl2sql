from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "assets" / "opening_template.pptx"
OUT = ROOT / "docs" / "opening" / "opening_ppt_template_version.pptx"

TITLE = "面向 SQL+ 简化查询表达的多智能体自然语言数据库查询生成与反馈修正方法研究"
FONT = "Microsoft YaHei"
PRIMARY = RGBColor(31, 78, 121)
DARK = RGBColor(45, 45, 45)
MUTED = RGBColor(100, 100, 100)
ACCENT = RGBColor(198, 89, 17)
LIGHT = RGBColor(238, 242, 247)


def clear_template_slides(prs: Presentation) -> None:
    """Keep the template theme/layouts, remove existing sample slides."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_text_frame(tf, text: str, font_size: int = 20, bold: bool = False, color=DARK) -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_title(slide, title: str) -> None:
    if slide.shapes.title is not None:
        set_text_frame(slide.shapes.title.text_frame, title, 28, True, PRIMARY)


def add_footer(slide, idx: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.92), Inches(9.0), Inches(0.25))
    tf = box.text_frame
    set_text_frame(tf, f"开题汇报 · SQL+ / NL2SQL / Multi-Agent  ·  {idx}", 8, False, MUTED)


def add_bullets(slide, items: list[str], x=0.75, y=1.75, w=8.6, h=4.6, font_size=18) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_note(slide, text: str, x=0.75, y=6.25, w=8.6, h=0.45) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    set_text_frame(tf, text, 11, False, ACCENT)


def add_table(slide, rows: list[list[str]], x=0.55, y=1.45, w=9.0, h=4.8, font_size=10) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = PRIMARY if r == 0 else (LIGHT if r % 2 == 0 else RGBColor(255, 255, 255))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size if r else font_size + 1)
                    run.font.bold = r == 0
                    run.font.color.rgb = RGBColor(255, 255, 255) if r == 0 else DARK


def add_flow(slide, steps: list[str], x=0.35, y=2.1, box_w=1.1, box_h=0.62, gap=0.12) -> None:
    for i, step in enumerate(steps):
        left = Inches(x + i * (box_w + gap))
        shape = slide.shapes.add_shape(1, left, Inches(y), Inches(box_w), Inches(box_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY if i % 2 == 0 else ACCENT
        shape.line.color.rgb = RGBColor(255, 255, 255)
        tf = shape.text_frame
        set_text_frame(tf, step, 10, True, RGBColor(255, 255, 255))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_metric_cards(slide, cards: list[tuple[str, str, str]], x=0.55, y=1.55) -> None:
    for i, (label, value, desc) in enumerate(cards):
        col = i % 3
        row = i // 3
        left = Inches(x + col * 3.1)
        top = Inches(y + row * 1.45)
        shape = slide.shapes.add_shape(1, left, top, Inches(2.85), Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = PRIMARY
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = FONT
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.name = FONT
        p2.font.size = Pt(19)
        p2.font.bold = True
        p2.font.color.rgb = PRIMARY
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = FONT
        p3.font.size = Pt(8)
        p3.font.color.rgb = DARK


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(slide, title)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            set_text_frame(shape.text_frame, subtitle, 18, False, DARK)
    add_footer(slide, 1)


def add_content_slide(prs: Presentation, idx: int, title: str, bullets: list[str], note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    add_bullets(slide, bullets)
    if note:
        add_note(slide, note)
    add_footer(slide, idx)


def add_table_slide(prs: Presentation, idx: int, title: str, rows: list[list[str]], note: str = "", font_size: int = 10) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9.0), Inches(0.55))
    set_text_frame(title_box.text_frame, title, 25, True, PRIMARY)
    add_table(slide, rows, font_size=font_size)
    if note:
        add_note(slide, note)
    add_footer(slide, idx)


def build() -> None:
    prs = Presentation(str(TEMPLATE))
    clear_template_slides(prs)

    add_title_slide(
        prs,
        TITLE,
        "研究方向：达梦 SQL+ / Text-to-SQL / 多智能体 / 执行反馈修正\n汇报人：待填写    导师：待填写    时间：待填写",
    )

    add_content_slide(
        prs,
        2,
        "汇报提纲",
        [
            "研究背景与问题提出",
            "国内外研究现状与不足",
            "研究目标、关键问题与技术路线",
            "SQL+ 表达设计与多智能体反馈修正框架",
            "实验设计、当前初步结果与后续计划",
        ],
    )

    add_content_slide(
        prs,
        3,
        "研究背景：自然语言数据库查询仍然不稳定",
        [
            "业务人员希望直接用自然语言完成数据查询和统计分析，降低 SQL 使用门槛。",
            "LLM 提升了 Text-to-SQL 能力，但真实数据库中仍存在 schema 复杂、多表 join、聚合口径、排序逻辑和方言差异问题。",
            "标准 SQL 出错后，错误位置往往不清楚，直接整条重生成容易破坏原本正确的部分。",
            "本课题关注的不只是生成 SQL，而是生成、执行、诊断和局部修复的闭环能力。",
        ],
    )

    add_content_slide(
        prs,
        4,
        "核心动机：为什么引入 SQL+",
        [
            "标准 SQL 的书写顺序与自然语言意图构造顺序不完全一致，复杂查询存在较强跨子句耦合。",
            "SQL+ 借鉴 GoogleSQL Pipe Syntax 的线性数据流思想，但目标不是替代 SQL，而是作为 NL2SQL 的中间表示。",
            "SQL+ 将 FROM、JOIN、WHERE、GROUP、AGG、SELECT、ORDER 等操作显式拆成步骤。",
            "步骤化表达可以为 Critic Agent 定位错误、Skill Router 路由修复、Repair Skill 局部 patch 提供载体。",
        ],
        "关键表述：SQL+ 不一定更短，但更线性、更低耦合、更适合错误定位和局部修复。",
    )

    add_content_slide(
        prs,
        5,
        "国内外研究现状：Text-to-SQL",
        [
            "传统方法：Seq2SQL、SyntaxSQLNet、RAT-SQL、PICARD、RESDSQL 等关注结构解码、schema linking 和语法约束。",
            "LLM 方法：DAIL-SQL、DIN-SQL、C3、CHESS、CHASE-SQL 等引入 prompt、任务分解、候选生成和执行验证。",
            "Benchmark 发展：Spider 关注跨数据库泛化，BIRD 强调真实数据库、外部知识和效率，Spider 2.0 关注企业级工作流。",
            "评价趋势：从 exact match / execution accuracy 扩展到 valid rate、schema/value linking、token cost、latency 和修复能力。",
        ],
    )

    add_content_slide(
        prs,
        6,
        "国内外研究现状：中间表示与 Agent",
        [
            "SemQL 和 NatSQL 说明标准 SQL 并不是唯一生成目标，中间表示可以降低初次生成难度。",
            "GoogleSQL Pipe Syntax 说明 SQL 生态内部也可以引入线性表达，改善可读性和维护性。",
            "MAC-SQL、Tool-Assisted Agent、LEVER、SQLCritic 等表明多角色协作、执行反馈和候选验证具有价值。",
            "不足：多数工作仍偏向标准 SQL 层候选选择或整条修复，较少研究“中间步骤级错误定位 + 局部 repair skill”。",
        ],
    )

    add_content_slide(
        prs,
        7,
        "拟解决的关键问题",
        [
            "SQL+ 必要性：相比标准 SQL、SemQL、NatSQL、Pipe-style 表示，SQL+ 在生成、转换、诊断和修复上新增什么价值？",
            "SQL+ 设计边界：如何在简化表达、可转换、可解释和局部修复之间取得平衡？",
            "SQL+ 生成：如何处理 schema linking、value linking、join path、aggregation planning 和 projection？",
            "SQL+ 修复：如何把执行反馈和结果异常映射到 SQL+ 步骤，并路由到对应 repair skill？",
            "实验验证：如何从准确率、成本、错误定位、patch 范围和修复轮数多维度证明方法有效？",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9), Inches(0.55))
    set_text_frame(title_box.text_frame, "总体技术路线", 25, True, PRIMARY)
    add_flow(
        slide,
        ["自然语言", "SQL+生成", "SQL转换", "执行反馈", "Critic", "Router", "Repair", "Executor"],
        x=0.35,
        y=2.0,
        box_w=1.05,
    )
    add_bullets(
        slide,
        [
            "核心闭环：Natural language -> SQL+ -> SQL -> Execution feedback -> Critic Agent -> Skill Router -> Repair Skill -> Executor",
            "Schema Agent / Planner Agent 负责前端理解与查询规划，Critic / Router / Repair Skill 负责后端诊断与修复。",
            "每个模块都保留可观察输出，避免多智能体退化为多个 prompt 的简单串联。",
        ],
        y=3.15,
        font_size=16,
    )
    add_footer(slide, 8)

    add_table_slide(
        prs,
        9,
        "SQL+ 最小语法子集",
        [
            ["步骤", "作用", "易错位置"],
            ["FROM", "指定数据源", "主表选择错误"],
            ["JOIN", "指定连接路径", "连接方向、冗余 join、缺 join"],
            ["WHERE", "过滤条件", "值链接、日期边界、隐含条件"],
            ["GROUP / AGG", "分组与聚合", "COUNT/SUM 口径、聚合别名"],
            ["HAVING", "聚合后过滤", "聚合别名引用"],
            ["SELECT", "最终投影", "列多、列少、顺序不一致"],
            ["ORDER / LIMIT", "排序与 top-k", "排序字段、方向、数量"],
        ],
        "设计原则：不追求重新发明 SQL，而是保留可转换性和步骤级可修复性。",
        font_size=11,
    )

    add_content_slide(
        prs,
        10,
        "多智能体框架设计",
        [
            "Schema Agent：检索相关表、字段、外键关系、候选字段值和 join 路径。",
            "Planner Agent：把自然语言问题拆成 SQL+ 步骤草图。",
            "SQL+ Generator / Translator：生成 SQL+ 并转换为可执行 SQL。",
            "Critic Agent：输出 likely_error_type、可疑步骤和诊断证据。",
            "Skill Router：把错误类型路由到 value-linking、ORDER、aggregation、join、projection 等局部 repair skill。",
            "Executor：执行候选 patch，过滤语法错误和执行错误，选择最终结果。",
        ],
    )

    add_table_slide(
        prs,
        11,
        "错误类型与 Repair Skill",
        [
            ["错误类型", "典型表现", "对应修复能力"],
            ["value-linking", "值拼写、日期边界、状态值不匹配", "候选值检索、过滤条件替换"],
            ["ORDER/LIMIT", "排序遗漏、字段错误、方向错误", "排序字段与 top-k 修复"],
            ["aggregation", "GROUP 维度、COUNT 口径、AGG 别名", "聚合口径和别名引用修复"],
            ["join", "join 路径、冗余 join、缺 join", "连接路径补全与规范化"],
            ["projection", "结果列多、少或列顺序错误", "SELECT 投影局部修复"],
        ],
        font_size=10,
    )

    add_content_slide(
        prs,
        12,
        "实验数据、对比方法与指标",
        [
            "数据层次：自建订单分析数据集、SQL+ 已知失败集、Spider dev 小规模受支持子集，后续扩展 BIRD / 达梦样例。",
            "对比方法：Direct NL2SQL、NL2SQL+ single agent、SemQL-style / NatSQL-style / Pipe-style proxy、SQL 层修复、SQL+ Skill Router + Repair Skills。",
            "结果质量：execution accuracy、valid SQL rate、SQL+ valid rate、SQL executable rate。",
            "表达复杂度：token length、step/clause count、nesting depth、alias dependency、cross-clause reference、conversion time。",
            "修复能力：repair success、error localization accuracy、router accuracy、patch minimality、repair rounds、repair token cost、latency。",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9), Inches(0.55))
    set_text_frame(title_box.text_frame, "当前初步实验总览", 25, True, PRIMARY)
    add_metric_cards(
        slide,
        [
            ("SQL+ 转换", "30/30", "自建订单数据集执行一致"),
            ("Direct NL2SQL", "16/30", "直接生成 SQL baseline"),
            ("SQL+ prompt v2", "17/30", "略高于 Direct，但差距小"),
            ("SQL+ 单 Refiner", "4/13", "非 gold 执行反馈"),
            ("Direct SQL Refiner", "6/14", "SQL 层修复对照"),
            ("Skill Router v3", "13/13", "当前已知失败集"),
            ("Spider smoke", "20/20", "受支持小子集，不是完整跑分"),
            ("SQL+ conversion", "0.007ms", "平均确定性转换时间"),
            ("Repair skill", "5 类", "value/order/agg/join/projection"),
        ],
    )
    add_footer(slide, 13)

    add_table_slide(
        prs,
        14,
        "SQL+ 转换实验",
        [
            ["指标", "结果", "说明"],
            ["SQL+ 样例数", "30", "自建企业订单分析样例"],
            ["SQL+ 语法通过", "30/30", "parser 可解析"],
            ["转换 SQL 可执行", "30/30", "SQLite 执行成功"],
            ["执行结果一致", "30/30", "与 gold SQL 结果一致"],
        ],
        "结论：SQL+ 已形成可解析、可转换、可执行的最小闭环。",
        font_size=12,
    )

    add_table_slide(
        prs,
        15,
        "IR 表达复杂度对比",
        [
            ["表示形式", "平均 token", "步骤/子句", "嵌套深度", "别名依赖", "跨子句引用", "转换成功"],
            ["Standard SQL", "31.5333", "5.9", "0.6667", "2.0333", "2.3333", "30/30"],
            ["SQL+", "35.0333", "6.1333", "0.6667", "0.7", "1.0", "30/30"],
            ["SemQL-style proxy", "50.5667", "10.7333", "3.6667", "0.9", "1.2", "N/A"],
            ["NatSQL-style proxy", "31.5", "5.4333", "0.9667", "1.3667", "1.6667", "N/A"],
            ["Pipe-style proxy", "40.8", "6.1333", "0.6667", "1.3667", "1.6667", "N/A"],
        ],
        "结论：SQL+ 不更短，但别名依赖和跨子句引用更低，更适合错误定位与局部修复。",
        font_size=8,
    )

    add_table_slide(
        prs,
        16,
        "IR 生成成本与执行效果",
        [
            ["方法", "表示有效", "SQL 可执行", "执行一致", "平均总 token", "平均延迟"],
            ["Direct SQL", "30/30", "30/30", "12/30", "599.1667", "6.5851s"],
            ["SQL+", "28/30", "28/30", "14/30", "813.0333", "9.2197s"],
            ["NatSQL-style proxy", "30/30", "30/30", "13/30", "740.7667", "6.2802s"],
            ["SemQL-style proxy", "30/30", "25/30", "12/30", "1028.9667", "9.9684s"],
        ],
        "结论：SQL+ 初次生成并无显著成本优势，后续重点验证修复收益是否抵消额外开销。",
        font_size=10,
    )

    add_table_slide(
        prs,
        17,
        "Baseline 与失败类型分析",
        [
            ["方法", "样例数", "SQL/SQL+ 有效", "SQL 可执行", "执行一致"],
            ["Direct NL2SQL", "30", "-", "30/30", "16/30"],
            ["NL2SQL+ prompt v1", "30", "30/30", "30/30", "13/30"],
            ["NL2SQL+ prompt v2", "30", "30/30", "30/30", "17/30"],
            ["失败类型", "数量", "说明", "", ""],
            ["filter/value-linking", "5", "字段值、日期、状态值错误", "", ""],
            ["ORDER/LIMIT", "3", "排序字段、方向或 top-k 错误", "", ""],
            ["aggregation / join / projection", "5", "聚合口径、连接路径、投影列问题", "", ""],
        ],
        "结论：失败主要是语义错误，不是语法错误。",
        font_size=9,
    )

    add_table_slide(
        prs,
        18,
        "反馈修正实验结果",
        [
            ["方法", "失败样例", "SQL+ 有效", "SQL 可执行", "修复成功"],
            ["SQL+ 非 gold Refiner v2", "13", "13/13", "12/13", "4/13"],
            ["Direct SQL 非 gold Refiner", "14", "-", "14/14", "6/14"],
            ["SQL+ Schema-Critic-Refiner", "13", "13/13", "13/13", "3/13"],
            ["SQL+ Step-wise Critic-Refiner", "13", "13/13", "12/13", "3/13"],
            ["SQL+ Skill Router + Repair Skills v3", "13", "13/13", "13/13", "13/13"],
        ],
        "结论：简单多 Agent 串联不一定提升，Critic + Router + 局部 skill 的闭环更稳定。",
        font_size=9,
    )

    add_table_slide(
        prs,
        19,
        "Repairability 指标对比",
        [
            ["方法", "样例数", "修复成功", "定位准确率", "最小 patch 率", "平均 repair token"],
            ["Direct SQL Refiner", "14", "6/14", "0.8571", "0.8571", "1609.3571"],
            ["SQL+ Critic Router Skills", "13", "13/13", "0.7692", "0.9231", "3813.9231"],
            ["Direct SQL overlap", "9", "4/9", "0.8889", "0.8889", "1583.2222"],
            ["SQL+ Router overlap", "9", "9/9", "0.7778", "0.8889", "4001.7778"],
        ],
        "结论：SQL+ 修复成功率和 patch 可控性更好，但 Critic token 成本更高。",
        font_size=9,
    )

    add_table_slide(
        prs,
        20,
        "分治 Repair Skill 实验",
        [
            ["Repair Skill", "样例数", "修复成功", "覆盖问题"],
            ["value-linking", "3", "3/3", "候选值、日期边界、状态值"],
            ["ORDER", "3", "3/3", "排序字段、排序方向、LIMIT"],
            ["aggregation", "3", "3/3", "COUNT 口径、GROUP 维度、AGG 别名"],
            ["join", "3", "3/3", "JOIN 方向、冗余 JOIN、缺失 JOIN、paid 过滤"],
            ["projection", "1", "1/1", "结果列多、列少或列顺序错误"],
        ],
        "结论：不同错误类型需要不同的局部修复策略，整体重写不是最优路径。",
        font_size=10,
    )

    add_table_slide(
        prs,
        21,
        "Spider 小规模 Benchmark Smoke Test",
        [
            ["指标", "结果", "边界说明"],
            ["数据集", "Spider dev 子集", "只验证迁移可行性"],
            ["数据库", "concert_singer", "当前 SQL+ 子集可覆盖"],
            ["样例数", "20", "非完整 Spider"],
            ["SQL+ 有效", "20/20", "转换链路可用"],
            ["SQL 可执行", "20/20", "执行无错误"],
            ["执行一致", "20/20", "与 gold SQL 结果一致"],
        ],
        "结论：这是公开 benchmark 子集迁移证据，不是完整 Spider benchmark 分数。",
        font_size=10,
    )

    add_content_slide(
        prs,
        22,
        "当前结果对研究假设的支撑",
        [
            "SQL+ 转换 30/30 说明中间表示和确定性转换具备基础可行性。",
            "IR 对比说明 SQL+ 不应被表述为更短，而应表述为更低耦合、更适合步骤级定位。",
            "Direct SQL 与 SQL+ 初次生成差距很小，说明仅替换输出格式不足以解决 NL2SQL。",
            "Skill Router v3 的 13/13 说明错误类型路由和局部 repair skill 是值得继续研究的方向。",
            "Spider smoke test 20/20 说明当前 SQL+ 子集具备初步公开数据集迁移可能。",
        ],
    )

    add_content_slide(
        prs,
        23,
        "当前不足与风险",
        [
            "自建数据集规模仍小，查询类型以常见分析查询为主。",
            "SemQL-style / NatSQL-style 目前是 controlled proxy，不是完整系统复现。",
            "Skill Router v3 的 13/13 来自当前已知 SQL+ 失败集，不能外推为大规模 benchmark 性能。",
            "旧 Direct SQL 与 SQL+ Critic 输出没有完整 API latency，端到端修复延迟仍需重跑。",
            "达梦 SQL 方言适配尚未完成，当前主要在 SQLite / 标准 SQL 链路上验证机制。",
        ],
    )

    add_content_slide(
        prs,
        24,
        "后续研究计划",
        [
            "扩展 SQL+ 语法：覆盖复杂子查询、LEFT JOIN、窗口函数、复杂布尔条件等结构。",
            "完善 Schema Agent 与 Critic Agent：引入 schema/value lookup、SQL+ parser、SQL executor 和结果预览工具。",
            "扩展无报错语义错：重点处理 projection mismatch、复合聚合错误和多步修复顺序。",
            "扩大公开数据集验证：从 Spider 单库 smoke test 扩展到多数据库、多难度子集，再考虑 BIRD。",
            "适配达梦 SQL 方言：验证 SQL+ 到达梦 SQL 的函数、分页、日期和方言兼容性。",
        ],
    )

    add_content_slide(
        prs,
        25,
        "预期创新点",
        [
            "面向生成和修复的 SQL+ 中间查询表示：不是替代 SQL，而是服务于 NL2SQL 生成、转换、诊断和局部修复的步骤化表达层。",
            "面向 SQL+ 的多智能体诊断与技能路由机制：将 Critic Agent、Skill Router、Repair Skill 和 Executor 组织为反馈闭环。",
            "面向 repairability 的评价体系：除 execution accuracy 外，引入 error localization accuracy、router accuracy、patch minimality、repair rounds、token cost 和 latency。",
        ],
    )

    add_content_slide(
        prs,
        26,
        "总结",
        [
            "本课题研究重点不是一个简单 Text-to-SQL Demo，而是 SQL+ 中间表示与多智能体反馈修正机制。",
            "当前实验证明 SQL+ 表达与转换链路可行，且局部 repair skill 在已知失败集上具有初步效果。",
            "SQL+ 的优势不能概括为更短或初次生成更准，而应落在低耦合、可定位、可路由和可修复上。",
            "后续重点是扩大样例规模、增强真实非 gold 诊断、补齐 latency 成本，并适配公开数据集和达梦 SQL 方言。",
        ],
    )

    add_content_slide(
        prs,
        27,
        "主要参考文献",
        [
            "Spider, BIRD, Spider 2.0：Text-to-SQL benchmark 与真实场景评估。",
            "IRNet/SemQL, NatSQL：Text-to-SQL 中间表示研究。",
            "RAT-SQL, PICARD, RESDSQL：schema linking、语法约束与结构化解析。",
            "DAIL-SQL, DIN-SQL, CHESS, CHASE-SQL：LLM-based Text-to-SQL、任务分解与候选选择。",
            "MAC-SQL, SQLCritic, Tool-Assisted Agent, LEVER：多智能体、执行反馈和修复评估。",
            "GoogleSQL Pipe Syntax：SQL 线性表达与 SQL 语言扩展。",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
