from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
SRC = ROOT / "docs" / "opening" / "opening_ppt_template_version_v4.pptx"
OUT = ROOT / "docs" / "opening" / "opening_ppt_template_version_v5_charts.pptx"
CHART_DIR = ROOT / "docs" / "opening" / "assets" / "charts"

SLIDE_CHARTS = {
    8: ("动机测试一：SQL+ 表达复杂度是否更利于修复", "chart_ir_complexity.png"),
    9: ("动机测试二：生成成本与执行效果", "chart_ir_generation_cost.png"),
    17: ("可行性测试一：初次生成 baseline 对比", "chart_baseline_generation.png"),
    18: ("可行性测试二：反馈修正策略对比", "chart_repair_success.png"),
    20: ("可行性测试四：Spider 小规模子集验证", "chart_spider_subset.png"),
}

EXPERIMENT_NOTES = {
    8: """【汇报讲稿】
这一页回答第一个动机问题：为什么需要 SQL+ 这种中间表示，而不是直接生成标准 SQL。实验数据来自自建订单分析数据集的 30 条查询。对同一批 gold 查询，我们构造 Standard SQL、SQL+、SemQL-style proxy、NatSQL-style proxy 和 Pipe-style proxy 五种表示。这里的 proxy 是受控近似表示，用于比较表达形态，不等同于完整复现 SemQL、NatSQL 或 GoogleSQL Pipe Syntax 系统。
实验变量是中间表示形式；固定条件是同一批查询、同一 schema、同一套统计脚本。图中展示两个直接影响 repairability 的指标：别名依赖和跨子句引用。别名依赖指 SELECT、ORDER、HAVING 等位置对 AS 别名的后续引用次数；跨子句引用指一个逻辑字段或表达式跨多个 SQL 区域反复出现的次数。结果显示，SQL+ 的平均别名依赖为 0.7，跨子句引用为 1.0，低于 Standard SQL 的 2.03 和 2.33。
结论不是“SQL+ 更短”，而是 SQL+ 把查询拆成了更清晰的步骤边界，降低了跨区域耦合。这个结果支持后续把错误定位到 WHERE、JOIN、AGG、ORDER 等局部步骤中修复。
【答辩备注】
可能问题：proxy 是否等同于原始 SemQL 或 NatSQL？回答：不是。这里是开题阶段的 controlled proxy，只用于同批样例下比较表达形态和结构依赖，不能声称复现原系统 benchmark。后续如时间允许，会接入原始系统或公开实现做更严格比较。
可能问题：为什么不把 token 更短作为 SQL+ 目标？回答：实验已经显示 SQL+ token 并不更短，因此本课题的论证重点调整为步骤边界、错误定位和局部修复收益。""",
    9: """【汇报讲稿】
这一页继续回答 SQL+ 的代价问题。实验仍使用同一批 30 条自建订单分析查询，模型为 gpt-5-mini。比较 Direct SQL、SQL+、NatSQL-style proxy 和 SemQL-style proxy 四种生成目标。实验过程是：模型生成目标表示；脚本检查表示是否有效；若需要则转换为 SQL；在同一 SQLite 数据库上执行；将执行结果与 gold SQL 结果比较；同时记录 OpenAI API usage 中的 token 和本地脚本记录的 latency。
图中柱状图是 execution match，折线是平均 token。SQL+ 的 execution match 为 14/30，Direct SQL 为 12/30，NatSQL-style proxy 为 13/30，SemQL-style proxy 为 12/30。SQL+ 结果略高，但 token 成本为 813.03，高于 Direct SQL 的 599.17，平均延迟也更高。
因此这里的结论要保守：SQL+ 在初次生成阶段没有明显成本优势，也不能仅凭 14/30 说明显著提高准确率。它的价值需要在后续错误定位、patch minimality、repair rounds 和 repair success 中继续验证。
【答辩备注】
可能问题：SQL+ 成本更高，为什么还要用？回答：因为本课题不是把 SQL+ 作为更省 token 的语言，而是作为可诊断、可局部修复的中间表示。如果修复阶段能减少整条 SQL 重写和定位困难，额外生成成本才有研究价值。
可能问题：为什么用 gpt-5-mini？回答：开题阶段选择常见 API 调用形式的小模型做可重复 smoke test，重点验证方法链路和指标设计。后续可以扩展到更强模型做横向对比。""",
    17: """【汇报讲稿】
这一页展示初次生成 baseline。数据集是自建订单分析数据集的 30 条自然语言问题，每条都有 gold SQL 和 gold SQL+。Direct SQL 让模型直接生成 SQL；NL2SQL+ prompt v1 和 v2 让模型先生成 SQL+，再由本地 converter 转成 SQL。评价指标是生成表示是否有效、SQL 是否可执行、执行结果是否与 gold SQL 一致。
图中只展示最关键的 execution match。Direct SQL 为 16/30，NL2SQL+ prompt v1 为 13/30，prompt v2 为 17/30。v2 比 Direct SQL 高 1 条，但幅度很小。
这个实验的作用不是证明 SQL+ 初次生成显著优于 SQL，而是说明仅换输出格式不够。后续必须引入 Schema Agent、Critic Agent、Skill Router 和 Repair Skill，把 SQL+ 的步骤化优势用于诊断与修复。
【答辩备注】
可能问题：为什么初次生成提升不明显还继续做 SQL+？回答：因为当前研究假设已经从“初次生成更准”调整为“错误更容易定位和修复”。这个 baseline 反而说明研究重点不能停留在 prompt，而应进入反馈修正闭环。""",
    18: """【汇报讲稿】
这一页展示反馈修正策略对比。样例来自当前 SQL+ prompt v2 的 13 条已知失败样例，以及 Direct SQL 的 14 条失败样例。失败类型包括 value-linking、ORDER、aggregation、join 和 projection。对比方法包括 SQL+ 单 Refiner、Direct SQL 单 Refiner、Schema-Critic-Refiner、Step-wise Critic-Refiner 和 SQL+ Skill Router + Repair Skills v3。
修复成功的判定方式是：修复后的 SQL 或 SQL+ 转 SQL 在 SQLite 中执行，并且执行结果与 gold SQL 一致。图中可以看到，SQL+ 单 Refiner 为 4/13，Direct SQL Refiner 为 6/14，Schema-Critic 和 Step-wise Critic 都是 3/13，而 SQL+ Skill Router + Repair Skills v3 在当前已知失败集上达到 13/13。
结论是，多 Agent 数量本身不是关键，关键在于分工是否形成可检查中间产物：Critic 负责定位错误类型和步骤，Skill Router 负责选择修复技能，Repair Skill 只做局部 patch，Executor 负责验证候选。
【答辩备注】
可能问题：13/13 是否能代表方法已经有效？回答：不能。它只代表当前 13 条已知失败集上的链路可行性，不能表述为大规模 benchmark 结果。开题中应把它作为 repairability 方向成立的初步证据，后续需要扩展未知失败样例和多数据库 Spider 子集。""",
    20: """【汇报讲稿】
这一页展示 Spider 小规模子集验证。数据来自 Spider dev 中 concert_singer 数据库，目前选取 20 条 SQL+ 子集能够覆盖的查询。这里必须区分两类实验：conversion smoke test 和端到端生成实验。
conversion smoke test 是从 Spider gold SQL 出发，人工或规则映射为 SQL+，再由 SQL+ converter 转回 SQL，并比较执行结果，因此 20/20 只说明 SQL+ 表达与转换机制可迁移，不是模型端到端生成准确率。端到端实验中，fresh e2e 让系统从自然语言和 schema 生成 SQL+，结果是 19/20；同一批 fresh 输出经过 Skill Router 到 semantic repair skill 后达到 20/20。
结论是，SQL+ 链路在公开 benchmark 小子集上已经具备初步可迁移性，但目前仍是单数据库、小规模、受支持 SQL+ 子集，不应写成完整 Spider benchmark 成绩。
【答辩备注】
可能问题：为什么不直接跑完整 Spider？回答：开题阶段目标是验证方法可行性和实验设计，不是追榜。完整 Spider 需要更多数据库、更多 SQL 结构覆盖和稳定 schema linking。后续计划先扩展多数据库小子集，再逐步扩大样例规模。""",
}

GENERAL_NOTES = {
    1: """【汇报讲稿】
这一页介绍课题题目和定位。本课题不是做一个单轮 Text-to-SQL demo，而是研究 SQL+ 中间表示、多智能体诊断、技能路由和执行反馈如何形成可修复闭环。
【答辩备注】
如果被问 SQL+ 是否替代 SQL，回答：不是。SQL+ 是面向生成与修复的中间表示，最终仍转换为标准 SQL 或达梦 SQL 执行。""",
    2: """【汇报讲稿】
这一页说明汇报结构：背景动机、研究现状、方法设计、初步实验、后续计划。后面的实验会围绕“为什么 SQL+”和“为什么多 Agent 修复”展开。
【答辩备注】
注意控制节奏，实验部分重点讲变量、指标、边界，不只报数字。""",
    3: """【汇报讲稿】
这一页说明自然语言数据库查询的真实困难：直接生成 SQL 容易遇到 schema linking、join 路径、聚合口径和语义偏差问题。
【答辩备注】
强调“能执行”不等于“语义正确”，这正是需要反馈修正的原因。""",
    4: """【汇报讲稿】
这一页梳理研究发展：从 Spider 等 benchmark，到 SemQL/NatSQL 等中间表示，再到 DIN-SQL、MAC-SQL、CHASE-SQL、SQLCritic 等过程化与 agentic 方法。
【答辩备注】
研究空间来自已有工作对分解、候选、反馈的证明，但对 SQL+ 步骤级定位、Skill Router 和局部 repair skill 的系统研究仍不足。""",
    5: """【汇报讲稿】
这一页总结当前研究缺口：任务越来越真实，系统越来越复杂，评价也需要从准确率扩展到可定位、可修复、成本和延迟。
【答辩备注】
不要说已有工作“不好”，要说它们主要围绕最终 SQL，本课题补充中间步骤层的反馈修正。""",
    6: """【汇报讲稿】
这一页说明动机测试设计。实验先回答 SQL+ 是否有结构优势，再回答这种结构优势能否转化为修复收益。
【答辩备注】
强调实验不是为了证明 SQL+ 更短，而是验证 SQL+ 是否更适合错误定位和局部修复。""",
    7: """【汇报讲稿】
这一页说明统一实验口径：同一批样例、同一 schema、同一执行评价方式，对比不同表示和修复策略。
【答辩备注】
指标包括 valid rate、execution match、token、latency、patch minimality 和 repair rounds。""",
    10: """【汇报讲稿】
这一页把前面的动机测试收束为系统设计：SQL+ 负责步骤化表达，多 Agent 负责意图、schema、规划、诊断、路由和修复。
【答辩备注】
突出每个 agent 都要产出可检查中间产物，而不是简单 prompt 串联。""",
    11: """【汇报讲稿】
这一页给出研究不足与本课题切入点。中间表示、Agent 框架、反馈修正都已有基础，但 SQL+ 层局部修复仍有研究空间。
【答辩备注】
可用一句话概括：把错误从最终 SQL 压回 SQL+ 步骤中修。""",
    12: """【汇报讲稿】
这一页说明 SQL+ 与反馈修正的关系。SQL+ 的每个步骤都可以成为 Critic 定位和 Repair Skill 修改的最小单元。
【答辩备注】
强调不是多写一层格式，而是为了让错误有明确落点。""",
    13: """【汇报讲稿】
这一页通过原始 SQL 与 SQL+ 对照说明复杂查询如何被拆成可诊断步骤。SQL+ 保留 FROM、JOIN、WHERE、GROUP、AGG、SELECT、ORDER、LIMIT 等步骤边界。
【答辩备注】
讲重点：WHERE 查过滤，JOIN 查连接，AGG 查聚合口径，ORDER 查 top-k。""",
    14: """【汇报讲稿】
这一页展示系统架构：自然语言到 SQL+，SQL+ 到 SQL，执行后反馈给 Critic，Router 选择 skill，Repair Skill 产生局部 patch，再回到转换和执行。
【答辩备注】
强调闭环位置：失败反馈不是回到最开始整条重写，而是回到 SQL+ 局部步骤。""",
    15: """【汇报讲稿】
这一页说明可行性测试路径。先证明转换链路，再证明修复链路，最后用 Spider 小子集检查迁移可能性。
【答辩备注】
如果被问规模小，回答：开题阶段是机制验证，后续计划扩展多数据库与更复杂 SQL。""",
    16: """【汇报讲稿】
这一页说明实验数据、模型和指标。自建 30 条样例用于机制对比，已知失败集用于修复链路，Spider concert_singer 20 条用于公开子集迁移验证。
【答辩备注】
gold SQL 用于评价，不用于真实生成；conversion smoke 和 fresh e2e 必须区分。""",
    19: """【汇报讲稿】
这一页展示分治 repair skill：value、order、aggregation、join、projection 等不同错误类型分别用不同局部技能处理。
【答辩备注】
强调分治的意义是缩小 patch 范围，提高可解释性，而不是把 prompt 拆得更多。""",
    21: """【汇报讲稿】
这一页总结当前可行性：SQL+ 转换链路跑通，修复链路在已知失败集上有效，Spider 小子集端到端也有初步结果。
【答辩备注】
主动说明限制：小样例、小子集、成本仍高，后续要扩大实验。""",
    22: """【汇报讲稿】
这一页说明后续实验计划和指标。后续不只看 execution accuracy，还会看 error localization、router accuracy、patch minimality、repair rounds、token 和 latency。
【答辩备注】
这页回应导师关于“研究内容不能像工程步骤”的反馈：每个方向都对应技术难点和评价指标。""",
    23: """【汇报讲稿】
这一页说明进度和风险控制。当前完成了 SQL+ parser、转换器、baseline、repair skills、Skill Router 和 Spider 小子集验证。
【答辩备注】
风险包括数据规模、proxy 对比边界、达梦方言适配和模型成本；对应都有后续计划。""",
    24: """【汇报讲稿】
这一页总结预期创新：面向修复的 SQL+ 中间表示、多 Agent 反馈修正闭环，以及 repairability 评价体系。
【答辩备注】
表述要稳：是“尝试提出”和“验证一种方法”，不要夸大为已经全面优于已有系统。""",
    25: """【汇报讲稿】
这一页列主要参考文献。文献覆盖 benchmark、中间表示、LLM Text-to-SQL、多智能体和反馈修正。
【答辩备注】
如果被问具体支撑：Spider/BIRD 支撑 benchmark，SemQL/NatSQL/Pipe Syntax 支撑表示对比，DIN-SQL/MAC-SQL/CHASE-SQL/SQLCritic 支撑分解、agent 与反馈修正。""",
}


def remove_body_shapes(slide):
    keep = []
    for shape in list(slide.shapes):
        top = getattr(shape, "top", 0)
        # Keep template header and footer only. Body content will be rebuilt.
        if top < Inches(0.95) or top > Inches(6.65):
            keep.append(shape)
            continue
        slide.shapes._spTree.remove(shape._element)
    return keep


def add_title(slide, title: str):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.05), Inches(11.9), Inches(0.48))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(23, 50, 77)


def add_chart(slide, chart_file: str):
    path = CHART_DIR / chart_file
    if not path.exists():
        raise FileNotFoundError(path)
    slide.shapes.add_picture(str(path), Inches(0.72), Inches(1.65), width=Inches(11.9), height=Inches(5.35))


def set_notes(slide, note: str):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = "\n".join(line.strip() for line in note.splitlines() if line.strip())


def main():
    if not SRC.exists():
        raise FileNotFoundError(SRC)
    shutil.copyfile(SRC, OUT)
    prs = Presentation(str(OUT))

    for slide_no, (title, chart_file) in SLIDE_CHARTS.items():
        slide = prs.slides[slide_no - 1]
        remove_body_shapes(slide)
        add_title(slide, title)
        add_chart(slide, chart_file)

    for idx, slide in enumerate(prs.slides, start=1):
        note = EXPERIMENT_NOTES.get(idx) or GENERAL_NOTES.get(idx)
        if note:
            set_notes(slide, note)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
